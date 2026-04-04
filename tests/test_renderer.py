from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from ppt_creator.preview import (
    PREVIEW_HEIGHT,
    PREVIEW_MANIFEST_FILENAME,
    PREVIEW_WIDTH,
    PreviewRenderer,
    compare_preview_directories,
    format_visual_regression_failure,
    promote_preview_baseline,
    render_previews,
    visual_regression_has_failures,
)
from ppt_creator.renderer import PresentationRenderer
from ppt_creator.schema import PresentationInput


def test_render_example_creates_pptx(tmp_path: Path) -> None:
    spec = PresentationInput.from_path("examples/ai_sales.json")
    output = tmp_path / "example.pptx"

    renderer = PresentationRenderer(asset_root="examples")
    rendered = renderer.render(spec, output)

    assert rendered.exists()
    presentation = Presentation(str(rendered))
    assert len(presentation.slides) == 10


def test_renderer_requires_pptx_output_extension(tmp_path: Path) -> None:
    spec = PresentationInput.from_path("examples/ai_sales.json")
    renderer = PresentationRenderer(asset_root="examples")

    with pytest.raises(ValueError):
        renderer.render(spec, tmp_path / "example.txt")


def test_panel_content_box_uses_theme_padding() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    renderer = PresentationRenderer(asset_root="examples")

    box = renderer.panel_content_box(slide, left=1.0, top=1.0, width=4.0, height=2.0)
    padding = renderer.theme.components.panel_padding

    assert box.left == Inches(1.0 + padding)
    assert box.width == Inches(4.0 - (padding * 2))


def test_collect_missing_assets_reports_missing_image() -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Deck", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "image_text",
                    "title": "Image",
                    "body": "Body",
                    "image_path": "missing-image.png",
                }
            ],
        }
    )
    renderer = PresentationRenderer(asset_root="examples")

    missing_assets = renderer.collect_missing_assets(spec)

    assert missing_assets == ["slide 01 (Image): missing asset 'missing-image.png'"]


def test_collect_missing_assets_reports_missing_brand_logo() -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {
                "title": "Deck",
                "theme": "executive_premium_minimal",
                "logo_path": "missing-logo.png",
            },
            "slides": [{"type": "title", "title": "Hello"}],
        }
    )
    renderer = PresentationRenderer(asset_root="examples")

    missing_assets = renderer.collect_missing_assets(spec)

    assert missing_assets == ["presentation branding: missing asset 'missing-logo.png'"]


def test_fit_text_frame_reduces_font_size_for_tight_box() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    renderer = PresentationRenderer(asset_root="examples")

    shape = renderer.textbox(slide, 1.0, 1.0, 2.4, 0.45)
    renderer.write_paragraph(
        shape.text_frame,
        "This is a deliberately long title that should shrink to fit the box.",
        size=28,
        color=renderer.theme.colors.navy,
        bold=True,
    )
    renderer.fit_text_frame(shape.text_frame, max_size=28, bold=True)

    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.font.size is not None
    assert run.font.size.pt <= 28


def test_rebalance_flexible_regions_applies_stronger_second_pass_for_heavy_imbalance() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    regions = renderer.rebalance_flexible_regions(
        [
            {"kind": "light", "min_height": 0.6, "flex": 1.0, "content_weight": 1.0},
            {"kind": "heavy", "min_height": 0.6, "flex": 1.0, "content_weight": 6.0},
        ],
        min_flex=0.9,
        max_flex=1.35,
    )

    assert float(regions[1]["flex"]) > 1.35


def test_describe_visual_placeholder_differentiates_screenshot_and_diagram() -> None:
    renderer = PresentationRenderer(asset_root="examples")
    screenshot_spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Deck", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "image_text",
                    "title": "Product dashboard screenshot",
                    "body": "UI walkthrough",
                }
            ],
        }
    ).slides[0]
    diagram_spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Deck", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "image_text",
                    "title": "Workflow diagram",
                    "body": "Process architecture",
                }
            ],
        }
    ).slides[0]

    screenshot_copy = renderer.describe_visual_placeholder(screenshot_spec)
    diagram_copy = renderer.describe_visual_placeholder(diagram_spec)

    assert screenshot_copy["kind"] == "screenshot"
    assert diagram_copy["kind"] == "diagram"


def test_stack_vertical_regions_distributes_flexible_space() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    regions = renderer.stack_vertical_regions(
        top=1.0,
        height=2.0,
        regions=[
            {"kind": "title", "height": 0.3},
            {"kind": "body", "min_height": 0.4, "flex": 1.0},
            {"kind": "footer", "height": 0.2},
        ],
        gap=0.1,
    )

    assert len(regions) == 3
    body_bounds = regions[1][1]
    assert body_bounds[1] > 0.4


def test_stack_horizontal_regions_distributes_flexible_space() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    regions = renderer.stack_horizontal_regions(
        left=1.0,
        width=4.0,
        regions=[
            {"kind": "left", "width": 0.8},
            {"kind": "center", "min_width": 0.8, "flex": 1.0},
            {"kind": "right", "width": 0.8},
        ],
        gap=0.1,
    )

    assert len(regions) == 3
    center_bounds = regions[1][1]
    assert center_bounds[1] > 0.8


def test_build_grid_bounds_returns_expected_matrix() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    grid = renderer.build_grid_bounds(
        left=1.0,
        top=2.0,
        width=4.0,
        height=2.0,
        column_regions=[
            {"kind": "left", "min_width": 1.0, "flex": 1.0},
            {"kind": "right", "min_width": 1.0, "flex": 1.0},
        ],
        row_regions=[
            {"kind": "top", "min_height": 0.6, "flex": 1.0},
            {"kind": "bottom", "min_height": 0.6, "flex": 1.0},
        ],
        column_gap=0.1,
        row_gap=0.2,
    )

    assert len(grid) == 2
    assert len(grid[0]) == 2
    left, top, width, height = grid[0][0]
    assert left >= 1.0
    assert top >= 2.0
    assert width > 1.0
    assert height > 0.6


def test_build_columns_returns_expected_bounds() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    columns = renderer.build_columns(left=1.0, width=4.0, gap=0.1, count=3, min_width=0.8)

    assert len(columns) == 3
    assert columns[1][1] >= 0.8


def test_build_weighted_columns_allocates_more_width_to_heavier_content() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    columns = renderer.build_weighted_columns(
        left=1.0,
        width=4.0,
        gap=0.1,
        weights=[1.0, 3.0],
        min_width=1.0,
    )

    assert columns[1][1] > columns[0][1]


def test_build_constrained_columns_respects_fixed_sidebar_width() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    columns = renderer.build_constrained_columns(
        left=1.0,
        width=10.0,
        gap=0.2,
        regions=[
            {"kind": "main", "min_width": 5.0, "target_share": 3.0},
            {"kind": "sidebar", "width": 3.0, "min_width": 3.0},
        ],
    )

    assert len(columns) == 2
    assert columns[1][1] == pytest.approx(3.0)
    assert columns[0][1] == pytest.approx(6.8)


def test_build_panel_row_bounds_returns_rectangles() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    bounds = renderer.build_panel_row_bounds(left=1.0, top=2.0, width=4.0, height=1.5, gap=0.1, count=2, min_width=1.2)

    assert len(bounds) == 2
    assert bounds[0][1] == 2.0
    assert bounds[0][3] == 1.5


def test_build_panel_grid_creates_matrix_from_counts() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    grid = renderer.build_panel_grid(
        left=1.0,
        top=2.0,
        width=4.0,
        height=2.0,
        column_gap=0.1,
        row_gap=0.2,
        column_count=2,
        row_count=2,
        column_min_width=1.0,
        row_min_height=0.6,
    )

    assert len(grid) == 2
    assert len(grid[0]) == 2


def test_build_weighted_rows_allocates_more_height_to_heavier_content() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    rows = renderer.build_weighted_rows(
        top=1.0,
        height=3.0,
        gap=0.1,
        weights=[1.0, 3.0],
        min_height=0.6,
    )

    assert rows[1][1] > rows[0][1]


def test_build_constrained_rows_respects_max_height_caps() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    rows = renderer.build_constrained_rows(
        top=1.0,
        height=4.0,
        gap=0.1,
        regions=[
            {"kind": "intro", "min_height": 0.6, "target_share": 1.0, "max_height": 1.0},
            {"kind": "chart", "min_height": 1.5, "target_share": 3.0},
        ],
    )

    assert rows[0][1] <= 1.0
    assert rows[1][1] > rows[0][1]


def test_build_weighted_panel_grid_allocates_more_space_to_heavier_regions() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    grid = renderer.build_weighted_panel_grid(
        left=1.0,
        top=2.0,
        width=5.0,
        height=3.0,
        column_gap=0.1,
        row_gap=0.2,
        column_weights=[1.0, 2.5],
        row_weights=[1.0, 2.0],
        column_min_width=1.0,
        row_min_height=0.8,
    )

    assert len(grid) == 2
    assert len(grid[0]) == 2
    assert grid[0][1][2] > grid[0][0][2]
    assert grid[1][0][3] > grid[0][0][3]


def test_build_weighted_panel_row_content_bounds_returns_panel_and_inner_bounds() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    bounds = renderer.build_weighted_panel_row_content_bounds(
        left=1.0,
        top=2.0,
        width=5.0,
        height=1.8,
        gap=0.1,
        weights=[1.0, 2.0],
        min_width=1.2,
        padding=0.2,
    )

    assert len(bounds) == 2
    first_panel, first_content = bounds[0]
    assert first_panel[0] == 1.0
    assert first_content[0] == pytest.approx(first_panel[0] + 0.2)


def test_build_panel_grid_content_bounds_returns_panel_and_inner_bounds() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    grid = renderer.build_panel_grid_content_bounds(
        left=1.0,
        top=2.0,
        width=4.0,
        height=2.0,
        column_gap=0.1,
        row_gap=0.1,
        column_count=2,
        row_count=2,
        column_min_width=1.0,
        row_min_height=0.6,
        padding=0.12,
    )

    assert len(grid) == 2
    panel_bounds, content_bounds = grid[0][0]
    assert content_bounds[0] == pytest.approx(panel_bounds[0] + 0.12)


def test_build_panel_content_stack_bounds_uses_inner_padding() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    regions = renderer.build_panel_content_stack_bounds(
        left=1.0,
        top=2.0,
        width=4.0,
        height=2.0,
        padding=0.2,
        regions=[
            {"kind": "title", "height": 0.3},
            {"kind": "body", "min_height": 0.6, "flex": 1.0, "content_weight": 2.0},
        ],
        gap=0.1,
    )

    assert len(regions) == 2
    assert regions[0][1][0] == 1.2
    assert regions[0][1][2] == 3.6


def test_build_constrained_panel_content_stack_bounds_uses_target_share_and_padding() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    regions = renderer.build_constrained_panel_content_stack_bounds(
        left=1.0,
        top=2.0,
        width=4.0,
        height=2.4,
        padding=0.2,
        regions=[
            {"kind": "intro", "min_height": 0.4, "target_share": 1.0, "max_height": 0.8},
            {"kind": "body", "min_height": 0.9, "target_share": 3.0},
        ],
        gap=0.1,
    )

    assert len(regions) == 2
    assert regions[0][1][0] == pytest.approx(1.2)
    assert regions[0][1][2] == pytest.approx(3.6)
    assert regions[0][1][3] <= 0.8


def test_build_constrained_panel_grid_content_bounds_returns_matrix_with_inner_bounds() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    grid = renderer.build_constrained_panel_grid_content_bounds(
        left=1.0,
        top=2.0,
        width=5.2,
        height=3.0,
        column_gap=0.1,
        row_gap=0.2,
        column_regions=[
            {"kind": "left", "min_width": 1.4, "target_share": 1.0, "max_width": 1.8},
            {"kind": "right", "min_width": 1.8, "target_share": 2.0},
        ],
        row_regions=[
            {"kind": "top", "min_height": 0.8, "target_share": 1.0, "max_height": 1.0},
            {"kind": "bottom", "min_height": 1.0, "target_share": 2.0},
        ],
        padding=0.15,
    )

    assert len(grid) == 2
    assert len(grid[0]) == 2
    first_panel, first_content = grid[0][0]
    assert first_content[0] == pytest.approx(first_panel[0] + 0.15)
    assert first_content[1] == pytest.approx(first_panel[1] + 0.15)
    assert first_panel[2] <= 1.8


def test_estimate_content_weight_increases_with_more_content() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    light = renderer.estimate_content_weight(title="Short")
    heavy = renderer.estimate_content_weight(
        title="Longer title for comparison",
        body="This body contains substantially more words and should therefore produce a higher estimated content weight.",
        bullets=["A long bullet explaining more context", "Another long bullet with more detail"],
    )

    assert heavy > light


def test_normalize_content_flexes_returns_scaled_values() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    flexes = renderer.normalize_content_flexes([1.0, 2.0, 3.0], min_flex=0.8, max_flex=1.4)

    assert len(flexes) == 3
    assert min(flexes) >= 0.8
    assert max(flexes) <= 1.4
    assert flexes[0] < flexes[-1]


def test_build_content_stack_allocates_more_height_to_heavier_region() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    regions = renderer.build_content_stack(
        top=1.0,
        height=3.0,
        regions=[
            {"kind": "body", "min_height": 0.6, "flex": 1.0, "content_weight": 1.0},
            {"kind": "bullets", "min_height": 0.6, "flex": 1.0, "content_weight": 3.0},
        ],
        gap=0.1,
        min_flex=0.9,
        max_flex=1.4,
    )

    assert regions[1][1][1] > regions[0][1][1]


def test_compute_cover_crop_crops_wide_images_horizontally() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    crop_left, crop_top, crop_right, crop_bottom = renderer.compute_cover_crop(
        image_width_px=400,
        image_height_px=200,
        box_width=2.0,
        box_height=2.0,
    )

    assert crop_left > 0
    assert crop_right == pytest.approx(crop_left)
    assert crop_top == 0
    assert crop_bottom == 0


def test_compute_cover_crop_crops_tall_images_vertically() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    crop_left, crop_top, crop_right, crop_bottom = renderer.compute_cover_crop(
        image_width_px=200,
        image_height_px=400,
        box_width=2.0,
        box_height=1.0,
    )

    assert crop_top > 0
    assert crop_bottom == pytest.approx(crop_top)
    assert crop_left == 0
    assert crop_right == 0


def test_compute_cover_crop_respects_horizontal_focal_point() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    crop_left, _, crop_right, _ = renderer.compute_cover_crop(
        image_width_px=400,
        image_height_px=200,
        box_width=2.0,
        box_height=2.0,
        focal_x=0.15,
    )

    assert crop_left == pytest.approx(0.0)
    assert crop_right > 0.0


def test_compute_cover_crop_respects_vertical_focal_point() -> None:
    renderer = PresentationRenderer(asset_root="examples")

    _, crop_top, _, crop_bottom = renderer.compute_cover_crop(
        image_width_px=200,
        image_height_px=400,
        box_width=2.0,
        box_height=1.0,
        focal_y=0.85,
    )

    assert crop_bottom < crop_top
    assert crop_bottom >= 0.0
    assert crop_top > 0.0


def test_add_image_cover_applies_crop_for_mismatched_aspect_ratio(tmp_path: Path) -> None:
    image_path = tmp_path / "wide.png"
    Image.new("RGB", (400, 200), (120, 140, 180)).save(image_path)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    renderer = PresentationRenderer(asset_root="examples")

    picture = renderer.add_image_cover(slide, image_path, left=1.0, top=1.0, width=2.0, height=2.0)

    assert picture.crop_left > 0
    assert picture.crop_right == pytest.approx(picture.crop_left)


def test_preview_artifact_review_flags_body_content_packed_into_corner(tmp_path: Path) -> None:
    renderer = PreviewRenderer()
    image = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), (246, 244, 240))
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 64, 64), fill=(20, 30, 40))

    preview_path = tmp_path / "corner-packed.png"
    image.save(preview_path)

    review = renderer.build_preview_artifact_review([str(preview_path)])

    assert review["corner_density_warning_count"] == 1
    assert review["body_edge_contact_count"] == 1
    assert review["slides"][0]["corner_density_warning"] is True


def test_preview_artifact_review_flags_body_content_near_footer_boundary(tmp_path: Path) -> None:
    renderer = PreviewRenderer()
    image = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), (246, 244, 240))
    draw = ImageDraw.Draw(image)
    footer_line_y = renderer._y(renderer.theme.grid.footer_line_y)
    draw.rectangle((220, footer_line_y - 20, 1060, footer_line_y - 4), fill=(20, 30, 40))

    preview_path = tmp_path / "footer-intrusion.png"
    image.save(preview_path)

    review = renderer.build_preview_artifact_review([str(preview_path)])

    assert review["footer_intrusion_count"] == 1
    assert review["slides"][0]["footer_intrusion_warning"] is True


def test_preview_manifest_ignores_unlisted_pngs_during_regression(tmp_path: Path) -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Manifest deck", "theme": "executive_premium_minimal"},
            "slides": [{"type": "title", "title": "Only slide"}],
        }
    )
    baseline_dir = tmp_path / "baseline"
    current_dir = tmp_path / "current"
    comparison_dir = tmp_path / "comparison"

    render_previews(spec, baseline_dir, basename="baseline-manifest")
    render_previews(spec, current_dir, basename="current-manifest")
    Image.new("RGB", (1280, 720), (255, 0, 0)).save(baseline_dir / "stray-extra.png")

    comparison = compare_preview_directories(current_dir, baseline_dir, comparison_dir)

    assert (baseline_dir / PREVIEW_MANIFEST_FILENAME).exists()
    assert comparison["baseline_preview_count"] == 1
    assert comparison["comparison"]["extra_baseline_count"] == 0
    assert comparison["comparison"]["diff_count"] == 0


def test_preview_require_real_previews_fails_without_office_runtime(monkeypatch, tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Require real", "theme": "executive_premium_minimal"},
            "slides": [{"type": "title", "title": "Only slide"}],
        }
    )
    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: None)

    with pytest.raises(RuntimeError, match="real preview was required"):
        render_previews(
            spec,
            tmp_path / "require-real",
            backend="auto",
            require_real_previews=True,
        )


def test_preview_regression_report_includes_guidance_and_top_regressions(tmp_path: Path) -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Regression deck", "theme": "executive_premium_minimal"},
            "slides": [{"type": "title", "title": "Original slide"}],
        }
    )
    baseline_dir = tmp_path / "baseline-guidance"
    current_dir = tmp_path / "current-guidance"
    report_dir = tmp_path / "report-guidance"

    render_previews(spec, baseline_dir, basename="baseline-guidance")
    current_result = render_previews(spec, current_dir, basename="current-guidance")
    changed_preview = Path(current_result["previews"][0])
    Image.new("RGB", (1280, 720), (0, 0, 0)).save(changed_preview)

    comparison = compare_preview_directories(current_dir, baseline_dir, report_dir, write_diff_images=True)
    regression = comparison["comparison"]

    assert regression["status"] == "review"
    assert regression["diff_count"] == 1
    assert regression["top_regressions"]
    assert regression["top_regressions"][0]["slide_number"] == 1
    assert regression["guidance"]
    assert visual_regression_has_failures(regression) is True
    assert "failed visual regression" in format_visual_regression_failure(regression, context="preview regression")


def test_promote_preview_baseline_copies_manifest_and_preview_files(tmp_path: Path) -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Promote deck", "theme": "executive_premium_minimal"},
            "slides": [{"type": "title", "title": "Only slide"}],
        }
    )
    source_dir = tmp_path / "source-baseline"
    baseline_dir = tmp_path / "promoted-baseline"

    source_result = render_previews(spec, source_dir, basename="promote-source")
    promotion = promote_preview_baseline(source_dir, baseline_dir)

    assert promotion["mode"] == "promote-baseline"
    assert promotion["preview_count"] == len(source_result["previews"])
    assert Path(promotion["preview_manifest"]).exists()
    assert len(list(baseline_dir.glob("*.png"))) >= len(source_result["previews"])
    assert promotion["copied_thumbnail_sheets"]


def test_preview_image_text_respects_focal_point_when_cover_cropping(tmp_path: Path) -> None:
    asset_path = tmp_path / "focus.png"
    image = Image.new("RGB", (400, 200), (0, 0, 255))
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 199, 199), fill=(255, 0, 0))
    image.save(asset_path)

    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Deck", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "image_text",
                    "title": "Focused image",
                    "body": "Body",
                    "image_path": str(asset_path),
                    "image_focal_x": 0.15,
                }
            ],
        }
    )

    renderer = PreviewRenderer(asset_root=tmp_path)
    rendered = renderer.render_slide(spec.presentation, spec.slides[0], 1, 1)

    sampled = rendered.getpixel((974, 370))
    assert sampled[0] > sampled[2]


def test_preview_title_hero_cover_respects_focal_point_when_image_present(tmp_path: Path) -> None:
    asset_path = tmp_path / "title_focus.png"
    image = Image.new("RGB", (400, 200), (0, 0, 255))
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 199, 199), fill=(255, 0, 0))
    image.save(asset_path)

    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Deck", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "title",
                    "title": "Focused cover",
                    "layout_variant": "hero_cover",
                    "image_path": str(asset_path),
                    "image_focal_x": 0.15,
                }
            ],
        }
    )

    renderer = PreviewRenderer(asset_root=tmp_path)
    rendered = renderer.render_slide(spec.presentation, spec.slides[0], 1, 1)

    sampled = rendered.getpixel((960, 230))
    assert sampled[0] > sampled[2]


def test_preview_chart_supports_negative_values_without_crashing() -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Negative chart", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "chart",
                    "title": "Signal comparison",
                    "layout_variant": "bar",
                    "chart_categories": ["A", "B", "C"],
                    "chart_series": [
                        {"name": "Signal", "values": [25.0, -30.0, 20.0]},
                    ],
                }
            ],
        }
    )

    renderer = PreviewRenderer(asset_root="examples")
    rendered = renderer.render_slide(spec.presentation, spec.slides[0], 1, 1)

    assert rendered.size == (PREVIEW_WIDTH, PREVIEW_HEIGHT)


def test_title_layout_with_cover_image_renders_without_crashing(tmp_path: Path) -> None:
    asset_path = tmp_path / "title_cover.png"
    Image.new("RGB", (400, 200), (120, 140, 180)).save(asset_path)
    spec = PresentationInput.model_validate(
        {
            "presentation": {"title": "Deck", "theme": "executive_premium_minimal"},
            "slides": [
                {
                    "type": "title",
                    "title": "Cover with image",
                    "layout_variant": "hero_cover",
                    "image_path": str(asset_path),
                }
            ],
        }
    )
    output = tmp_path / "title-cover-image.pptx"

    renderer = PresentationRenderer(asset_root=tmp_path)
    rendered = renderer.render(spec, output)

    assert rendered.exists()

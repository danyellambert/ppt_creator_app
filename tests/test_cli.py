from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation as PptxPresentation

from ppt_creator.cli import main
from ppt_creator.renderer import PresentationRenderer
from ppt_creator.schema import PresentationInput


def test_cli_render_generates_file(tmp_path: Path) -> None:
    output = tmp_path / "cli_output.pptx"
    result = main(["render", "examples/ai_sales.json", str(output)])
    assert result == 0
    assert output.exists()


def test_cli_returns_error_for_missing_input(capsys) -> None:
    result = main(["validate", "examples/missing.json"])
    captured = capsys.readouterr()

    assert result == 2
    assert "Input JSON not found" in captured.err


def test_cli_returns_error_for_invalid_output_extension(tmp_path: Path, capsys) -> None:
    output = tmp_path / "cli_output.txt"
    result = main(["render", "examples/ai_sales.json", str(output)])
    captured = capsys.readouterr()

    assert result == 2
    assert "must end with .pptx" in captured.err


def test_cli_validate_second_example() -> None:
    result = main(["validate", "examples/product_strategy.json"])
    assert result == 0


def test_cli_review_generates_report(tmp_path: Path, capsys) -> None:
    report_path = tmp_path / "review_report.json"
    result = main(["review", "examples/ai_sales.json", "--report-json", str(report_path)])
    captured = capsys.readouterr()

    assert result == 0
    assert report_path.exists()
    assert "Review completed" in captured.out
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert "severity_counts" in payload
    assert "overflow_risk_count" in payload
    assert "balance_warning_count" in payload


def test_cli_review_can_attach_preview_result_with_real_artifact_preference(tmp_path: Path, monkeypatch) -> None:
    from ppt_creator import cli as cli_module

    report_path = tmp_path / "review_with_preview_report.json"
    preview_dir = tmp_path / "review_with_preview_artifacts"

    monkeypatch.setattr(
        cli_module,
        "render_previews_for_rendered_artifact",
        lambda *args, **kwargs: (
            {
                "mode": "preview-pptx",
                "preview_count": 10,
                "previews": [],
                "thumbnail_sheet": str(preview_dir / "thumbs.png"),
                "preview_artifact_review": {"status": "ok"},
                "visual_regression": None,
                "backend_used": "office",
            },
            "rendered_pptx",
        ),
    )

    result = main(
        [
            "review",
            "examples/ai_sales.json",
            "--preview-dir",
            str(preview_dir),
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["preview_source"] == "rendered_pptx"
    assert payload["preview_result"]["mode"] == "preview-pptx"


def test_cli_validate_emits_informational_logs(capsys) -> None:
    result = main(["validate", "examples/product_strategy.json"])
    captured = capsys.readouterr()

    assert result == 0
    assert "[INFO] Loading input:" in captured.out
    assert "[OK] Valid JSON:" in captured.out


def test_cli_render_dry_run_does_not_create_file(tmp_path: Path, capsys) -> None:
    output = tmp_path / "dry_run_output.pptx"
    result = main(["render", "examples/ai_sales.json", str(output), "--dry-run"])
    captured = capsys.readouterr()

    assert result == 0
    assert not output.exists()
    assert "Dry run" in captured.out


def test_cli_render_dry_run_can_include_quality_review(tmp_path: Path) -> None:
    output = tmp_path / "dry_run_review_output.pptx"
    report_path = tmp_path / "dry_run_review_report.json"

    result = main(
        [
            "render",
            "examples/ai_sales.json",
            str(output),
            "--dry-run",
            "--review",
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["quality_review"] is not None
    assert "severity_counts" in payload["quality_review"]
    assert "overflow_risk_count" in payload["quality_review"]
    assert "clipping_risk_count" in payload["quality_review"]
    assert "collision_risk_count" in payload["quality_review"]
    assert "top_risk_slides" in payload["quality_review"]


def test_cli_render_can_generate_preview_report_from_rendered_pptx(tmp_path: Path, monkeypatch) -> None:
    from ppt_creator import preview as preview_module

    output = tmp_path / "render_with_preview.pptx"
    preview_dir = tmp_path / "render_previews"
    preview_report_path = tmp_path / "render_preview_report.json"
    report_path = tmp_path / "render_report.json"

    def _fake_run(command, capture_output, text, check):
        outdir = Path(command[command.index("--outdir") + 1])
        source_pptx = Path(command[-1])
        slide_count = len(PptxPresentation(str(source_pptx)).slides)
        outdir.mkdir(parents=True, exist_ok=True)
        for index in range(1, slide_count + 1):
            Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"render-preview-{index:02d}.png")

        class _Completed:
            returncode = 0
            stdout = ""
            stderr = ""

        return _Completed()

    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: "/usr/bin/soffice")
    monkeypatch.setattr(preview_module.subprocess, "run", _fake_run)

    result = main(
        [
            "render",
            "examples/ai_sales.json",
            str(output),
            "--preview-dir",
            str(preview_dir),
            "--preview-report-json",
            str(preview_report_path),
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    preview_payload = json.loads(preview_report_path.read_text(encoding="utf-8"))
    assert payload["rendered"] is True
    assert payload["preview_output_dir"] == str(preview_dir)
    assert payload["preview_source"] == "rendered_pptx"
    assert preview_payload["mode"] == "preview-pptx"


def test_cli_render_batch_review_report_includes_aggregate_risk_fields(tmp_path: Path) -> None:
    output_dir = tmp_path / "batch_review_output"
    report_path = tmp_path / "batch_review_report.json"

    result = main(
        [
            "render-batch",
            "examples",
            str(output_dir),
            "--pattern",
            "product_strategy.json",
            "--dry-run",
            "--review",
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["review_enabled"] is True
    assert "review_overflow_risk_count" in payload
    assert "review_clipping_risk_count" in payload
    assert "review_collision_risk_count" in payload
    assert "review_balance_warning_count" in payload


def test_cli_render_batch_generates_output_and_report(tmp_path: Path) -> None:
    output_dir = tmp_path / "batch_output"
    report_path = tmp_path / "batch_report.json"

    result = main(
        [
            "render-batch",
            "examples",
            str(output_dir),
            "--pattern",
            "product_strategy.json",
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    assert (output_dir / "product_strategy.pptx").exists()
    assert report_path.exists()


def test_cli_render_dry_run_accepts_theme_color_overrides(tmp_path: Path, capsys) -> None:
    output = tmp_path / "brand_override.pptx"
    result = main(
        [
            "render",
            "examples/ai_sales.json",
            str(output),
            "--dry-run",
            "--theme",
            "dark_boardroom",
            "--primary-color",
            "112233",
            "--secondary-color",
            "AABBCC",
        ]
    )
    captured = capsys.readouterr()

    assert result == 0
    assert not output.exists()
    assert "Dry run" in captured.out


def test_cli_template_generates_domain_json(tmp_path: Path, capsys) -> None:
    output = tmp_path / "sales_template.json"
    result = main(["template", "sales", str(output)])
    captured = capsys.readouterr()

    assert result == 0
    assert output.exists()
    assert "Generated template" in captured.out

    spec = PresentationInput.from_path(output)
    assert spec.presentation.title == "Sales operating review"
    assert len(spec.slides) >= 4


def test_cli_template_accepts_theme_override(tmp_path: Path) -> None:
    output = tmp_path / "strategy_template.json"
    result = main(["template", "strategy", str(output), "--theme", "consulting_clean"])

    assert result == 0
    spec = PresentationInput.from_path(output)
    assert spec.presentation.theme == "consulting_clean"


def test_cli_template_accepts_audience_profile(tmp_path: Path) -> None:
    output = tmp_path / "board_template.json"
    result = main(["template", "sales", str(output), "--audience-profile", "board"])

    assert result == 0
    spec = PresentationInput.from_path(output)
    assert spec.presentation.footer_text == "Board profile"


def test_cli_profiles_and_assets_commands_emit_reports(tmp_path: Path) -> None:
    profiles_report = tmp_path / "profiles.json"
    assets_report = tmp_path / "assets.json"

    assert main(["profiles", "--report-json", str(profiles_report)]) == 0
    assert main(["assets", "--report-json", str(assets_report)]) == 0

    profiles_payload = json.loads(profiles_report.read_text(encoding="utf-8"))
    assets_payload = json.loads(assets_report.read_text(encoding="utf-8"))
    assert profiles_payload["profiles"]
    assert assets_payload["collections"]


def test_cli_workflows_and_workflow_template_commands_emit_reports(tmp_path: Path) -> None:
    workflows_report = tmp_path / "workflows.json"
    workflow_template_output = tmp_path / "sales_qbr_template.json"

    assert main(["workflows", "--report-json", str(workflows_report)]) == 0
    assert main(["workflow-template", "sales_qbr", str(workflow_template_output)]) == 0

    workflows_payload = json.loads(workflows_report.read_text(encoding="utf-8"))
    assert workflows_payload["workflows"]

    spec = PresentationInput.from_path(workflow_template_output)
    assert spec.presentation.footer_text == "Sales profile"
    assert spec.presentation.title == "Sales operating review"


def test_cli_preview_generates_pngs_and_thumbnail_sheet(tmp_path: Path, capsys) -> None:
    output_dir = tmp_path / "preview_output"
    report_path = tmp_path / "preview_report.json"

    result = main(
        [
            "preview",
            "examples/ai_sales.json",
            str(output_dir),
            "--basename",
            "ai-sales-preview",
            "--report-json",
            str(report_path),
        ]
    )
    captured = capsys.readouterr()

    assert result == 0
    assert report_path.exists()
    assert "Generated previews" in captured.out
    report = report_path.read_text(encoding="utf-8")
    assert "quality_review" in report
    assert "severity_counts" in report
    assert "top_risk_slides" in report
    assert "preview_artifact_review" in report
    generated_pngs = sorted(output_dir.glob("*.png"))
    assert len(generated_pngs) == 11
    assert any(path.name.endswith("-thumbnails.png") for path in generated_pngs)


def test_cli_preview_accepts_debug_overlays(tmp_path: Path) -> None:
    output_dir = tmp_path / "preview_debug_output"
    result = main(
        [
            "preview",
            "examples/ai_sales.json",
            str(output_dir),
            "--debug-grid",
            "--debug-safe-areas",
        ]
    )

    assert result == 0
    assert sorted(output_dir.glob("*.png"))


def test_cli_preview_auto_backend_reports_synthetic_fallback(tmp_path: Path, monkeypatch) -> None:
    from ppt_creator import preview as preview_module

    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: None)
    output_dir = tmp_path / "preview_auto_output"
    report_path = tmp_path / "preview_auto_report.json"

    result = main(
        [
            "preview",
            "examples/ai_sales.json",
            str(output_dir),
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["backend_requested"] == "auto"
    assert payload["backend_used"] == "synthetic"


def test_cli_preview_supports_visual_regression_against_baseline(tmp_path: Path) -> None:
    baseline_dir = tmp_path / "baseline_previews"
    current_dir = tmp_path / "current_previews"
    report_path = tmp_path / "regression_report.json"

    baseline_result = main(
        [
            "preview",
            "examples/ai_sales.json",
            str(baseline_dir),
            "--basename",
            "baseline-deck",
        ]
    )
    assert baseline_result == 0

    result = main(
        [
            "preview",
            "examples/ai_sales.json",
            str(current_dir),
            "--basename",
            "current-deck",
            "--baseline-dir",
            str(baseline_dir),
            "--write-diff-images",
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["visual_regression"] is not None
    assert payload["visual_regression"]["status"] == "ok"
    assert payload["visual_regression"]["diff_count"] == 0
    assert payload["visual_regression"]["compared_preview_count"] == 10
    assert len(payload["visual_regression"]["diff_images"]) == 10


def test_cli_preview_pptx_generates_pngs_via_office_backend_mock(tmp_path: Path, monkeypatch) -> None:
    from ppt_creator import preview as preview_module

    pptx_path = tmp_path / "source_deck.pptx"
    output_dir = tmp_path / "pptx_preview_output"
    report_path = tmp_path / "pptx_preview_report.json"

    spec = PresentationInput.from_path("examples/ai_sales.json")
    PresentationRenderer(asset_root="examples").render(spec, pptx_path)

    def _fake_run(command, capture_output, text, check):
        outdir = Path(command[command.index("--outdir") + 1])
        source_pptx = Path(command[-1])
        slide_count = len(PptxPresentation(str(source_pptx)).slides)
        outdir.mkdir(parents=True, exist_ok=True)
        for index in range(1, slide_count + 1):
            Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"mock-{index:02d}.png")

        class _Completed:
            returncode = 0
            stdout = ""
            stderr = ""

        return _Completed()

    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: "/usr/bin/soffice")
    monkeypatch.setattr(preview_module.subprocess, "run", _fake_run)

    result = main(
        [
            "preview-pptx",
            str(pptx_path),
            str(output_dir),
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["mode"] == "preview-pptx"
    assert payload["preview_count"] == 10
    assert payload["backend_used"] == "office"
    assert payload["preview_artifact_review"]["status"] in {"ok", "review"}


def test_cli_preview_pptx_falls_back_to_pdf_rasterization_when_office_exports_single_png(
    tmp_path: Path,
    monkeypatch,
) -> None:
    from ppt_creator import preview as preview_module

    pptx_path = tmp_path / "source_deck_pdf_fallback.pptx"
    output_dir = tmp_path / "pptx_preview_pdf_fallback_output"
    report_path = tmp_path / "pptx_preview_pdf_fallback_report.json"

    spec = PresentationInput.from_path("examples/ai_sales.json")
    PresentationRenderer(asset_root="examples").render(spec, pptx_path)

    def _fake_run(command, capture_output, text, check):
        outdir = Path(command[command.index("--outdir") + 1]) if "--outdir" in command else tmp_path
        outdir.mkdir(parents=True, exist_ok=True)
        if command[0] == "/usr/bin/soffice" and "png" in command:
            Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / "single.png")
        elif command[0] == "/usr/bin/soffice" and "pdf" in command:
            (outdir / "source_deck_pdf_fallback.pdf").write_bytes(b"%PDF-1.4 mock")
        elif command[0] == "/usr/bin/gs":
            slide_count = len(PptxPresentation(str(pptx_path)).slides)
            pattern = next(argument.split("=", 1)[1] for argument in command if argument.startswith("-sOutputFile="))
            for index in range(1, slide_count + 1):
                target = Path(pattern.replace("%02d", f"{index:02d}"))
                Image.new("RGB", (1280, 720), (245, 245, 245)).save(target)

        class _Completed:
            returncode = 0
            stdout = ""
            stderr = ""

        return _Completed()

    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: "/usr/bin/soffice")
    monkeypatch.setattr(preview_module, "find_ghostscript_runtime", lambda: "/usr/bin/gs")
    monkeypatch.setattr(preview_module.subprocess, "run", _fake_run)

    result = main(
        [
            "preview-pptx",
            str(pptx_path),
            str(output_dir),
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["preview_count"] == 10
    assert payload["backend_used"] == "office"
    assert payload["office_conversion_strategy"] == "pdf_via_ghostscript"


def test_cli_compare_pptx_generates_visual_comparison_report(tmp_path: Path, monkeypatch) -> None:
    from ppt_creator import preview as preview_module

    before_pptx = tmp_path / "before_deck.pptx"
    after_pptx = tmp_path / "after_deck.pptx"
    output_dir = tmp_path / "compare_pptx_output"
    report_path = tmp_path / "compare_pptx_report.json"

    spec = PresentationInput.from_path("examples/ai_sales.json")
    PresentationRenderer(asset_root="examples").render(spec, before_pptx)
    PresentationRenderer(asset_root="examples").render(spec, after_pptx)

    def _fake_run(command, capture_output, text, check):
        outdir = Path(command[command.index("--outdir") + 1]) if "--outdir" in command else tmp_path
        source_pptx = Path(command[-1])
        slide_count = len(PptxPresentation(str(source_pptx)).slides)
        outdir.mkdir(parents=True, exist_ok=True)
        for index in range(1, slide_count + 1):
            Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"compare-mock-{index:02d}.png")

        class _Completed:
            returncode = 0
            stdout = ""
            stderr = ""

        return _Completed()

    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: "/usr/bin/soffice")
    monkeypatch.setattr(preview_module.subprocess, "run", _fake_run)

    result = main(
        [
            "compare-pptx",
            str(before_pptx),
            str(after_pptx),
            str(output_dir),
            "--write-diff-images",
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["mode"] == "compare-pptx"
    assert payload["comparison"]["status"] == "ok"
    assert payload["comparison"]["diff_count"] == 0


def test_cli_review_pptx_generates_real_artifact_review_report(tmp_path: Path, monkeypatch) -> None:
    from ppt_creator import preview as preview_module

    pptx_path = tmp_path / "review_source_deck.pptx"
    output_dir = tmp_path / "review_pptx_output"
    report_path = tmp_path / "review_pptx_report.json"

    spec = PresentationInput.from_path("examples/ai_sales.json")
    PresentationRenderer(asset_root="examples").render(spec, pptx_path)

    def _fake_run(command, capture_output, text, check):
        outdir = Path(command[command.index("--outdir") + 1])
        source_pptx = Path(command[-1])
        slide_count = len(PptxPresentation(str(source_pptx)).slides)
        outdir.mkdir(parents=True, exist_ok=True)
        for index in range(1, slide_count + 1):
            Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"review-mock-{index:02d}.png")

        class _Completed:
            returncode = 0
            stdout = ""
            stderr = ""

        return _Completed()

    monkeypatch.setattr(preview_module, "find_office_runtime", lambda: "/usr/bin/soffice")
    monkeypatch.setattr(preview_module.subprocess, "run", _fake_run)

    result = main(
        [
            "review-pptx",
            str(pptx_path),
            str(output_dir),
            "--report-json",
            str(report_path),
        ]
    )

    assert result == 0
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["mode"] == "review-pptx"
    assert payload["preview_result"]["mode"] == "preview-pptx"
    assert "preview_artifact_review" in payload

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ppt_creator.renderer import PresentationRenderer
from ppt_creator.schema import PresentationInput


def build_layout_smoke_spec() -> PresentationInput:
    return PresentationInput.model_validate(
        {
            "presentation": {
                "title": "Layout Smoke Test",
                "theme": "executive_premium_minimal",
            },
            "slides": [
                {
                    "type": "title",
                    "title": "Title Slide",
                    "subtitle": "Smoke coverage",
                    "layout_variant": "hero_cover",
                },
                {
                    "type": "section",
                    "title": "Section Slide",
                    "section_label": "Section",
                },
                {
                    "type": "agenda",
                    "title": "Agenda Slide",
                    "body": "Today we focus on the few decisions that matter most.",
                    "bullets": ["Context", "Options", "Risks", "Decision"],
                },
                {
                    "type": "bullets",
                    "title": "Bullets Slide",
                    "bullets": ["One", "Two", "Three"],
                },
                {
                    "type": "cards",
                    "title": "Cards Slide",
                    "cards": [
                        {"title": "Card 1", "body": "Body 1"},
                        {"title": "Card 2", "body": "Body 2"},
                        {"title": "Card 3", "body": "Body 3"},
                    ],
                },
                {
                    "type": "metrics",
                    "title": "Metrics Slide",
                    "metrics": [
                        {"value": "10%", "label": "lift"},
                        {"value": "2x", "label": "speed"},
                    ],
                },
                {
                    "type": "chart",
                    "title": "Chart Slide",
                    "chart_categories": ["Q1", "Q2", "Q3"],
                    "chart_series": [{"name": "Revenue", "values": [2.1, 3.0, 4.4]}],
                },
                {
                    "type": "image_text",
                    "title": "Image Slide",
                    "body": "Body copy",
                    "image_path": "missing-image.png",
                    "image_caption": "Placeholder caption",
                },
                {
                    "type": "timeline",
                    "title": "Timeline Slide",
                    "timeline_items": [
                        {"title": "Discover", "body": "Frame the problem", "tag": "Week 1"},
                        {"title": "Pilot", "body": "Run the first workflow", "tag": "Week 2"},
                        {"title": "Scale", "body": "Operationalize the rollout", "tag": "Week 3"},
                    ],
                },
                {
                    "type": "comparison",
                    "title": "Comparison Slide",
                    "comparison_columns": [
                        {
                            "title": "Current state",
                            "body": "Inconsistent execution and fragmented tooling.",
                        },
                        {
                            "title": "Target state",
                            "bullets": ["Clear workflow", "Structured outputs", "Measurable lift"],
                        },
                    ],
                },
                {
                    "type": "two_column",
                    "title": "Two Column Slide",
                    "two_column_columns": [
                        {"title": "Current narrative", "body": "Too many competing priorities obscure the core decision."},
                        {"title": "Target narrative", "bullets": ["Sequence the work", "Clarify ownership", "Measure outcomes"]},
                    ],
                },
                {
                    "type": "table",
                    "title": "Table Slide",
                    "table_columns": ["Area", "Status", "Action"],
                    "table_rows": [
                        ["Pipeline", "Stable", "Maintain cadence"],
                        ["Enablement", "Lagging", "Refresh assets"],
                    ],
                },
                {
                    "type": "faq",
                    "title": "FAQ Slide",
                    "faq_items": [
                        {"title": "What changes first?", "body": "Start with the highest-friction workflow."},
                        {"title": "How do we measure success?", "body": "Track adoption and operational lift."},
                    ],
                },
                {
                    "type": "summary",
                    "title": "Summary Slide",
                    "body": "The overall recommendation is to narrow the workflow, measure impact, and scale only after the process is stable.",
                    "bullets": ["Keep scope tight", "Measure adoption", "Scale with discipline"],
                },
                {
                    "type": "closing",
                    "title": "Closing Slide",
                    "quote": "Stay structured.",
                },
            ],
        }
    )


def test_all_layouts_render_without_crashing(tmp_path: Path) -> None:
    spec = build_layout_smoke_spec()
    output = tmp_path / "layouts-smoke.pptx"

    renderer = PresentationRenderer(asset_root="examples")
    rendered = renderer.render(spec, output)

    assert rendered.exists()
    presentation = Presentation(str(rendered))
    assert len(presentation.slides) == 15


def test_missing_image_uses_placeholder_text(tmp_path: Path) -> None:
    spec = build_layout_smoke_spec()
    output = tmp_path / "layouts-placeholder.pptx"

    renderer = PresentationRenderer(asset_root="examples")
    rendered = renderer.render(spec, output)
    presentation = Presentation(str(rendered))

    image_slide = presentation.slides[7]
    texts = [shape.text for shape in image_slide.shapes if hasattr(shape, "text")]
    joined = "\n".join(texts)

    assert "VISUAL PLACEHOLDER" in joined
    assert "Image unavailable" in joined
    assert "Missing asset: missing-image.png" in joined


def test_layout_variants_render_without_crashing(tmp_path: Path) -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {
                "title": "Variant Deck",
                "theme": "executive_premium_minimal",
            },
            "slides": [
                {
                    "type": "title",
                    "title": "Hero Cover",
                    "layout_variant": "hero_cover",
                },
                {
                    "type": "bullets",
                    "title": "Full Width Bullets",
                    "layout_variant": "full_width",
                    "bullets": ["One", "Two", "Three"],
                },
                {
                    "type": "metrics",
                    "title": "Compact KPIs",
                    "layout_variant": "compact",
                    "metrics": [
                        {"value": "10%", "label": "lift"},
                        {"value": "2x", "label": "speed"},
                        {"value": "98%", "label": "quality"},
                        {"value": "4h", "label": "saved"},
                    ],
                },
                {
                    "type": "chart",
                    "title": "Bar Chart",
                    "layout_variant": "bar",
                    "chart_categories": ["A", "B", "C"],
                    "chart_series": [{"name": "Series", "values": [3, 2, 4]}],
                },
                {
                    "type": "image_text",
                    "title": "Image Left",
                    "layout_variant": "image_left",
                    "body": "Body copy",
                    "image_path": "missing-image.png",
                },
            ],
        }
    )
    output = tmp_path / "variant-layouts.pptx"

    renderer = PresentationRenderer(asset_root="examples")
    rendered = renderer.render(spec, output)

    assert rendered.exists()
    presentation = Presentation(str(rendered))
    assert len(presentation.slides) == 5


def test_new_layout_types_render_without_crashing(tmp_path: Path) -> None:
    spec = PresentationInput.model_validate(
        {
            "presentation": {
                "title": "New Layout Deck",
                "theme": "executive_premium_minimal",
            },
            "slides": [
                {
                    "type": "timeline",
                    "title": "Rollout timeline",
                    "timeline_items": [
                        {"title": "Diagnose", "body": "Find the highest-value workflow"},
                        {"title": "Deploy", "body": "Launch a controlled pilot"},
                        {"title": "Measure", "body": "Track adoption and impact"},
                    ],
                },
                {
                    "type": "comparison",
                    "title": "Before vs after",
                    "comparison_columns": [
                        {
                            "title": "Before",
                            "bullets": ["Manual prep", "Uneven quality"],
                        },
                        {
                            "title": "After",
                            "bullets": ["Structured workflow", "Better consistency"],
                        },
                    ],
                },
                {
                    "type": "chart",
                    "title": "Trend",
                    "layout_variant": "line",
                    "chart_categories": ["Jan", "Feb", "Mar"],
                    "chart_series": [{"name": "Adoption", "values": [1, 3, 5]}],
                },
                {
                    "type": "agenda",
                    "title": "Agenda",
                    "bullets": ["Context", "Decision", "Next steps"],
                },
                {
                    "type": "two_column",
                    "title": "Narrative",
                    "two_column_columns": [
                        {"title": "Before", "body": "The story is fragmented."},
                        {"title": "After", "bullets": ["Clear message", "Clear owner"]},
                    ],
                },
                {
                    "type": "table",
                    "title": "Executive table",
                    "table_columns": ["Area", "Status"],
                    "table_rows": [["Scope", "Tight"], ["Adoption", "Measured"]],
                },
                {
                    "type": "faq",
                    "title": "FAQ",
                    "faq_items": [
                        {"title": "Why now?", "body": "The workflow pressure is already here."},
                        {"title": "What next?", "body": "Pilot, measure, then scale."},
                    ],
                },
                {
                    "type": "summary",
                    "title": "Summary",
                    "bullets": ["Stay focused", "Sequence the rollout"],
                },
            ],
        }
    )
    output = tmp_path / "new-layouts.pptx"

    renderer = PresentationRenderer(asset_root="examples")
    rendered = renderer.render(spec, output)

    assert rendered.exists()
    presentation = Presentation(str(rendered))
    assert len(presentation.slides) == 8


def test_long_content_layouts_render_with_initial_autofit(tmp_path: Path) -> None:
    long_bullet = "This bullet is intentionally long so the renderer needs to shrink the text and preserve a cleaner layout without immediately overflowing the expected content area."
    spec = PresentationInput.model_validate(
        {
            "presentation": {
                "title": "Autofit Stress Deck",
                "theme": "executive_premium_minimal",
            },
            "slides": [
                {
                    "type": "agenda",
                    "title": "A very long agenda title that should still remain usable after autofit is applied to the heading and content rows",
                    "body": "This introduction paragraph is intentionally verbose so we can validate that the first layer of autofit handles agenda narrative copy without crashing the render pipeline.",
                    "bullets": [long_bullet, long_bullet, long_bullet],
                },
                {
                    "type": "metrics",
                    "title": "Metrics under pressure",
                    "subtitle": "Labels and details here are long on purpose to exercise the new layout hardening rules.",
                    "metrics": [
                        {
                            "value": "128%",
                            "label": "Expansion revenue influenced by automation-assisted planning workflows",
                            "detail": "Very long detail text that should be compressed to fit inside the KPI card without destroying the rest of the composition.",
                            "trend": "Up meaningfully vs prior baseline",
                        },
                        {
                            "value": "4.3x",
                            "label": "Time saved in preparation-heavy executive communication loops",
                            "detail": "Another intentionally verbose detail line for layout resilience.",
                        },
                    ],
                },
                {
                    "type": "table",
                    "title": "Dense executive table",
                    "table_columns": ["Area with a longer label", "Status summary", "Next action and implication"],
                    "table_rows": [
                        [
                            "Pipeline governance and operating cadence",
                            "Stable but still verbose",
                            "Clarify owners and tighten the weekly review motion",
                        ],
                        [
                            "Enablement and training rhythm",
                            "Lagging in parts",
                            "Refresh collateral and sequence leadership communication better",
                        ],
                    ],
                },
                {
                    "type": "faq",
                    "title": "FAQ under stress",
                    "faq_items": [
                        {
                            "title": "Why should this initiative start now if the organization is already overloaded?",
                            "body": "Because the current operating friction is already imposing a tax on execution quality, and delaying a narrower rollout only preserves that cost structure.",
                        },
                        {
                            "title": "How do we ensure the first release remains constrained enough to be credible?",
                            "body": "Pick one workflow, one audience, one measurable outcome, and one governance owner before expanding scope.",
                        },
                    ],
                },
                {
                    "type": "image_text",
                    "title": "Image text stress case",
                    "body": "This narrative paragraph is intentionally long so the renderer needs to reduce text size and keep the surrounding composition more stable than before.",
                    "bullets": [long_bullet, long_bullet],
                    "image_path": "missing-image.png",
                    "image_caption": "A deliberately long placeholder caption describing the intended image and how it supports the business narrative.",
                },
                {
                    "type": "cards",
                    "title": "Cards stress case",
                    "cards": [
                        {
                            "title": "A longer card title that should remain visually stable",
                            "body": "This card body is intentionally verbose so the renderer needs to reduce type size and keep the panel composition balanced across cards.",
                            "footer": "Important footer note",
                        },
                        {
                            "title": "Second card with long narrative",
                            "body": "Another long card body used to exercise the first generation of stacked panel primitives in executive card layouts.",
                        },
                        {
                            "title": "Third card",
                            "body": "Still intentionally long for stress testing the layout behavior under denser content conditions.",
                        },
                    ],
                },
                {
                    "type": "two_column",
                    "title": "Two column stress case",
                    "two_column_columns": [
                        {
                            "title": "Current narrative under pressure",
                            "body": "This body text is deliberately long so the renderer needs to distribute vertical space more intelligently.",
                            "bullets": [long_bullet, long_bullet],
                            "footer": "Current state footer",
                        },
                        {
                            "title": "Target narrative under pressure",
                            "body": "A second long body helps validate that both columns keep a more coherent composition.",
                            "bullets": [long_bullet, long_bullet],
                            "footer": "Target state footer",
                        },
                    ],
                },
                {
                    "type": "bullets",
                    "title": "A bullets slide title that is intentionally long so the layout needs stronger autofit behavior in the heading boxes",
                    "subtitle": "A similarly verbose subtitle makes sure the manual bullets heading path also participates in the broader autofit rollout.",
                    "eyebrow": "Executive framing under pressure",
                    "layout_variant": "full_width",
                    "bullets": [long_bullet, long_bullet, long_bullet, long_bullet],
                },
                {
                    "type": "title",
                    "layout_variant": "hero_cover",
                    "title": "A longer title cover that should still stay readable after more semantic split logic and autofit are applied",
                    "subtitle": "A similarly long subtitle helps exercise the evolved title layout.",
                    "body": "This body text is intentionally verbose so the title slide also participates in the latest semantic layout rollout.",
                },
                {
                    "type": "section",
                    "title": "A section divider with a long title that should still feel visually balanced",
                    "subtitle": "Section subtitle with additional context for stress testing.",
                    "section_label": "Strategic section",
                },
                {
                    "type": "chart",
                    "title": "Chart stress case",
                    "body": "Longer narrative context above the chart validates that the chart region and intro region now rebalance more semantically.",
                    "chart_categories": ["Q1", "Q2", "Q3", "Q4", "Q5", "Q6"],
                    "chart_series": [
                        {"name": "Adoption", "values": [1, 2, 3, 5, 8, 13]},
                        {"name": "Revenue", "values": [2, 3, 5, 8, 13, 21]},
                    ],
                },
                {
                    "type": "timeline",
                    "title": "Timeline stress case",
                    "timeline_items": [
                        {
                            "title": "Diagnose the real operational bottleneck",
                            "body": "A longer body helps validate adaptive panel widths and internal vertical balancing.",
                            "tag": "Week 1",
                            "footer": "Start narrow",
                        },
                        {
                            "title": "Pilot the first constrained workflow",
                            "body": "Another intentionally verbose body ensures the panels cannot rely on rigid geometry alone.",
                            "tag": "Week 2",
                            "footer": "Measure fast",
                        },
                        {
                            "title": "Scale only after signal is clear",
                            "body": "The final step also includes more text to keep the stress test realistic.",
                            "tag": "Week 3",
                            "footer": "Codify learnings",
                        },
                    ],
                },
                {
                    "type": "closing",
                    "title": "Closing stress case",
                    "quote": "This closing quote is intentionally much longer than ideal so the shared quote block now needs to rely on stronger autofit behavior even in the final composed slide.",
                    "attribution": "PPT Creator QA",
                },
            ],
        }
    )
    output = tmp_path / "autofit-layouts.pptx"

    renderer = PresentationRenderer(asset_root="examples")
    rendered = renderer.render(spec, output)

    assert rendered.exists()
    presentation = Presentation(str(rendered))
    assert len(presentation.slides) == 13

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ppt_creator.renderer import PresentationRenderer
from ppt_creator.schema import PresentationInput


def test_all_example_json_files_validate() -> None:
    example_paths = sorted(Path("examples").glob("*.json"))
    assert example_paths

    for example_path in example_paths:
        spec = PresentationInput.from_path(example_path)
        assert spec.presentation.title
        assert spec.slides


def test_all_examples_render_successfully(tmp_path: Path) -> None:
    example_paths = sorted(Path("examples").glob("*.json"))
    renderer = PresentationRenderer(asset_root="examples")

    for example_path in example_paths:
        spec = PresentationInput.from_path(example_path)
        output_path = tmp_path / f"{example_path.stem}.pptx"
        rendered = renderer.render(spec, output_path)

        assert rendered.exists()
        presentation = Presentation(str(rendered))
        assert len(presentation.slides) == len(spec.slides)

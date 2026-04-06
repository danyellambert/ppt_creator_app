from __future__ import annotations

import json
import re
import shutil
import subprocess
from pathlib import Path
from threading import Thread
from urllib import request

from PIL import Image
from pptx import Presentation as PptxPresentation

from ppt_creator.api import _build_playground_html, build_api_server
from ppt_creator.renderer import PresentationRenderer
from ppt_creator.schema import PresentationInput


def _request_json(url: str, payload: dict[str, object] | None = None, *, method: str = "GET"):
    data = None if payload is None else json.dumps(payload).encode("utf-8")
    headers = {"Content-Type": "application/json"} if data is not None else {}
    req = request.Request(url, data=data, headers=headers, method=method)
    with request.urlopen(req, timeout=5) as response:
        return response.status, json.loads(response.read().decode("utf-8"))


def test_api_health_and_templates_endpoints(monkeypatch) -> None:
    monkeypatch.delenv("PPT_CREATOR_AI_SERVICE_PROVIDER", raising=False)
    monkeypatch.delenv("PPT_CREATOR_AI_SERVICE_MODEL", raising=False)

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        status, health_payload = _request_json(f"{base_url}/health")
        assert status == 200
        assert health_payload["status"] == "ok"

        status, templates_payload = _request_json(f"{base_url}/templates")
        assert status == 200
        assert templates_payload["domains"] == ["consulting", "product", "proposal", "sales", "strategy"]

        status, profiles_payload = _request_json(f"{base_url}/profiles")
        assert status == 200
        assert profiles_payload["profiles"]

        status, assets_payload = _request_json(f"{base_url}/assets")
        assert status == 200
        assert assets_payload["collections"]

        status, brand_packs_payload = _request_json(f"{base_url}/brand-packs")
        assert status == 200
        assert brand_packs_payload["brand_packs"]

        status, workflows_payload = _request_json(f"{base_url}/workflows")
        assert status == 200
        assert workflows_payload["workflows"]
        assert any(item["name"] == "commercial_proposal" for item in workflows_payload["workflows"])

        status, marketplace_payload = _request_json(f"{base_url}/marketplace")
        assert status == 200
        assert marketplace_payload["summary"]["workflow_count"] >= 5
        assert any(item["name"] == "commercial_proposal" for item in marketplace_payload["workflows"])
        assert any(item["type"] == "summary" for item in marketplace_payload["layouts"])

        status, providers_payload = _request_json(f"{base_url}/ai/providers")
        assert status == 200
        assert providers_payload["providers"] == ["heuristic", "local_service", "ollama_local"]
        ollama_provider = next(
            item for item in providers_payload["provider_details"] if item["name"] == "ollama_local"
        )
        assert ollama_provider["supports_model_listing"] is True

        status, ai_status_payload = _request_json(f"{base_url}/ai/status")
        assert status == 200
        assert ai_status_payload["result"]["mode"] == "ai-status"
        assert ai_status_payload["result"]["selected_provider"] == "local_service"
        assert ai_status_payload["result"]["provider_status"]["service_url"]
        assert ai_status_payload["result"]["provider_status"]["provider_name"] == "ollama"
        assert ai_status_payload["result"]["provider_status"]["provider_source"] == "app_default"
        assert ai_status_payload["result"]["provider_status"]["model_name"] == "nemotron-3-nano:30b-cloud"
        assert ai_status_payload["result"]["provider_status"]["model_source"] == "app_default"

        req = request.Request(f"{base_url}/playground", method="GET")
        with request.urlopen(req, timeout=5) as response:
            html = response.read().decode("utf-8")
        assert "PPT Creator Playground" in html
        assert "workflowPreset" in html
        assert "brandPack" in html
        assert "autoValidate" in html
        assert "autoReview" in html
        assert "compareBeforePptx" in html
        assert "Compare PPTX" in html
        assert "Review rendered PPTX" in html
        assert "Iterate flow" in html
        assert "Focus top risk slide" in html
        assert "Export current deck" in html
        assert "liveReviewSummary" in html
        assert "slideSelector" in html
        assert "Apply guided edits" in html
        assert "guidedImagePath" in html
        assert "guidedImageCaption" in html
        assert "guidedImageFocalX" in html
        assert "guidedImageFocalY" in html
        assert "Promote baseline" in html
        assert "requireRealPreviews" in html
        assert "failOnRegression" in html
        assert "aiProvider" in html
        assert "aiAuthoringMode" in html
        assert "briefingInput" in html
        assert "intentInput" in html
        assert "aiUseIntentMode" in html
        assert "Generate from AI briefing" in html
        assert "Generate + Render with AI" in html
        assert "/ai/providers" in html
        assert "/ai/models" in html
        assert "/ai/status" in html
        assert "Refresh AI status" in html
        assert "aiRuntimeStatus" in html
        assert "aiActionStatus" in html
        assert "aiModelOptions" in html
        assert "aiAutoPreview" in html
        assert "aiAutoRender" in html
        assert "themeToggleButton" in html
        assert "toggleTheme()" in html
        assert "themeStorageKey" in html
        assert "Deck studio + AI copilot" in html
        assert "Default trusted workflow" in html
        assert "Request provider:" in html
        assert "Request model:" in html
        assert "Describe the deck you want" in html
        assert "AI-first (send raw prompt to the AI author)" in html
        assert "AI request / response" in html
        assert "Prompt sent to AI" in html
        assert "Raw AI response" in html
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_exposes_ollama_model_listing(monkeypatch) -> None:
    from ppt_creator_ai.providers import get_provider

    provider = get_provider("ollama_local")
    monkeypatch.setattr(
        provider,
        "list_models",
        lambda: {
            "provider": "ollama_local",
            "service_url": "http://127.0.0.1:11434",
            "model_count": 2,
            "models": [
                {"name": "qwen2.5:7b", "size": 123},
                {"name": "llama3.2:3b", "size": 456},
            ],
        },
    )
    monkeypatch.setattr(
        provider,
        "status_payload",
        lambda: {
            "service_url": "http://127.0.0.1:11434",
            "provider_name": "ollama",
            "provider_source": "app_default",
            "model_name": "qwen2.5:7b",
            "model_source": "auto_discovered",
            "health_status": "ok",
            "model_count": 2,
            "models": [{"name": "qwen2.5:7b"}, {"name": "llama3.2:3b"}],
            "supports_model_listing": True,
        },
    )

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        status, models_payload = _request_json(f"{base_url}/ai/models?provider_name=ollama_local")
        assert status == 200
        assert models_payload["result"]["provider"] == "ollama_local"
        assert [item["name"] for item in models_payload["result"]["models"]] == ["qwen2.5:7b", "llama3.2:3b"]

        status, ai_status_payload = _request_json(f"{base_url}/ai/status?provider_name=ollama_local")
        assert status == 200
        assert ai_status_payload["result"]["selected_provider"] == "ollama_local"
        assert ai_status_payload["result"]["provider_status"]["supports_model_listing"] is True
        assert ai_status_payload["result"]["provider_status"]["model_count"] == 2
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_playground_html_embedded_script_remains_valid(tmp_path: Path) -> None:
    html = _build_playground_html()
    assert "join('\\n')" in html
    assert "split(/\\n+/)" in html
    assert "attachGenerationContextToActionData" in html
    assert "runIterateFlow()" in html
    assert "exportCurrentDeck()" in html
    assert "focusSlide(index" in html
    assert "focusTopRiskSlide()" in html
    assert "syncLiveReviewSummary()" in html
    assert "latestActionResult = null" in html
    assert "optionalUnitFloat(" in html
    assert "Top risk right now:" in html
    assert "document.addEventListener('keydown'" in html
    assert "await runAction('/review', { generationResult: result });" in html
    assert "await runAction('/render', { generationResult: result });" in html
    assert "const merged = attachGenerationContextToActionData(response.data, options.generationResult);" in html

    match = re.search(r"<script>(.*?)</script>", html, re.S)
    assert match is not None

    node_binary = shutil.which("node")
    if node_binary is None:
        return

    script_path = tmp_path / "playground.js"
    script_path.write_text(match.group(1), encoding="utf-8")
    result = subprocess.run(
        [node_binary, "--check", str(script_path)],
        capture_output=True,
        text=True,
        check=False,
    )
    assert result.returncode == 0, result.stderr


def test_api_generate_and_generate_and_render_endpoints(tmp_path: Path) -> None:
    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        briefing_payload = json.loads(Path("examples/briefing_sales.json").read_text(encoding="utf-8"))

        status, generate_payload = _request_json(
            f"{base_url}/generate",
            {"briefing": briefing_payload},
            method="POST",
        )
        assert status == 200
        assert generate_payload["result"]["mode"] == "generate"
        assert generate_payload["result"]["provider"] == "heuristic"
        assert generate_payload["result"]["transport_provider"] == "heuristic"
        assert generate_payload["result"]["slide_count"] >= 4

        status, intent_generate_payload = _request_json(
            f"{base_url}/generate",
            {
                "intent_text": "Quero um deck para o board explicando por que um copiloto de vendas deve ser lançado agora, com três benefícios, métricas, comparação de opções e fechamento forte.",
                "provider_name": "heuristic",
            },
            method="POST",
        )
        assert status == 200
        assert intent_generate_payload["result"]["mode"] == "generate"
        assert intent_generate_payload["result"]["slide_count"] >= 5

        preview_dir = tmp_path / "generated_from_intent_previews"
        output_path_from_intent = tmp_path / "generated_from_intent.pptx"
        status, intent_generate_render_payload = _request_json(
            f"{base_url}/generate-and-render",
            {
                "intent_text": "Monte uma apresentação de product operating review mostrando onde o roadmap está diluído, quais decisões precisam ser tomadas, métricas, riscos, trade-offs e a sequência recomendada para o próximo trimestre.",
                "provider_name": "heuristic",
                "output_path": str(output_path_from_intent),
                "include_review": True,
                "preview_output_dir": str(preview_dir),
                "preview_backend": "synthetic",
            },
            method="POST",
        )
        assert status == 200
        assert intent_generate_render_payload["result"]["generation"]["mode"] == "generate"
        assert intent_generate_render_payload["result"]["render"]["rendered"] is True
        assert intent_generate_render_payload["result"]["render"]["preview_result"]["preview_count"] >= 1
        assert Path(intent_generate_render_payload["result"]["render"]["preview_result"]["thumbnail_sheet"]).exists()

        output_path = tmp_path / "generated_from_api.pptx"
        status, generate_render_payload = _request_json(
            f"{base_url}/generate-and-render",
            {
                "briefing": briefing_payload,
                "output_path": str(output_path),
                "include_review": True,
            },
            method="POST",
        )
        assert status == 200
        assert generate_render_payload["result"]["generation"]["mode"] == "generate"
        assert generate_render_payload["result"]["render"]["rendered"] is True
        assert Path(generate_render_payload["result"]["render"]["output_path"]).exists()
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_generate_intent_with_local_service_defaults_to_ai_first_authoring(monkeypatch) -> None:
    from ppt_creator.api import generate_briefing_payload
    from ppt_creator_ai.providers import get_provider
    from ppt_creator_ai.providers.base import BriefingGenerationResult

    provider = get_provider("local_service")
    captured: dict[str, object] = {}

    def _fake_generate(briefing, *, theme_name=None, feedback_messages=None):
        captured["briefing"] = briefing.model_dump(mode="json")
        return BriefingGenerationResult(
            provider_name="local_service",
            payload={
                "presentation": {"title": briefing.title, "theme": "executive_premium_minimal"},
                "slides": [
                    {"type": "title", "title": briefing.title},
                    {"type": "closing", "title": "Closing", "quote": "Done."},
                ],
            },
            analysis={"provider": "local_service"},
        )

    monkeypatch.setattr(provider, "generate", _fake_generate)

    result = generate_briefing_payload(
        provider_name="local_service",
        intent_text="Quero uma apresentação para entrevista de AI Engineer mostrando minha trajetória, stack e projetos em IA.",
    )

    assert result["provider"] == "local_service"
    assert result["transport_provider"] == "local_service"
    assert result["authoring_mode"] == "ai_first"
    assert captured["briefing"]["outline"] == []
    assert captured["briefing"]["recommendations"] == []
    assert captured["briefing"]["briefing_text"]


def test_api_generate_intent_without_provider_defaults_to_model_backed_ai_first(monkeypatch) -> None:
    from ppt_creator.api import generate_briefing_payload
    from ppt_creator_ai.providers import get_provider
    from ppt_creator_ai.providers.base import BriefingGenerationResult

    provider = get_provider("local_service")
    captured: dict[str, object] = {}

    def _fake_generate(briefing, *, theme_name=None, feedback_messages=None):
        captured["briefing"] = briefing.model_dump(mode="json")
        return BriefingGenerationResult(
            provider_name="local_service",
            payload={
                "presentation": {"title": briefing.title, "theme": "executive_premium_minimal"},
                "slides": [
                    {"type": "title", "title": briefing.title},
                    {"type": "closing", "title": "Closing", "quote": "Done."},
                ],
            },
            analysis={"provider": "local_service"},
        )

    monkeypatch.setattr(provider, "generate", _fake_generate)

    result = generate_briefing_payload(
        intent_text="Quero uma apresentação para entrevista de AI Engineer mostrando minha trajetória, stack e projetos em IA.",
    )

    assert result["provider"] == "local_service"
    assert result["transport_provider"] == "local_service"
    assert result["authoring_mode"] == "ai_first"
    assert captured["briefing"]["outline"] == []
    assert captured["briefing"]["recommendations"] == []
    assert captured["briefing"]["briefing_text"]


def test_api_validate_render_and_template_endpoints(tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")

        status, validate_payload = _request_json(
            f"{base_url}/validate",
            {"spec": spec_payload, "check_assets": True},
            method="POST",
        )
        assert status == 200
        assert validate_payload["result"]["valid"] is True
        assert validate_payload["result"]["slide_count"] == len(spec_payload["slides"])

        output_path = tmp_path / "api_render_output.pptx"
        status, render_payload = _request_json(
            f"{base_url}/render",
            {
                "spec": spec_payload,
                "output_path": str(output_path),
                "dry_run": True,
                "include_review": True,
            },
            method="POST",
        )
        assert status == 200
        assert render_payload["result"]["rendered"] is False
        assert render_payload["result"]["output_path"].endswith("api_render_output.pptx")
        assert render_payload["result"]["quality_review"] is not None
        assert "clipping_risk_count" in render_payload["result"]["quality_review"]

        status, review_payload = _request_json(
            f"{base_url}/review",
            {"spec": spec_payload},
            method="POST",
        )
        assert status == 200
        assert review_payload["result"]["slide_count"] == len(spec_payload["slides"])
        assert review_payload["result"]["status"] in {"ok", "review", "attention"}
        assert "severity_counts" in review_payload["result"]
        assert "overflow_risk_count" in review_payload["result"]
        assert "clipping_risk_count" in review_payload["result"]
        assert "collision_risk_count" in review_payload["result"]
        assert "balance_warning_count" in review_payload["result"]
        assert "top_risk_slides" in review_payload["result"]

        status, template_payload = _request_json(
            f"{base_url}/template",
            {"domain": "sales", "theme_name": "consulting_clean", "audience_profile": "board", "brand_pack": "board_navy"},
            method="POST",
        )
        assert status == 200
        assert template_payload["template"]["presentation"]["theme"] == "consulting_clean"
        assert template_payload["template"]["presentation"]["footer_text"] == "Board brand pack"
        assert template_payload["packet"]["asset_collections"]
        assert template_payload["packet"]["asset_strategy"]["cover_asset_collection"] == "boardroom_backdrops"
        assert template_payload["packet"]["slide_asset_suggestions"]

        status, proposal_template_payload = _request_json(
            f"{base_url}/template",
            {"domain": "proposal", "audience_profile": "proposal"},
            method="POST",
        )
        assert status == 200
        assert proposal_template_payload["template"]["presentation"]["title"] == "Commercial proposal"
        assert proposal_template_payload["template"]["presentation"]["footer_text"] == "Proposal profile"

        status, workflow_template_payload = _request_json(
            f"{base_url}/workflow-template",
            {"workflow_name": "sales_qbr", "brand_pack": "sales_pipeline"},
            method="POST",
        )
        assert status == 200
        assert workflow_template_payload["packet"]["workflow"]["name"] == "sales_qbr"
        assert workflow_template_payload["packet"]["template"]["presentation"]["footer_text"] == "Sales pipeline brand pack"
        assert workflow_template_payload["packet"]["preview_recommendation"]["recommended_source"] == "rendered_pptx"
        assert workflow_template_payload["packet"]["preview_recommendation"]["require_real_previews"] is True
        assert str(workflow_template_payload["packet"]["preview_recommendation"]["baseline_dir"]).endswith("sales_qbr_baseline")
        assert workflow_template_payload["packet"]["asset_strategy"]["placeholder_style"] == "analytical_visual"
        assert workflow_template_payload["packet"]["slide_asset_suggestions"]

        status, proposal_workflow_template_payload = _request_json(
            f"{base_url}/workflow-template",
            {"workflow_name": "commercial_proposal"},
            method="POST",
        )
        assert status == 200
        assert proposal_workflow_template_payload["packet"]["workflow"]["name"] == "commercial_proposal"
        assert proposal_workflow_template_payload["packet"]["workflow"]["domain"] == "proposal"
        assert proposal_workflow_template_payload["packet"]["template"]["presentation"]["title"] == "Commercial proposal"

        preview_module.find_office_runtime = lambda: None
        preview_dir = tmp_path / "api_previews"
        status, preview_payload = _request_json(
            f"{base_url}/preview",
            {
                "spec": spec_payload,
                "output_dir": str(preview_dir),
                "basename": "api-preview",
                "debug_grid": True,
                "debug_safe_areas": True,
            },
            method="POST",
        )
        assert status == 200
        assert preview_payload["result"]["preview_count"] == len(spec_payload["slides"])
        assert preview_payload["result"]["quality_review"]["status"] in {"ok", "review"}
        assert "severity_counts" in preview_payload["result"]["quality_review"]
        assert "clipping_risk_count" in preview_payload["result"]["quality_review"]
        assert preview_payload["result"]["preview_artifact_review"]["status"] in {"ok", "review"}
        assert preview_payload["result"]["backend_requested"] == "auto"
        assert preview_payload["result"]["backend_used"] in {"synthetic", "office"}
        assert Path(preview_payload["result"]["thumbnail_sheet"]).exists()

        artifact_request = request.Request(
            f"{base_url}/artifact?path={request.pathname2url(str(Path(preview_payload['result']['thumbnail_sheet'])))}",
            method="GET",
        )
        with request.urlopen(artifact_request, timeout=5) as response:
            artifact_bytes = response.read()
        assert artifact_bytes

        baseline_dir = tmp_path / "api_baseline_previews"
        status, baseline_preview_payload = _request_json(
            f"{base_url}/preview",
            {
                "spec": spec_payload,
                "output_dir": str(baseline_dir),
                "basename": "api-baseline",
            },
            method="POST",
        )
        assert status == 200
        assert baseline_preview_payload["result"]["preview_count"] == len(spec_payload["slides"])

        regression_dir = tmp_path / "api_regression_previews"
        status, regression_preview_payload = _request_json(
            f"{base_url}/preview",
            {
                "spec": spec_payload,
                "output_dir": str(regression_dir),
                "basename": "api-regression",
                "baseline_dir": str(baseline_dir),
                "write_diff_images": True,
            },
            method="POST",
        )
        assert status == 200
        assert regression_preview_payload["result"]["visual_regression"] is not None
        assert regression_preview_payload["result"]["visual_regression"]["diff_count"] == 0
        assert regression_preview_payload["result"]["visual_regression"]["current_manifest"]
        assert regression_preview_payload["result"]["visual_regression"]["baseline_manifest"]
        assert regression_preview_payload["result"]["visual_regression"]["source_mismatch"] is False
        assert regression_preview_payload["result"]["visual_regression"]["guidance"] == []

        pptx_path = tmp_path / "api_preview_source.pptx"
        PresentationRenderer(asset_root="examples").render(PresentationInput.model_validate(spec_payload), pptx_path)

        def _fake_run(command, capture_output, text, check):
            outdir = Path(command[command.index("--outdir") + 1])
            source_pptx = Path(command[-1])
            slide_count = len(PptxPresentation(str(source_pptx)).slides)
            outdir.mkdir(parents=True, exist_ok=True)
            for index in range(1, slide_count + 1):
                Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"api-mock-{index:02d}.png")

            class _Completed:
                returncode = 0
                stdout = ""
                stderr = ""

            return _Completed()

        preview_module.find_office_runtime = lambda: "/usr/bin/soffice"
        preview_module.subprocess.run = _fake_run

        pptx_preview_dir = tmp_path / "api_preview_pptx_output"
        status, pptx_preview_payload = _request_json(
            f"{base_url}/preview-pptx",
            {
                "input_pptx": str(pptx_path),
                "output_dir": str(pptx_preview_dir),
            },
            method="POST",
        )
        assert status == 200
        assert pptx_preview_payload["result"]["mode"] == "preview-pptx"
        assert pptx_preview_payload["result"]["preview_count"] == len(spec_payload["slides"])
        assert pptx_preview_payload["result"]["backend_used"] == "office"
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_render_can_generate_previews_from_rendered_pptx(tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")

        def _fake_run(command, capture_output, text, check):
            outdir = Path(command[command.index("--outdir") + 1])
            source_pptx = Path(command[-1])
            slide_count = len(PptxPresentation(str(source_pptx)).slides)
            outdir.mkdir(parents=True, exist_ok=True)
            for index in range(1, slide_count + 1):
                Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"api-render-preview-{index:02d}.png")

            class _Completed:
                returncode = 0
                stdout = ""
                stderr = ""

            return _Completed()

        preview_module.find_office_runtime = lambda: "/usr/bin/soffice"
        preview_module.subprocess.run = _fake_run

        output_path = tmp_path / "api_render_with_preview_output.pptx"
        preview_dir = tmp_path / "api_render_previews"
        status, render_payload = _request_json(
            f"{base_url}/render",
            {
                "spec": spec_payload,
                "output_path": str(output_path),
                "preview_output_dir": str(preview_dir),
            },
            method="POST",
        )

        assert status == 200
        assert render_payload["result"]["rendered"] is True
        assert render_payload["result"]["preview_output_dir"] == str(preview_dir)
        assert render_payload["result"]["preview_source"] == "rendered_pptx"
        assert render_payload["result"]["preview_result"]["mode"] == "preview-pptx"
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_review_can_attach_preview_result_with_real_artifact_preference(tmp_path: Path) -> None:
    from ppt_creator import api as api_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")
        preview_dir = tmp_path / "api_review_previews"

        api_module.render_previews_for_rendered_artifact = lambda *args, **kwargs: (
            {
                "mode": "preview-pptx",
                "preview_count": 10,
                "previews": [],
                "thumbnail_sheet": str(preview_dir / "thumbs.png"),
                "preview_artifact_review": {"status": "ok"},
                "visual_regression": None,
                "backend_used": "office",
            },
            "rendered_pptx",
        )

        status, review_payload = _request_json(
            f"{base_url}/review",
            {
                "spec": spec_payload,
                "preview_output_dir": str(preview_dir),
            },
            method="POST",
        )
        assert status == 200
        assert review_payload["result"]["preview_source"] == "rendered_pptx"
        assert review_payload["result"]["preview_result"]["mode"] == "preview-pptx"
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_preview_pptx_falls_back_to_pdf_rasterization_when_needed(tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")
        pptx_path = tmp_path / "api_preview_pdf_fallback_source.pptx"
        PresentationRenderer(asset_root="examples").render(PresentationInput.model_validate(spec_payload), pptx_path)

        def _fake_run(command, capture_output, text, check):
            outdir = Path(command[command.index("--outdir") + 1]) if "--outdir" in command else tmp_path
            outdir.mkdir(parents=True, exist_ok=True)
            if command[0] == "/usr/bin/soffice" and "png" in command:
                Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / "single.png")
            elif command[0] == "/usr/bin/soffice" and "pdf" in command:
                (outdir / "api_preview_pdf_fallback_source.pdf").write_bytes(b"%PDF-1.4 mock")
            elif command[0] == "/usr/bin/gs":
                slide_count = len(PptxPresentation(str(pptx_path)).slides)
                pattern = next(argument.split("=", 1)[1] for argument in command if argument.startswith("-sOutputFile="))
                for index in range(1, slide_count + 1):
                    target = Path(pattern.replace("%02d", f"{index:02d}"))
                    Image.new("RGB", (1280, 720), (245, 245, 245)).save(target)

            class _Completed:
                returncode = 0
                stdout = ""
                stderr = ""

            return _Completed()

        preview_module.find_office_runtime = lambda: "/usr/bin/soffice"
        preview_module.find_ghostscript_runtime = lambda: "/usr/bin/gs"
        preview_module.subprocess.run = _fake_run

        pptx_preview_dir = tmp_path / "api_preview_pdf_fallback_output"
        status, pptx_preview_payload = _request_json(
            f"{base_url}/preview-pptx",
            {
                "input_pptx": str(pptx_path),
                "output_dir": str(pptx_preview_dir),
            },
            method="POST",
        )
        assert status == 200
        assert pptx_preview_payload["result"]["preview_count"] == len(spec_payload["slides"])
        assert pptx_preview_payload["result"]["office_conversion_strategy"] == "pdf_via_ghostscript"
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_compare_pptx_endpoint_generates_visual_comparison(tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec = PresentationInput.from_path("examples/ai_sales.json")
        before_pptx = tmp_path / "api_compare_before.pptx"
        after_pptx = tmp_path / "api_compare_after.pptx"
        PresentationRenderer(asset_root="examples").render(spec, before_pptx)
        PresentationRenderer(asset_root="examples").render(spec, after_pptx)

        def _fake_run(command, capture_output, text, check):
            outdir = Path(command[command.index("--outdir") + 1]) if "--outdir" in command else tmp_path
            source_pptx = Path(command[-1])
            slide_count = len(PptxPresentation(str(source_pptx)).slides)
            outdir.mkdir(parents=True, exist_ok=True)
            for index in range(1, slide_count + 1):
                Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"api-compare-{index:02d}.png")

            class _Completed:
                returncode = 0
                stdout = ""
                stderr = ""

            return _Completed()

        preview_module.find_office_runtime = lambda: "/usr/bin/soffice"
        preview_module.subprocess.run = _fake_run

        comparison_dir = tmp_path / "api_compare_output"
        status, comparison_payload = _request_json(
            f"{base_url}/compare-pptx",
            {
                "before_pptx": str(before_pptx),
                "after_pptx": str(after_pptx),
                "output_dir": str(comparison_dir),
                "write_diff_images": True,
            },
            method="POST",
        )
        assert status == 200
        assert comparison_payload["result"]["mode"] == "compare-pptx"
        assert comparison_payload["result"]["comparison"]["status"] == "ok"
        assert comparison_payload["result"]["comparison"]["diff_count"] == 0
        assert comparison_payload["result"]["before_preview_manifest"]
        assert comparison_payload["result"]["after_preview_manifest"]
        assert comparison_payload["result"]["comparison"]["current_real_preview"] is True
        assert comparison_payload["result"]["comparison"]["baseline_real_preview"] is True
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_preview_can_require_real_previews(tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")
        preview_module.find_office_runtime = lambda: None

        try:
            _request_json(
                f"{base_url}/preview",
                {
                    "spec": spec_payload,
                    "output_dir": str(tmp_path / "api_require_real_preview"),
                    "require_real_previews": True,
                },
                method="POST",
            )
            assert False, "request should fail when real previews are required without office runtime"
        except Exception as exc:  # noqa: BLE001
            assert "HTTP Error 500" in str(exc)
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_preview_fail_on_regression_returns_conflict(tmp_path: Path) -> None:
    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")
        baseline_dir = tmp_path / "api_fail_baseline"

        status, _ = _request_json(
            f"{base_url}/preview",
            {
                "spec": spec_payload,
                "output_dir": str(baseline_dir),
                "basename": "api-fail-baseline",
            },
            method="POST",
        )
        assert status == 200

        try:
            _request_json(
                f"{base_url}/preview",
                {
                    "spec": spec_payload,
                    "output_dir": str(tmp_path / "api_fail_current"),
                    "basename": "api-fail-current",
                    "theme_name": "dark_boardroom",
                    "baseline_dir": str(baseline_dir),
                    "fail_on_regression": True,
                },
                method="POST",
            )
            assert False, "preview should fail when fail_on_regression is enabled and diffs exist"
        except Exception as exc:  # noqa: BLE001
            assert "HTTP Error 409" in str(exc)
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_promote_baseline_endpoint_copies_preview_set(tmp_path: Path) -> None:
    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec_payload = PresentationInput.from_path("examples/ai_sales.json").model_dump(mode="json")
        source_dir = tmp_path / "api_promote_source"
        baseline_dir = tmp_path / "api_promote_target"

        status, preview_payload = _request_json(
            f"{base_url}/preview",
            {
                "spec": spec_payload,
                "output_dir": str(source_dir),
                "basename": "api-promote-source",
            },
            method="POST",
        )
        assert status == 200

        status, promote_payload = _request_json(
            f"{base_url}/promote-baseline",
            {
                "source_dir": str(source_dir),
                "baseline_dir": str(baseline_dir),
            },
            method="POST",
        )
        assert status == 200
        assert promote_payload["result"]["mode"] == "promote-baseline"
        assert promote_payload["result"]["preview_count"] == preview_payload["result"]["preview_count"]
        assert Path(promote_payload["result"]["preview_manifest"]).exists()
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


def test_api_review_pptx_endpoint_generates_real_artifact_review(tmp_path: Path) -> None:
    from ppt_creator import preview as preview_module

    server = build_api_server("127.0.0.1", 0, asset_root="examples")
    thread = Thread(target=server.serve_forever, daemon=True)
    thread.start()

    try:
        base_url = f"http://127.0.0.1:{server.server_address[1]}"
        spec = PresentationInput.from_path("examples/ai_sales.json")
        pptx_path = tmp_path / "api_review_source.pptx"
        PresentationRenderer(asset_root="examples").render(spec, pptx_path)

        def _fake_run(command, capture_output, text, check):
            outdir = Path(command[command.index("--outdir") + 1])
            source_pptx = Path(command[-1])
            slide_count = len(PptxPresentation(str(source_pptx)).slides)
            outdir.mkdir(parents=True, exist_ok=True)
            for index in range(1, slide_count + 1):
                Image.new("RGB", (1280, 720), (245, 245, 245)).save(outdir / f"api-review-{index:02d}.png")

            class _Completed:
                returncode = 0
                stdout = ""
                stderr = ""

            return _Completed()

        preview_module.find_office_runtime = lambda: "/usr/bin/soffice"
        preview_module.subprocess.run = _fake_run

        review_dir = tmp_path / "api_review_pptx_output"
        status, review_payload = _request_json(
            f"{base_url}/review-pptx",
            {
                "input_pptx": str(pptx_path),
                "output_dir": str(review_dir),
            },
            method="POST",
        )
        assert status == 200
        assert review_payload["result"]["mode"] == "review-pptx"
        assert review_payload["result"]["preview_result"]["mode"] == "preview-pptx"
        assert "preview_artifact_review" in review_payload["result"]
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)

from __future__ import annotations

import json
import re
import shutil
import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory

from PIL import Image, ImageChops, ImageDraw, ImageFont, ImageOps, ImageStat
from pptx import Presentation

from ppt_creator.qa import review_presentation, review_preview_artifacts
from ppt_creator.renderer import PresentationRenderer, infer_contextual_image_focal_point
from ppt_creator.schema import PresentationInput, PresentationMeta, Slide, SlideType
from ppt_creator.theme import get_theme

PREVIEW_WIDTH = 1280
PREVIEW_HEIGHT = 720
THUMBNAIL_WIDTH = 320
THUMBNAIL_HEIGHT = 180
CONTACT_SHEET_COLUMNS = 3
PREVIEW_MANIFEST_FILENAME = "preview-manifest.json"


def _risk_badge_fill(*, status: str | None, risk_level: str | None) -> tuple[int, int, int]:
    if risk_level == "high" or status == "attention":
        return (184, 103, 87)
    if risk_level == "medium" or status == "review":
        return (176, 139, 91)
    return (97, 114, 135)


def find_office_runtime() -> str | None:
    for candidate in ("soffice", "libreoffice"):
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    return None


def find_ghostscript_runtime() -> str | None:
    return shutil.which("gs")


def _rgb_tuple(hex_value: str) -> tuple[int, int, int]:
    normalized = hex_value.replace("#", "")
    return tuple(int(normalized[i : i + 2], 16) for i in (0, 2, 4))


def _safe_basename(value: str) -> str:
    normalized = re.sub(r"[^a-z0-9]+", "_", value.lower()).strip("_")
    return normalized or "deck_preview"


def _truncate_text(value: str, *, max_chars: int) -> str:
    if len(value) <= max_chars:
        return value
    return value[: max_chars - 3].rstrip() + "..."


def _load_font(size: int, *, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def _list_preview_images(path: str | Path) -> list[Path]:
    root = Path(path)
    return sorted(
        file_path
        for file_path in root.glob("*.png")
        if not file_path.name.endswith("-thumbnails.png") and "-diff-" not in file_path.stem
    )


def _preview_manifest_path(path: str | Path) -> Path:
    return Path(path) / PREVIEW_MANIFEST_FILENAME


def _load_preview_manifest(path: str | Path) -> dict[str, object] | None:
    manifest_path = _preview_manifest_path(path)
    if not manifest_path.exists() or not manifest_path.is_file():
        return None
    try:
        payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    return payload if isinstance(payload, dict) else None


def _resolve_preview_sequence(path: str | Path) -> tuple[list[Path], dict[str, object] | None]:
    root = Path(path)
    manifest = _load_preview_manifest(root)
    if manifest is not None:
        slides = manifest.get("slides")
        if isinstance(slides, list):
            resolved: list[Path] = []
            for item in slides:
                if not isinstance(item, dict):
                    continue
                filename = item.get("filename")
                if not isinstance(filename, str) or not filename.strip():
                    continue
                candidate = root / filename
                if candidate.exists() and candidate.is_file():
                    resolved.append(candidate)
            if resolved:
                return resolved, manifest
    return _list_preview_images(root), manifest


def _is_real_preview_source(value: object) -> bool:
    if not isinstance(value, str):
        return False
    return value in {"rendered_pptx", "input_pptx"}


def _write_preview_manifest(
    output_dir: str | Path,
    *,
    mode: str,
    preview_paths: list[str],
    basename: str,
    preview_source: str,
    presentation_title: str | None,
    theme_name: str | None,
    backend_requested: str,
    backend_used: str,
    office_runtime: str | None,
    office_conversion_strategy: str | None,
    input_pptx: str | None = None,
) -> Path:
    destination = Path(output_dir)
    destination.mkdir(parents=True, exist_ok=True)
    manifest_path = _preview_manifest_path(destination)
    payload = {
        "manifest_version": 1,
        "mode": mode,
        "basename": basename,
        "output_dir": str(destination.resolve()),
        "presentation_title": presentation_title,
        "theme": theme_name,
        "preview_count": len(preview_paths),
        "preview_source": preview_source,
        "real_preview": _is_real_preview_source(preview_source),
        "backend_requested": backend_requested,
        "backend_used": backend_used,
        "office_runtime": office_runtime,
        "office_conversion_strategy": office_conversion_strategy,
        "input_pptx": input_pptx,
        "slides": [
            {
                "slide_number": index,
                "filename": Path(preview_path).name,
                "path": str(Path(preview_path).resolve()),
            }
            for index, preview_path in enumerate(preview_paths, start=1)
        ],
    }
    manifest_path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    return manifest_path


def _build_visual_regression_guidance(
    *,
    diff_count: int,
    missing_baseline_count: int,
    extra_baseline_count: int,
    source_mismatch: bool,
    current_source: str | None,
    baseline_source: str | None,
) -> list[str]:
    guidance: list[str] = []
    if diff_count:
        guidance.append(
            f"{diff_count} slide(s) differ from the baseline; inspect top_regressions and any generated diff images."
        )
    if missing_baseline_count:
        guidance.append(
            f"{missing_baseline_count} current slide(s) have no matching baseline slide; promote a refreshed baseline if this change is intentional."
        )
    if extra_baseline_count:
        guidance.append(
            f"The baseline contains {extra_baseline_count} extra slide(s); refresh the golden set so the comparison uses the same slide count."
        )
    if source_mismatch:
        guidance.append(
            "Current and baseline previews were generated from different provenance "
            f"({current_source or 'unknown'} vs {baseline_source or 'unknown'}); regenerate both from the same source before trusting the diff."
        )
    return guidance


def visual_regression_has_failures(visual_regression: dict[str, object] | None) -> bool:
    if not isinstance(visual_regression, dict):
        return False
    return str(visual_regression.get("status") or "ok") != "ok"


def format_visual_regression_failure(
    visual_regression: dict[str, object] | None,
    *,
    context: str,
) -> str:
    if not isinstance(visual_regression, dict):
        return f"{context} did not produce a visual regression report. Provide a baseline first."

    parts: list[str] = [f"{context} failed visual regression"]
    diff_count = int(visual_regression.get("diff_count") or 0)
    missing_baseline_count = int(visual_regression.get("missing_baseline_count") or 0)
    extra_baseline_count = int(visual_regression.get("extra_baseline_count") or 0)
    if diff_count:
        parts.append(f"{diff_count} diff(s)")
    if missing_baseline_count:
        parts.append(f"{missing_baseline_count} missing baseline slide(s)")
    if extra_baseline_count:
        parts.append(f"{extra_baseline_count} extra baseline slide(s)")
    if visual_regression.get("source_mismatch"):
        parts.append("preview provenance mismatch")

    message = "; ".join(parts)
    guidance = visual_regression.get("guidance") or []
    if isinstance(guidance, list) and guidance:
        first_guidance = next((item for item in guidance if isinstance(item, str) and item.strip()), None)
        if first_guidance:
            message = f"{message}. {first_guidance}"
    return message


def _normalized_diff_score(generated: Image.Image, baseline: Image.Image) -> tuple[float, float, Image.Image]:
    left = generated.convert("RGB")
    right = baseline.convert("RGB")
    size_mismatch = left.size != right.size
    if size_mismatch:
        right = right.resize(left.size)

    diff = ImageChops.difference(left, right)
    grayscale = diff.convert("L")
    mean_diff = float(ImageStat.Stat(grayscale).mean[0]) / 255.0
    histogram = grayscale.histogram()
    changed_pixels = sum(count for index, count in enumerate(histogram) if index > 0)
    total_pixels = max(1, grayscale.size[0] * grayscale.size[1])
    changed_ratio = changed_pixels / total_pixels
    if size_mismatch:
        mean_diff = max(mean_diff, 1.0)
        changed_ratio = max(changed_ratio, 1.0)
    return mean_diff, changed_ratio, diff


def _changed_ratio_against_background(
    image: Image.Image,
    background_rgb: tuple[int, int, int],
    *,
    box: tuple[int, int, int, int] | None = None,
) -> float:
    region = image.crop(box) if box else image
    background = Image.new("RGB", region.size, background_rgb)
    diff = ImageChops.difference(region.convert("RGB"), background)
    histogram = diff.convert("L").histogram()
    changed_pixels = sum(count for index, count in enumerate(histogram) if index > 0)
    total_pixels = max(1, region.size[0] * region.size[1])
    return changed_pixels / total_pixels


def _foreground_bbox_against_background(
    image: Image.Image,
    background_rgb: tuple[int, int, int],
    *,
    box: tuple[int, int, int, int] | None = None,
) -> tuple[int, int, int, int] | None:
    region = image.crop(box) if box else image
    background = Image.new("RGB", region.size, background_rgb)
    return ImageChops.difference(region.convert("RGB"), background).getbbox()


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for paragraph in text.splitlines() or [text]:
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            bbox = draw.textbbox((0, 0), candidate, font=font)
            if bbox[2] - bbox[0] <= max_width:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


class PreviewRenderer:
    def __init__(
        self,
        theme_name: str | None = None,
        asset_root: str | Path | None = None,
        primary_color: str | None = None,
        secondary_color: str | None = None,
        debug_grid: bool = False,
        debug_safe_areas: bool = False,
        backend: str = "auto",
        baseline_dir: str | Path | None = None,
        diff_threshold: float = 0.01,
        write_diff_images: bool = False,
        require_real_previews: bool = False,
    ):
        self.requested_theme_name = theme_name
        self.requested_primary_color = primary_color
        self.requested_secondary_color = secondary_color
        self.asset_root = Path(asset_root or ".").resolve()
        self.debug_grid = debug_grid
        self.debug_safe_areas = debug_safe_areas
        self.backend = backend
        self.baseline_dir = Path(baseline_dir).resolve() if baseline_dir else None
        self.diff_threshold = diff_threshold
        self.write_diff_images = write_diff_images
        self.require_real_previews = require_real_previews
        self.theme = get_theme(theme_name, primary_color=primary_color, secondary_color=secondary_color)

    def _run_office_convert(
        self,
        runtime: str,
        *,
        source_path: Path,
        outdir: Path,
        target_format: str,
    ) -> subprocess.CompletedProcess[str]:
        command = [
            runtime,
            "--headless",
            "--convert-to",
            target_format,
            "--outdir",
            str(outdir),
            str(source_path),
        ]
        completed = subprocess.run(command, capture_output=True, text=True, check=False)
        if completed.returncode != 0:
            stderr = (completed.stderr or completed.stdout or "").strip()
            raise RuntimeError(stderr or f"office conversion failed with exit code {completed.returncode}")
        return completed

    def _rasterize_pdf_with_ghostscript(
        self,
        pdf_path: Path,
        outdir: Path,
        *,
        basename: str,
        expected_count: int,
    ) -> list[Path]:
        runtime = find_ghostscript_runtime()
        if not runtime:
            raise RuntimeError(
                "office conversion produced a single export and Ghostscript ('gs') is not available for PDF-to-PNG page rasterization"
            )

        output_pattern = outdir / f"{basename}-%02d.png"
        completed = subprocess.run(
            [
                runtime,
                "-dSAFER",
                "-dBATCH",
                "-dNOPAUSE",
                "-sDEVICE=png16m",
                "-r160",
                f"-sOutputFile={output_pattern}",
                str(pdf_path),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        if completed.returncode != 0:
            stderr = (completed.stderr or completed.stdout or "").strip()
            raise RuntimeError(stderr or f"ghostscript rasterization failed with exit code {completed.returncode}")

        rasterized = sorted(path for path in outdir.glob(f"{basename}-*.png") if path.is_file())
        if len(rasterized) != expected_count:
            raise RuntimeError(
                f"ghostscript rasterization produced {len(rasterized)} PNG(s) for {expected_count} slide(s)"
            )
        return rasterized

    def _convert_pptx_to_png_paths(
        self,
        *,
        runtime: str,
        pptx_path: Path,
        outdir: Path,
        basename: str,
        expected_count: int,
    ) -> tuple[list[Path], str]:
        self._run_office_convert(runtime, source_path=pptx_path, outdir=outdir, target_format="png")
        direct_pngs = sorted(path for path in outdir.glob("*.png") if path.is_file())
        if len(direct_pngs) == expected_count:
            return direct_pngs, "direct_png"

        for path in direct_pngs:
            path.unlink(missing_ok=True)

        self._run_office_convert(runtime, source_path=pptx_path, outdir=outdir, target_format="pdf")
        pdf_candidates = sorted(path for path in outdir.glob("*.pdf") if path.is_file())
        if not pdf_candidates:
            raise RuntimeError(
                f"office conversion produced {len(direct_pngs)} PNG(s) for {expected_count} slide(s) and did not produce an intermediate PDF fallback"
            )

        rasterized = self._rasterize_pdf_with_ghostscript(
            pdf_candidates[0],
            outdir,
            basename=basename,
            expected_count=expected_count,
        )
        return rasterized, "pdf_via_ghostscript"

    def render(self, spec: PresentationInput, output_dir: str | Path, *, basename: str | None = None) -> dict[str, object]:
        self.theme = get_theme(
            self.requested_theme_name or spec.presentation.theme,
            primary_color=self.requested_primary_color or spec.presentation.primary_color,
            secondary_color=self.requested_secondary_color or spec.presentation.secondary_color,
        )
        destination = Path(output_dir)
        destination.mkdir(parents=True, exist_ok=True)
        base = basename or _safe_basename(spec.presentation.title)
        runtime = find_office_runtime() if self.backend in {"auto", "office"} else None
        backend_used = "synthetic"
        fallback_reason: str | None = None
        office_conversion_strategy: str | None = None
        preview_source = "spec"

        if self.backend == "office":
            if not runtime:
                raise RuntimeError("Office preview backend requested but no 'soffice' or 'libreoffice' executable was found")
            preview_paths, office_conversion_strategy = self.render_office_previews(spec, destination, base, runtime=runtime)
            backend_used = "office"
            preview_source = "rendered_pptx"
        elif self.backend == "auto" and runtime:
            try:
                preview_paths, office_conversion_strategy = self.render_office_previews(spec, destination, base, runtime=runtime)
                backend_used = "office"
                preview_source = "rendered_pptx"
            except Exception as exc:  # noqa: BLE001
                if self.require_real_previews:
                    raise RuntimeError(
                        f"real preview was required but Office-backed preview failed: {exc}"
                    ) from exc
                fallback_reason = f"Office preview unavailable, falling back to synthetic preview: {exc}"
                preview_paths = self.render_synthetic_previews(spec, destination, base)
                backend_used = "synthetic"
        else:
            if self.backend == "auto" and not runtime:
                if self.require_real_previews:
                    raise RuntimeError(
                        "real preview was required but no 'soffice' or 'libreoffice' runtime was found"
                    )
                fallback_reason = "Office runtime not found; using synthetic preview backend"
            preview_paths = self.render_synthetic_previews(spec, destination, base)

        quality_review = self.build_preview_quality_review(spec)
        contact_sheet_path = destination / f"{base}-thumbnails.png"
        self.render_contact_sheet(
            spec.presentation,
            spec.slides,
            preview_paths,
            contact_sheet_path,
            quality_review=quality_review,
        )

        manifest_path = _write_preview_manifest(
            destination,
            mode="preview",
            preview_paths=preview_paths,
            basename=base,
            preview_source=preview_source,
            presentation_title=spec.presentation.title,
            theme_name=self.theme.name,
            backend_requested=self.backend,
            backend_used=backend_used,
            office_runtime=runtime,
            office_conversion_strategy=office_conversion_strategy,
        )
        manifest = _load_preview_manifest(destination)

        visual_regression = self.build_visual_regression_report(
            preview_paths,
            destination,
            base,
            current_preview_source=preview_source,
            current_manifest=manifest,
        )
        preview_artifact_review = self.build_preview_artifact_review(preview_paths)

        return {
            "mode": "preview",
            "output_dir": str(destination),
            "preview_count": len(preview_paths),
            "previews": preview_paths,
            "thumbnail_sheet": str(contact_sheet_path),
            "quality_review": quality_review,
            "preview_artifact_review": preview_artifact_review,
            "visual_regression": visual_regression,
            "debug_grid": self.debug_grid,
            "debug_safe_areas": self.debug_safe_areas,
            "preview_source": preview_source,
            "real_preview": _is_real_preview_source(preview_source),
            "preview_manifest": str(manifest_path),
            "backend_requested": self.backend,
            "backend_used": backend_used,
            "backend_fallback_reason": fallback_reason,
            "office_runtime": runtime,
            "office_conversion_strategy": office_conversion_strategy,
        }

    def render_synthetic_previews(
        self,
        spec: PresentationInput,
        destination: Path,
        base: str,
    ) -> list[str]:
        preview_paths: list[str] = []
        total_slides = len(spec.slides)
        for index, slide_spec in enumerate(spec.slides, start=1):
            image = self.render_slide(spec.presentation, slide_spec, index, total_slides)
            preview_path = destination / f"{base}-{index:02d}.png"
            image.save(preview_path)
            preview_paths.append(str(preview_path))
        return preview_paths

    def render_office_previews(
        self,
        spec: PresentationInput,
        destination: Path,
        base: str,
        *,
        runtime: str,
    ) -> tuple[list[str], str]:
        with TemporaryDirectory(prefix="ppt_creator_preview_") as tmpdir:
            temp_root = Path(tmpdir)
            pptx_path = temp_root / f"{base}.pptx"
            renderer = PresentationRenderer(
                theme_name=self.requested_theme_name or spec.presentation.theme,
                asset_root=self.asset_root,
                primary_color=self.requested_primary_color or spec.presentation.primary_color,
                secondary_color=self.requested_secondary_color or spec.presentation.secondary_color,
            )
            renderer.render(spec, pptx_path)

            converted_pngs, strategy = self._convert_pptx_to_png_paths(
                runtime=runtime,
                pptx_path=pptx_path,
                outdir=temp_root,
                basename=base,
                expected_count=len(spec.slides),
            )

            preview_paths: list[str] = []
            for index, source_path in enumerate(converted_pngs, start=1):
                target_path = destination / f"{base}-{index:02d}.png"
                shutil.copy2(source_path, target_path)
                preview_paths.append(str(target_path))
            return preview_paths, strategy

    def render_pptx_previews(
        self,
        input_pptx: str | Path,
        output_dir: str | Path,
        *,
        basename: str | None = None,
    ) -> dict[str, object]:
        runtime = find_office_runtime()
        if not runtime:
            raise RuntimeError("PPTX preview requires a local 'soffice' or 'libreoffice' runtime")

        input_path = Path(input_pptx).resolve()
        if not input_path.exists():
            raise FileNotFoundError(f"Input PPTX not found: {input_path}")
        if input_path.suffix.lower() != ".pptx":
            raise ValueError(f"Input PPTX path must end with .pptx: {input_path}")

        destination = Path(output_dir)
        destination.mkdir(parents=True, exist_ok=True)
        base = basename or _safe_basename(input_path.stem)
        expected_slide_count = len(Presentation(str(input_path)).slides)
        office_conversion_strategy: str | None = None

        with TemporaryDirectory(prefix="ppt_creator_preview_pptx_") as tmpdir:
            temp_root = Path(tmpdir)
            temp_pptx = temp_root / f"{base}.pptx"
            shutil.copy2(input_path, temp_pptx)

            converted_pngs, office_conversion_strategy = self._convert_pptx_to_png_paths(
                runtime=runtime,
                pptx_path=temp_pptx,
                outdir=temp_root,
                basename=base,
                expected_count=expected_slide_count,
            )

            preview_paths: list[str] = []
            for index, source_path in enumerate(converted_pngs, start=1):
                target_path = destination / f"{base}-{index:02d}.png"
                shutil.copy2(source_path, target_path)
                preview_paths.append(str(target_path))

        contact_sheet_path = destination / f"{base}-thumbnails.png"
        meta = PresentationMeta.model_construct(title=input_path.stem, subtitle="PPTX preview")
        slides = [
            Slide.model_construct(type=SlideType.TITLE, title=f"Slide {index:02d}")
            for index in range(1, len(preview_paths) + 1)
        ]
        self.render_contact_sheet(meta, slides, preview_paths, contact_sheet_path)

        manifest_path = _write_preview_manifest(
            destination,
            mode="preview-pptx",
            preview_paths=preview_paths,
            basename=base,
            preview_source="input_pptx",
            presentation_title=input_path.stem,
            theme_name=self.theme.name,
            backend_requested="office",
            backend_used="office",
            office_runtime=runtime,
            office_conversion_strategy=office_conversion_strategy,
            input_pptx=str(input_path),
        )
        manifest = _load_preview_manifest(destination)

        visual_regression = self.build_visual_regression_report(
            preview_paths,
            destination,
            base,
            current_preview_source="input_pptx",
            current_manifest=manifest,
        )
        preview_artifact_review = self.build_preview_artifact_review(preview_paths)

        return {
            "mode": "preview-pptx",
            "input_pptx": str(input_path),
            "output_dir": str(destination),
            "preview_count": len(preview_paths),
            "previews": preview_paths,
            "thumbnail_sheet": str(contact_sheet_path),
            "quality_review": None,
            "preview_artifact_review": preview_artifact_review,
            "visual_regression": visual_regression,
            "preview_source": "input_pptx",
            "real_preview": True,
            "preview_manifest": str(manifest_path),
            "backend_requested": "office",
            "backend_used": "office",
            "backend_fallback_reason": None,
            "office_runtime": runtime,
            "office_conversion_strategy": office_conversion_strategy,
        }

    def render_slide(self, meta: PresentationMeta, slide_spec: Slide, index: int, total_slides: int) -> Image.Image:
        colors = self.theme.colors
        image = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), _rgb_tuple(colors.background))
        draw = ImageDraw.Draw(image)

        renderers = {
            SlideType.TITLE: self._render_title,
            SlideType.SECTION: self._render_section,
            SlideType.AGENDA: self._render_agenda,
            SlideType.BULLETS: self._render_bullets,
            SlideType.CARDS: self._render_cards,
            SlideType.METRICS: self._render_metrics,
            SlideType.CHART: self._render_chart,
            SlideType.IMAGE_TEXT: self._render_image_text,
            SlideType.TIMELINE: self._render_timeline,
            SlideType.COMPARISON: self._render_comparison,
            SlideType.TWO_COLUMN: self._render_two_column,
            SlideType.TABLE: self._render_table,
            SlideType.FAQ: self._render_faq,
            SlideType.SUMMARY: self._render_summary,
            SlideType.CLOSING: self._render_closing,
        }
        renderers[slide_spec.type](draw, image, slide_spec, meta)
        self._render_footer(draw, meta, index, total_slides)
        if self.debug_safe_areas or self.debug_grid:
            self._render_debug_overlay(draw)
        return image

    def render_contact_sheet(
        self,
        meta: PresentationMeta,
        slides: list[Slide],
        preview_paths: list[str],
        output_path: str | Path,
        *,
        quality_review: dict[str, object] | None = None,
    ) -> Path:
        columns = min(CONTACT_SHEET_COLUMNS, max(1, len(preview_paths)))
        rows = (len(preview_paths) + columns - 1) // columns
        card_width = THUMBNAIL_WIDTH + 40
        card_height = THUMBNAIL_HEIGHT + 68
        header_height = 108
        margin = 28
        gutter = 24

        sheet = Image.new(
            "RGB",
            (
                margin * 2 + columns * card_width + (columns - 1) * gutter,
                header_height + margin + rows * card_height + max(rows - 1, 0) * gutter + margin,
            ),
            _rgb_tuple(self.theme.colors.background),
        )
        draw = ImageDraw.Draw(sheet)
        title_font = _load_font(26, bold=True)
        subtitle_font = _load_font(16)
        card_title_font = _load_font(16, bold=True)
        card_meta_font = _load_font(13)

        draw.text((margin, 28), meta.title, fill=_rgb_tuple(self.theme.colors.navy), font=title_font)
        subtitle = meta.subtitle or meta.footer_text or meta.client_name or "Preview contact sheet"
        draw.text(
            (margin, 64),
            f"{subtitle} • {len(preview_paths)} slide(s)",
            fill=_rgb_tuple(self.theme.colors.muted),
            font=subtitle_font,
        )
        draw.line(
            (margin, header_height - 10, sheet.width - margin, header_height - 10),
            fill=_rgb_tuple(self.theme.colors.line),
            width=2,
        )

        slide_review_lookup = {
            int(review["slide_number"]): review
            for review in (quality_review or {}).get("slides", [])
        }

        for idx, preview_path in enumerate(preview_paths):
            image = Image.open(preview_path).convert("RGB")
            image.thumbnail((THUMBNAIL_WIDTH, THUMBNAIL_HEIGHT))
            row = idx // columns
            col = idx % columns
            x = margin + col * (card_width + gutter)
            y = header_height + row * (card_height + gutter)

            shadow_box = (x + 4, y + 6, x + card_width + 4, y + card_height + 6)
            draw.rounded_rectangle(shadow_box, radius=20, fill=(232, 235, 240))
            card_box = (x, y, x + card_width, y + card_height)
            draw.rounded_rectangle(
                card_box,
                radius=20,
                fill=_rgb_tuple(self.theme.colors.surface),
                outline=_rgb_tuple(self.theme.colors.line),
                width=2,
            )

            image_x = x + (card_width - image.width) // 2
            image_y = y + 18
            sheet.paste(image, (image_x, image_y))

            badge_box = (x + 16, y + 14, x + 58, y + 40)
            draw.rounded_rectangle(badge_box, radius=12, fill=_rgb_tuple(self.theme.colors.accent))
            draw.text((x + 28, y + 19), f"{idx + 1:02d}", fill=_rgb_tuple(self.theme.colors.surface), font=card_meta_font)

            slide_title = _truncate_text(slides[idx].title or slides[idx].type.value, max_chars=28)
            slide_type = slides[idx].type.value.replace("_", " ")
            draw.text((x + 16, y + THUMBNAIL_HEIGHT + 24), slide_title, fill=_rgb_tuple(self.theme.colors.navy), font=card_title_font)
            draw.text((x + 16, y + THUMBNAIL_HEIGHT + 46), slide_type, fill=_rgb_tuple(self.theme.colors.muted), font=card_meta_font)

            review = slide_review_lookup.get(idx + 1)
            if review:
                status = str(review.get("status") or "ok")
                risk_level = str(review.get("risk_level") or "low")
                issue_count = len(review.get("issues") or [])
                badge_fill = _risk_badge_fill(status=status, risk_level=risk_level)
                badge_text = f"{risk_level.upper()} • {issue_count}"
                badge_box = (x + card_width - 120, y + THUMBNAIL_HEIGHT + 38, x + card_width - 16, y + THUMBNAIL_HEIGHT + 60)
                draw.rounded_rectangle(badge_box, radius=10, fill=badge_fill)
                draw.text((badge_box[0] + 10, badge_box[1] + 4), badge_text, fill=(255, 255, 255), font=card_meta_font)

                overflow_regions = review.get("likely_overflow_regions") or []
                if overflow_regions:
                    overflow_text = _truncate_text(
                        ", ".join(str(region) for region in overflow_regions),
                        max_chars=34,
                    )
                    draw.text(
                        (x + 16, y + THUMBNAIL_HEIGHT + 64),
                        overflow_text,
                        fill=_rgb_tuple(self.theme.colors.muted),
                        font=_load_font(11),
                    )

        destination = Path(output_path)
        destination.parent.mkdir(parents=True, exist_ok=True)
        sheet.save(destination)
        return destination

    def build_visual_regression_report(
        self,
        preview_paths: list[str],
        output_dir: Path,
        base: str,
        *,
        current_preview_source: str | None = None,
        current_manifest: dict[str, object] | None = None,
        current_manifest_path: str | Path | None = None,
    ) -> dict[str, object] | None:
        if self.baseline_dir is None:
            return None

        baseline_images, baseline_manifest = _resolve_preview_sequence(self.baseline_dir)
        baseline_manifest_path = _preview_manifest_path(self.baseline_dir)
        resolved_current_manifest_path = (
            Path(current_manifest_path)
            if current_manifest_path is not None
            else _preview_manifest_path(output_dir)
        )
        current_source = current_preview_source or (
            str(current_manifest.get("preview_source")) if current_manifest and current_manifest.get("preview_source") else None
        )
        baseline_source = (
            str(baseline_manifest.get("preview_source"))
            if baseline_manifest and baseline_manifest.get("preview_source")
            else None
        )
        current_real_preview = _is_real_preview_source(current_source)
        baseline_real_preview = _is_real_preview_source(baseline_source)
        if self.require_real_previews and not current_real_preview:
            raise ValueError("visual regression requires the current preview set to be based on a real PPTX preview")
        if self.require_real_previews and not baseline_real_preview:
            raise ValueError("visual regression requires the baseline preview set to be based on a real PPTX preview")

        compared = min(len(preview_paths), len(baseline_images))
        diff_images: list[str] = []
        slide_reports: list[dict[str, object]] = []
        diff_count = 0
        missing_baseline_count = max(0, len(preview_paths) - len(baseline_images))
        extra_baseline_count = max(0, len(baseline_images) - len(preview_paths))
        source_mismatch = bool(current_source and baseline_source and current_source != baseline_source)

        for index in range(compared):
            generated_path = Path(preview_paths[index])
            baseline_path = baseline_images[index]
            generated_image = Image.open(generated_path)
            baseline_image = Image.open(baseline_path)
            mean_diff, changed_ratio, diff_image = _normalized_diff_score(generated_image, baseline_image)
            regression = mean_diff > self.diff_threshold or changed_ratio > self.diff_threshold
            if regression:
                diff_count += 1

            diff_image_path: str | None = None
            if self.write_diff_images:
                target_path = output_dir / f"{base}-diff-{index + 1:02d}.png"
                diff_image.save(target_path)
                diff_image_path = str(target_path)
                diff_images.append(diff_image_path)

            slide_reports.append(
                {
                    "slide_number": index + 1,
                    "generated_path": str(generated_path),
                    "baseline_path": str(baseline_path),
                    "mean_diff": round(mean_diff, 6),
                    "changed_ratio": round(changed_ratio, 6),
                    "regression": regression,
                    "diff_image": diff_image_path,
                }
            )

        status = "ok"
        if diff_count or missing_baseline_count or extra_baseline_count or source_mismatch:
            status = "review"

        added_slide_numbers = list(range(compared + 1, len(preview_paths) + 1))
        removed_slide_numbers = list(range(compared + 1, len(baseline_images) + 1))
        top_regressions = sorted(
            (
                {
                    "slide_number": int(slide_report["slide_number"]),
                    "generated_path": slide_report["generated_path"],
                    "baseline_path": slide_report["baseline_path"],
                    "mean_diff": slide_report["mean_diff"],
                    "changed_ratio": slide_report["changed_ratio"],
                    "diff_score": round(
                        max(float(slide_report["mean_diff"]), float(slide_report["changed_ratio"])),
                        6,
                    ),
                    "diff_image": slide_report["diff_image"],
                }
                for slide_report in slide_reports
                if slide_report["regression"]
            ),
            key=lambda item: float(item["diff_score"]),
            reverse=True,
        )[:5]
        guidance = _build_visual_regression_guidance(
            diff_count=diff_count,
            missing_baseline_count=missing_baseline_count,
            extra_baseline_count=extra_baseline_count,
            source_mismatch=source_mismatch,
            current_source=current_source,
            baseline_source=baseline_source,
        )

        return {
            "status": status,
            "baseline_dir": str(self.baseline_dir),
            "current_manifest": str(resolved_current_manifest_path) if resolved_current_manifest_path.exists() else None,
            "baseline_manifest": str(baseline_manifest_path) if baseline_manifest_path.exists() else None,
            "threshold": self.diff_threshold,
            "compared_preview_count": compared,
            "current_preview_source": current_source,
            "baseline_preview_source": baseline_source,
            "current_real_preview": current_real_preview,
            "baseline_real_preview": baseline_real_preview,
            "source_mismatch": source_mismatch,
            "baseline_preview_count": len(baseline_images),
            "diff_count": diff_count,
            "missing_baseline_count": missing_baseline_count,
            "extra_baseline_count": extra_baseline_count,
            "added_slide_numbers": added_slide_numbers,
            "removed_slide_numbers": removed_slide_numbers,
            "top_regressions": top_regressions,
            "guidance": guidance,
            "slides": slide_reports,
            "diff_images": diff_images,
        }

    def build_preview_artifact_review(self, preview_paths: list[str]) -> dict[str, object]:
        background_rgb = _rgb_tuple(self.theme.colors.background)
        strip = 16
        corner_box = 40
        edge_contact_count = 0
        safe_margin_warning_count = 0
        edge_density_warning_count = 0
        body_edge_contact_count = 0
        safe_area_intrusion_count = 0
        footer_intrusion_count = 0
        corner_density_warning_count = 0
        slides: list[dict[str, object]] = []

        for index, preview_path in enumerate(preview_paths, start=1):
            image = Image.open(preview_path).convert("RGB")
            width, height = image.size
            bbox = _foreground_bbox_against_background(image, background_rgb)
            if bbox is None:
                bbox = (0, 0, 0, 0)

            footer_line_y = self._y(self.theme.grid.footer_line_y)
            safe_left = self._x(self.theme.grid.content_left)
            safe_right = self._x(self.theme.grid.content_right)
            safe_top = self._y(self.theme.canvas.margin_top)
            safe_bottom = min(self._y(self.theme.canvas.height - self.theme.canvas.margin_bottom), footer_line_y)

            body_box = (0, 0, width, max(1, footer_line_y))
            body_bbox = _foreground_bbox_against_background(image, background_rgb, box=body_box)
            if body_bbox is None:
                body_bbox = (0, 0, 0, 0)

            left_ratio = _changed_ratio_against_background(image, background_rgb, box=(0, 0, strip, height))
            right_ratio = _changed_ratio_against_background(image, background_rgb, box=(width - strip, 0, width, height))
            top_ratio = _changed_ratio_against_background(image, background_rgb, box=(0, 0, width, strip))
            bottom_ratio = _changed_ratio_against_background(image, background_rgb, box=(0, height - strip, width, height))
            max_edge_ratio = max(left_ratio, right_ratio, top_ratio, bottom_ratio)

            body_left_ratio = _changed_ratio_against_background(image, background_rgb, box=(0, 0, strip, footer_line_y))
            body_right_ratio = _changed_ratio_against_background(
                image,
                background_rgb,
                box=(width - strip, 0, width, footer_line_y),
            )
            body_top_ratio = _changed_ratio_against_background(image, background_rgb, box=(0, 0, width, strip))
            body_bottom_ratio = _changed_ratio_against_background(
                image,
                background_rgb,
                box=(0, max(0, footer_line_y - strip), width, footer_line_y),
            )
            body_max_edge_ratio = max(body_left_ratio, body_right_ratio, body_top_ratio, body_bottom_ratio)

            corner_ratios = {
                "top_left": _changed_ratio_against_background(image, background_rgb, box=(0, 0, corner_box, corner_box)),
                "top_right": _changed_ratio_against_background(
                    image,
                    background_rgb,
                    box=(width - corner_box, 0, width, corner_box),
                ),
                "bottom_left": _changed_ratio_against_background(
                    image,
                    background_rgb,
                    box=(0, max(0, footer_line_y - corner_box), corner_box, footer_line_y),
                ),
                "bottom_right": _changed_ratio_against_background(
                    image,
                    background_rgb,
                    box=(width - corner_box, max(0, footer_line_y - corner_box), width, footer_line_y),
                ),
            }
            max_corner_ratio = max(corner_ratios.values())

            edge_contact = bool(bbox != (0, 0, 0, 0) and (bbox[0] <= 8 or bbox[1] <= 8 or bbox[2] >= width - 8 or bbox[3] >= height - 8))
            safe_margin_warning = bool(
                bbox != (0, 0, 0, 0)
                and (bbox[0] <= 24 or bbox[1] <= 24 or bbox[2] >= width - 24 or bbox[3] >= height - 24)
            )
            edge_density_warning = max_edge_ratio >= 0.015
            body_edge_contact = bool(
                body_bbox != (0, 0, 0, 0)
                and (body_bbox[0] <= 8 or body_bbox[1] <= 8 or body_bbox[2] >= width - 8 or body_bbox[3] >= footer_line_y - 8)
            )
            safe_area_intrusion = bool(
                body_bbox != (0, 0, 0, 0)
                and (
                    body_bbox[0] < safe_left - 8
                    or body_bbox[1] < safe_top - 8
                    or body_bbox[2] > safe_right + 8
                    or body_bbox[3] > safe_bottom + 8
                )
            )
            footer_intrusion_warning = bool(body_bbox != (0, 0, 0, 0) and body_bbox[3] >= footer_line_y - 24)
            corner_density_warning = max_corner_ratio >= 0.05 or body_max_edge_ratio >= 0.03

            issues: list[str] = []
            if edge_contact:
                edge_contact_count += 1
                issues.append("foreground touches the slide edge")
            elif safe_margin_warning:
                safe_margin_warning_count += 1
                issues.append("foreground approaches the slide edge")
            if edge_density_warning:
                edge_density_warning_count += 1
                issues.append("edge density suggests possible clipping or aggressive packing")
            if body_edge_contact:
                body_edge_contact_count += 1
                issues.append("body content touches a slide boundary before the footer region")
            if safe_area_intrusion:
                safe_area_intrusion_count += 1
                issues.append("body content appears to intrude into an unsafe margin region")
            if footer_intrusion_warning:
                footer_intrusion_count += 1
                issues.append("body content approaches the footer boundary and may clip or collide")
            if corner_density_warning:
                corner_density_warning_count += 1
                issues.append("corner density suggests content is packed into unsafe corners")

            slides.append(
                {
                    "slide_number": index,
                    "preview_path": preview_path,
                    "foreground_bbox": {
                        "left": int(bbox[0]),
                        "top": int(bbox[1]),
                        "right": int(bbox[2]),
                        "bottom": int(bbox[3]),
                    },
                    "body_foreground_bbox": {
                        "left": int(body_bbox[0]),
                        "top": int(body_bbox[1]),
                        "right": int(body_bbox[2]),
                        "bottom": int(body_bbox[3]),
                    },
                    "edge_contact": edge_contact,
                    "safe_margin_warning": safe_margin_warning,
                    "edge_density_warning": edge_density_warning,
                    "body_edge_contact": body_edge_contact,
                    "safe_area_intrusion": safe_area_intrusion,
                    "footer_intrusion_warning": footer_intrusion_warning,
                    "corner_density_warning": corner_density_warning,
                    "max_edge_ratio": round(max_edge_ratio, 6),
                    "body_max_edge_ratio": round(body_max_edge_ratio, 6),
                    "max_corner_ratio": round(max_corner_ratio, 6),
                    "issues": issues,
                }
            )

        status = "review" if (
            edge_contact_count
            or safe_margin_warning_count
            or edge_density_warning_count
            or body_edge_contact_count
            or safe_area_intrusion_count
            or footer_intrusion_count
            or corner_density_warning_count
        ) else "ok"
        return {
            "status": status,
            "preview_count": len(preview_paths),
            "edge_contact_count": edge_contact_count,
            "safe_margin_warning_count": safe_margin_warning_count,
            "edge_density_warning_count": edge_density_warning_count,
            "body_edge_contact_count": body_edge_contact_count,
            "safe_area_intrusion_count": safe_area_intrusion_count,
            "footer_intrusion_count": footer_intrusion_count,
            "corner_density_warning_count": corner_density_warning_count,
            "slides": slides,
        }

    def _x(self, value_in_inches: float) -> int:
        return int((value_in_inches / self.theme.canvas.width) * PREVIEW_WIDTH)

    def _y(self, value_in_inches: float) -> int:
        return int((value_in_inches / self.theme.canvas.height) * PREVIEW_HEIGHT)

    def _render_debug_overlay(self, draw: ImageDraw.ImageDraw) -> None:
        grid = self.theme.grid
        canvas = self.theme.canvas

        if self.debug_safe_areas:
            safe_box = (
                self._x(grid.content_left),
                self._y(canvas.margin_top),
                self._x(grid.content_right),
                self._y(canvas.height - canvas.margin_bottom),
            )
            draw.rounded_rectangle(safe_box, radius=10, outline=(89, 124, 250), width=3)
            draw.line(
                (
                    self._x(grid.content_left),
                    self._y(grid.footer_line_y),
                    self._x(grid.content_right),
                    self._y(grid.footer_line_y),
                ),
                fill=(89, 124, 250),
                width=2,
            )

        if self.debug_grid:
            debug_lines = [
                ("header_top", self._y(grid.header_top)),
                ("title_top", self._y(grid.title_top)),
                ("body_top", self._y(grid.body_top)),
            ]
            for label, y in debug_lines:
                draw.line((40, y, PREVIEW_WIDTH - 40, y), fill=(91, 167, 120), width=1)
                draw.text((44, y - 16), label, fill=(91, 167, 120), font=_load_font(12))

            x_guides = [
                ("content_left", self._x(grid.content_left)),
                ("content_right", self._x(grid.content_right)),
                ("side_panel_left", self._x(grid.side_panel_left)),
                ("image_left", self._x(grid.image_left)),
            ]
            for label, x in x_guides:
                draw.line((x, 30, x, PREVIEW_HEIGHT - 30), fill=(91, 167, 120), width=1)
                draw.text((x + 6, 34), label, fill=(91, 167, 120), font=_load_font(12))

    def build_preview_quality_review(self, spec: PresentationInput) -> dict[str, object]:
        return review_presentation(spec, asset_root=self.asset_root, theme_name=self.theme.name)

    def _render_heading(self, draw: ImageDraw.ImageDraw, slide_spec: Slide, *, eyebrow_default: str | None = None) -> int:
        colors = self.theme.colors
        eyebrow_font = _load_font(20, bold=True)
        title_font = _load_font(44, bold=True)
        subtitle_font = _load_font(24)

        top = 74
        eyebrow = slide_spec.eyebrow or eyebrow_default
        if eyebrow:
            draw.text((92, top), eyebrow.upper(), fill=_rgb_tuple(colors.accent), font=eyebrow_font)
            top += 34

        draw.text((92, top), slide_spec.title or "", fill=_rgb_tuple(colors.navy), font=title_font)
        top += 58
        if slide_spec.subtitle:
            draw.text((92, top), slide_spec.subtitle, fill=_rgb_tuple(colors.muted), font=subtitle_font)
            top += 38
        return top

    def _render_footer(self, draw: ImageDraw.ImageDraw, meta: PresentationMeta, index: int, total_slides: int) -> None:
        colors = self.theme.colors
        font = _load_font(16)
        draw.line((92, 660, 1188, 660), fill=_rgb_tuple(colors.line), width=2)
        left_text = meta.footer_text or meta.client_name or meta.author or meta.title
        if meta.date:
            left_text = f"{left_text} • {meta.date}" if left_text else meta.date
        draw.text((92, 672), left_text or "", fill=_rgb_tuple(colors.muted), font=font)
        draw.text((1080, 672), f"{index:02d} / {total_slides:02d}", fill=_rgb_tuple(colors.muted), font=font)

    def _draw_panel(self, draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int]) -> None:
        colors = self.theme.colors
        draw.rounded_rectangle(box, radius=18, fill=_rgb_tuple(colors.surface), outline=_rgb_tuple(colors.line), width=2)

    def _draw_text_block(
        self,
        draw: ImageDraw.ImageDraw,
        text: str,
        box: tuple[int, int, int, int],
        *,
        font: ImageFont.ImageFont,
        fill: tuple[int, int, int],
        line_gap: int = 8,
    ) -> None:
        lines = _wrap_text(draw, text, font, box[2] - box[0])
        y = box[1]
        for line in lines:
            draw.text((box[0], y), line, fill=fill, font=font)
            y += font.size + line_gap if hasattr(font, "size") else 20 + line_gap
            if y > box[3]:
                break

    def _render_title(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        colors = self.theme.colors
        variant = slide_spec.layout_variant or "split_panel"
        asset = self.resolve_asset(slide_spec.image_path)
        focal_x, focal_y = infer_contextual_image_focal_point(slide_spec)
        if variant == "hero_cover":
            draw.rectangle((92, 72, 1188, 82), fill=_rgb_tuple(colors.accent))
            top = self._render_heading(draw, slide_spec, eyebrow_default=meta.client_name or meta.subtitle)
            if slide_spec.body:
                body_font = _load_font(24)
                self._draw_text_block(draw, slide_spec.body, (92, top + 20, 780, 470), font=body_font, fill=_rgb_tuple(colors.text))
            if asset:
                loaded = Image.open(asset).convert("RGB")
                resampling = getattr(getattr(Image, "Resampling", Image), "LANCZOS")
                fitted = ImageOps.fit(
                    loaded,
                    (288, 310),
                    method=resampling,
                    centering=(focal_x, focal_y),
                )
                image.paste(fitted, (900, 150))
                self._draw_panel(draw, (916, 352, 1172, 452))
            else:
                self._draw_panel(draw, (900, 150, 1188, 460))
            panel_title = _load_font(18, bold=True)
            body_font = _load_font(22, bold=True)
            small_font = _load_font(18)
            base_x = 934 if asset else 928
            base_y = 366 if asset else 178
            draw.text((base_x, base_y), "CONTEXT", fill=_rgb_tuple(colors.muted), font=panel_title)
            if meta.client_name:
                draw.text((base_x, base_y + 30), meta.client_name, fill=_rgb_tuple(colors.navy), font=body_font)
            if meta.author:
                draw.text((base_x, base_y + 72), meta.author, fill=_rgb_tuple(colors.text), font=small_font)
            if meta.date:
                draw.text((base_x, base_y + 98), meta.date, fill=_rgb_tuple(colors.muted), font=small_font)
            draw.text((base_x, base_y + 124), self.theme.name.replace("_", " ").title(), fill=_rgb_tuple(colors.accent), font=small_font)
            return

        draw.rectangle((92, 88, 104, 236), fill=_rgb_tuple(colors.accent))
        top = self._render_heading(draw, slide_spec, eyebrow_default=meta.subtitle)
        if slide_spec.body:
            body_font = _load_font(24)
            self._draw_text_block(draw, slide_spec.body, (128, top + 30, 720, 470), font=body_font, fill=_rgb_tuple(colors.text))
        if asset:
            loaded = Image.open(asset).convert("RGB")
            resampling = getattr(getattr(Image, "Resampling", Image), "LANCZOS")
            fitted = ImageOps.fit(
                loaded,
                (310, 350),
                method=resampling,
                centering=(focal_x, focal_y),
            )
            image.paste(fitted, (860, 120))
            self._draw_panel(draw, (880, 356, 1150, 456))
        else:
            self._draw_panel(draw, (860, 120, 1170, 470))
        small_font = _load_font(16, bold=True)
        body_font = _load_font(22, bold=True)
        text_font = _load_font(18)
        y = 370 if asset else 152
        x = 898 if asset else 890
        draw.text((x, y), "DECK", fill=_rgb_tuple(colors.accent), font=small_font)
        y += 36
        draw.text((x, y), meta.title, fill=_rgb_tuple(colors.navy), font=body_font)
        y += 52
        if meta.client_name:
            draw.text((x, y), "Client", fill=_rgb_tuple(colors.muted), font=small_font)
            y += 28
            draw.text((x, y), meta.client_name, fill=_rgb_tuple(colors.text), font=text_font)
            y += 40
        if meta.author:
            draw.text((x, y), "Author", fill=_rgb_tuple(colors.muted), font=small_font)
            y += 28
            draw.text((x, y), meta.author, fill=_rgb_tuple(colors.text), font=text_font)

    def _render_section(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        colors = self.theme.colors
        title_font = _load_font(54, bold=True)
        label_font = _load_font(22, bold=True)
        subtitle_font = _load_font(24)
        if slide_spec.section_label:
            draw.text((92, 170), slide_spec.section_label.upper(), fill=_rgb_tuple(colors.accent), font=label_font)
        draw.text((92, 245), slide_spec.title or "", fill=_rgb_tuple(colors.navy), font=title_font)
        if slide_spec.subtitle:
            draw.text((92, 320), slide_spec.subtitle, fill=_rgb_tuple(colors.muted), font=subtitle_font)
        draw.rectangle((92, 382, 430, 392), fill=_rgb_tuple(colors.accent))

    def _render_agenda(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        top = self._render_heading(draw, slide_spec, eyebrow_default="Agenda") + 18
        if slide_spec.body:
            self._draw_text_block(draw, slide_spec.body, (92, top, 1188, top + 70), font=_load_font(22), fill=_rgb_tuple(self.theme.colors.text))
            top += 82
        for idx, bullet in enumerate(slide_spec.bullets, start=1):
            box = (92, top, 1188, top + 50)
            self._draw_panel(draw, box)
            draw.rectangle((92, top, 106, top + 50), fill=_rgb_tuple(self.theme.colors.accent if idx == 1 else self.theme.colors.navy))
            draw.text((128, top + 12), f"{idx:02d}", fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(18, bold=True))
            draw.text((186, top + 12), bullet, fill=_rgb_tuple(self.theme.colors.text), font=_load_font(20, bold=idx == 1))
            top += 66

    def _render_bullets(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        top = self._render_heading(draw, slide_spec) + 18
        if slide_spec.body:
            self._draw_text_block(draw, slide_spec.body, (92, top, 780, top + 110), font=_load_font(22), fill=_rgb_tuple(self.theme.colors.text))
            top += 100
        for bullet in slide_spec.bullets:
            draw.text((112, top), f"• {bullet}", fill=_rgb_tuple(self.theme.colors.text), font=_load_font(22))
            top += 42

    def _render_cards(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec)
        positions = [(92, 210, 430, 470), (472, 210, 810, 470), (852, 210, 1190, 470)]
        for box, card in zip(positions, slide_spec.cards, strict=True):
            self._draw_panel(draw, box)
            draw.text((box[0] + 24, box[1] + 24), card.title, fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(24, bold=True))
            self._draw_text_block(draw, card.body, (box[0] + 24, box[1] + 70, box[2] - 24, box[3] - 60), font=_load_font(18), fill=_rgb_tuple(self.theme.colors.text))
            if card.footer:
                draw.text((box[0] + 24, box[3] - 34), card.footer, fill=_rgb_tuple(self.theme.colors.accent), font=_load_font(16, bold=True))

    def _render_metrics(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec)
        metrics = slide_spec.metrics
        gap = 24
        width = (1096 - gap * (len(metrics) - 1)) // len(metrics)
        for idx, metric in enumerate(metrics):
            left = 92 + idx * (width + gap)
            box = (left, 250, left + width, 450)
            self._draw_panel(draw, box)
            draw.rectangle((left, 250, left + width, 260), fill=_rgb_tuple(self.theme.colors.navy if idx % 2 == 0 else self.theme.colors.accent))
            draw.text((left + 24, 288), metric.value, fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(30, bold=True))
            draw.text((left + 24, 336), metric.label, fill=_rgb_tuple(self.theme.colors.text), font=_load_font(18, bold=True))
            if metric.trend:
                draw.text((left + 24, 402), metric.trend, fill=_rgb_tuple(self.theme.colors.accent), font=_load_font(16, bold=True))

    def _render_chart(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        top = self._render_heading(draw, slide_spec, eyebrow_default="Data view") + 24
        if slide_spec.body:
            self._draw_text_block(draw, slide_spec.body, (92, top, 1188, top + 70), font=_load_font(20), fill=_rgb_tuple(self.theme.colors.text))
            top += 80
        box = (92, top, 1188, 560)
        self._draw_panel(draw, box)
        self._draw_chart_area(draw, slide_spec, box)

    def _draw_chart_area(self, draw: ImageDraw.ImageDraw, slide_spec: Slide, box: tuple[int, int, int, int]) -> None:
        colors = self.theme.colors
        palette = [colors.navy, colors.accent, colors.text, colors.muted]
        left, top, right, bottom = box
        chart_left = left + 56
        chart_top = top + 30
        chart_right = right - 36
        chart_bottom = bottom - 64

        draw.line((chart_left, chart_top, chart_left, chart_bottom), fill=_rgb_tuple(colors.line), width=2)
        draw.line((chart_left, chart_bottom, chart_right, chart_bottom), fill=_rgb_tuple(colors.line), width=2)

        min_value = min(min(series.values) for series in slide_spec.chart_series)
        max_value = max(max(series.values) for series in slide_spec.chart_series)
        category_count = len(slide_spec.chart_categories)
        series_count = len(slide_spec.chart_series)
        variant = slide_spec.layout_variant or "column"
        label_font = _load_font(15)
        usable_height = max(1, chart_bottom - chart_top - 18)
        usable_width = max(1, chart_right - chart_left - 20)
        value_range = max(max_value - min_value, 1e-9)

        def _value_to_y(value: float) -> float:
            return chart_bottom - (((value - min_value) / value_range) * usable_height)

        def _value_to_x(value: float) -> float:
            return chart_left + (((value - min_value) / value_range) * usable_width)

        if min_value >= 0:
            zero_y = chart_bottom
            zero_x = chart_left
        elif max_value <= 0:
            zero_y = chart_top
            zero_x = chart_right
        else:
            zero_y = _value_to_y(0.0)
            zero_x = _value_to_x(0.0)

        if min_value < 0 < max_value:
            draw.line((chart_left, zero_y, chart_right, zero_y), fill=_rgb_tuple(colors.line), width=1)
            draw.line((zero_x, chart_top, zero_x, chart_bottom), fill=_rgb_tuple(colors.line), width=1)

        if variant == "line":
            step_x = (chart_right - chart_left) / max(category_count - 1, 1)
            for series_index, series in enumerate(slide_spec.chart_series):
                points: list[tuple[float, float]] = []
                for idx, value in enumerate(series.values):
                    x = chart_left + idx * step_x
                    y = _value_to_y(value)
                    points.append((x, y))
                draw.line(points, fill=_rgb_tuple(palette[series_index % len(palette)]), width=4)
                for x, y in points:
                    draw.ellipse((x - 5, y - 5, x + 5, y + 5), fill=_rgb_tuple(palette[series_index % len(palette)]))
            for idx, category in enumerate(slide_spec.chart_categories):
                x = chart_left + idx * step_x
                draw.text((x - 14, chart_bottom + 12), category, fill=_rgb_tuple(colors.muted), font=label_font)
            return

        category_width = (chart_right - chart_left) / category_count
        if variant == "bar":
            row_height = (chart_bottom - chart_top) / category_count
            for idx, category in enumerate(slide_spec.chart_categories):
                y = chart_top + idx * row_height + 10
                draw.text((chart_left - 48, y), category, fill=_rgb_tuple(colors.muted), font=label_font)
                for series_index, series in enumerate(slide_spec.chart_series):
                    value = series.values[idx]
                    value_x = _value_to_x(value)
                    bar_top = y + series_index * 14
                    x_start, x_end = sorted((zero_x, value_x))
                    draw.rectangle(
                        (x_start, bar_top, x_end, bar_top + 10),
                        fill=_rgb_tuple(palette[series_index % len(palette)]),
                    )
            return

        usable_width = category_width * 0.75
        series_bar_width = usable_width / series_count
        for idx, category in enumerate(slide_spec.chart_categories):
            base_x = chart_left + idx * category_width + (category_width - usable_width) / 2
            draw.text((base_x, chart_bottom + 12), category, fill=_rgb_tuple(colors.muted), font=label_font)
            for series_index, series in enumerate(slide_spec.chart_series):
                value = series.values[idx]
                x1 = base_x + series_index * series_bar_width
                x2 = x1 + series_bar_width - 6
                value_y = _value_to_y(value)
                y1, y2 = sorted((value_y, zero_y))
                draw.rectangle((x1, y1, x2, y2), fill=_rgb_tuple(palette[series_index % len(palette)]))

    def _render_image_text(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec)
        self._draw_text_block(draw, slide_spec.body or "", (92, 230, 610, 350), font=_load_font(20), fill=_rgb_tuple(self.theme.colors.text))
        top = 340
        for bullet in slide_spec.bullets:
            draw.text((112, top), f"• {bullet}", fill=_rgb_tuple(self.theme.colors.text), font=_load_font(18))
            top += 36

        box = (760, 220, 1188, 520)
        asset = self.resolve_asset(slide_spec.image_path)
        focal_x, focal_y = infer_contextual_image_focal_point(slide_spec)
        if asset:
            loaded = Image.open(asset).convert("RGB")
            resampling = getattr(getattr(Image, "Resampling", Image), "LANCZOS")
            fitted = ImageOps.fit(
                loaded,
                (box[2] - box[0], box[3] - box[1]),
                method=resampling,
                centering=(focal_x, focal_y),
            )
            image.paste(fitted, (box[0], box[1]))
        else:
            self._draw_panel(draw, box)
            inner_box = (box[0] + 44, box[1] + 46, box[2] - 44, box[1] + 182)
            draw.rounded_rectangle(inner_box, radius=16, fill=_rgb_tuple(self.theme.colors.soft_fill), outline=_rgb_tuple(self.theme.colors.line), width=2)
            draw.line((inner_box[0] + 16, inner_box[1] + 16, inner_box[2] - 16, inner_box[3] - 16), fill=_rgb_tuple(self.theme.colors.line), width=3)
            draw.line((inner_box[2] - 16, inner_box[1] + 16, inner_box[0] + 16, inner_box[3] - 16), fill=_rgb_tuple(self.theme.colors.line), width=3)
            draw.text((box[0] + 66, box[1] + 18), "VISUAL PLACEHOLDER", fill=_rgb_tuple(self.theme.colors.accent), font=_load_font(15, bold=True))
            draw.text((box[0] + 78, box[1] + 214), "Image unavailable", fill=_rgb_tuple(self.theme.colors.muted), font=_load_font(24, bold=True))
            if slide_spec.image_path:
                draw.text((box[0] + 54, box[1] + 294), f"Missing asset: {slide_spec.image_path}", fill=_rgb_tuple(self.theme.colors.muted), font=_load_font(16))
            draw.text(
                (box[0] + 54, box[1] + 248),
                _truncate_text(
                    slide_spec.image_caption
                    or "Use a contextual image, screenshot, or diagram once the final asset is available.",
                    max_chars=74,
                ),
                fill=_rgb_tuple(self.theme.colors.text),
                font=_load_font(15),
            )

    def _render_timeline(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec)
        y = 340
        draw.line((130, y, 1150, y), fill=_rgb_tuple(self.theme.colors.line), width=4)
        step = (1020) / max(len(slide_spec.timeline_items) - 1, 1)
        for idx, item in enumerate(slide_spec.timeline_items):
            x = 130 + idx * step
            draw.ellipse((x - 12, y - 12, x + 12, y + 12), fill=_rgb_tuple(self.theme.colors.accent if idx % 2 else self.theme.colors.navy))
            draw.text((x - 40, y + 28), item.title, fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(16, bold=True))
            if item.body:
                self._draw_text_block(draw, item.body, (x - 70, y + 56, x + 90, y + 132), font=_load_font(14), fill=_rgb_tuple(self.theme.colors.text))

    def _render_comparison(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec)
        for idx, column in enumerate(slide_spec.comparison_columns):
            left = 92 + idx * 556
            box = (left, 230, left + 520, 520)
            self._draw_panel(draw, box)
            draw.rectangle((left, 230, left + 12, 520), fill=_rgb_tuple(self.theme.colors.navy if idx == 0 else self.theme.colors.accent))
            draw.text((left + 28, 252), column.title, fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(22, bold=True))
            y = 294
            if column.body:
                self._draw_text_block(draw, column.body, (left + 28, y, left + 476, y + 96), font=_load_font(18), fill=_rgb_tuple(self.theme.colors.text))
                y += 96
            for bullet in column.bullets:
                draw.text((left + 28, y), f"• {bullet}", fill=_rgb_tuple(self.theme.colors.text), font=_load_font(16))
                y += 28

    def _render_two_column(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        columns = slide_spec.two_column_columns
        comparison_like = Slide.model_construct(  # type: ignore[attr-defined]
            type=slide_spec.type,
            title=slide_spec.title,
            subtitle=slide_spec.subtitle,
            eyebrow=slide_spec.eyebrow,
            comparison_columns=columns,
        )
        self._render_comparison(draw, image, comparison_like, meta)

    def _render_table(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec, eyebrow_default="Executive table")
        left = 92
        top = 240
        total_width = 1096
        columns = len(slide_spec.table_columns)
        cell_width = total_width // columns
        for idx, column in enumerate(slide_spec.table_columns):
            box = (left + idx * cell_width, top, left + (idx + 1) * cell_width, top + 42)
            draw.rectangle(box, fill=_rgb_tuple(self.theme.colors.soft_fill), outline=_rgb_tuple(self.theme.colors.line), width=2)
            draw.text((box[0] + 16, box[1] + 10), column, fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(16, bold=True))
        current_top = top + 42
        for row in slide_spec.table_rows:
            for idx, cell in enumerate(row):
                box = (left + idx * cell_width, current_top, left + (idx + 1) * cell_width, current_top + 46)
                draw.rectangle(box, fill=_rgb_tuple(self.theme.colors.surface), outline=_rgb_tuple(self.theme.colors.line), width=2)
                draw.text((box[0] + 14, box[1] + 12), cell, fill=_rgb_tuple(self.theme.colors.text), font=_load_font(16))
            current_top += 46

    def _render_faq(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        self._render_heading(draw, slide_spec, eyebrow_default="FAQ")
        positions = [(92, 232, 620, 370), (660, 232, 1188, 370), (92, 392, 620, 530), (660, 392, 1188, 530)]
        for box, item in zip(positions, slide_spec.faq_items):
            self._draw_panel(draw, box)
            draw.text((box[0] + 22, box[1] + 18), item.title, fill=_rgb_tuple(self.theme.colors.navy), font=_load_font(18, bold=True))
            self._draw_text_block(draw, item.body, (box[0] + 22, box[1] + 52, box[2] - 22, box[3] - 18), font=_load_font(16), fill=_rgb_tuple(self.theme.colors.text))

    def _render_summary(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        top = self._render_heading(draw, slide_spec, eyebrow_default="Executive summary") + 20
        asset = self.resolve_asset(slide_spec.image_path)
        focal_x, focal_y = infer_contextual_image_focal_point(slide_spec)

        body_box = (92, top, 700, top + 120)
        bullet_panel = (820, 220, 1188, 520)
        if asset:
            loaded = Image.open(asset).convert("RGB")
            resampling = getattr(getattr(Image, "Resampling", Image), "LANCZOS")
            fitted = ImageOps.fit(
                loaded,
                (292, 320),
                method=resampling,
                centering=(focal_x, focal_y),
            )
            image.paste(fitted, (896, 188))
            body_box = (92, top, 820, top + 140)
            bullet_panel = (616, 246, 864, 520)

        if slide_spec.body:
            self._draw_text_block(draw, slide_spec.body, body_box, font=_load_font(20), fill=_rgb_tuple(self.theme.colors.text))
        if slide_spec.bullets:
            self._draw_panel(draw, bullet_panel)
            draw.text((bullet_panel[0] + 28, bullet_panel[1] + 26), slide_spec.eyebrow or "Key takeaways", fill=_rgb_tuple(self.theme.colors.muted), font=_load_font(16, bold=True))
            y = bullet_panel[1] + 70
            for bullet in slide_spec.bullets:
                draw.text((bullet_panel[0] + 28, y), f"• {bullet}", fill=_rgb_tuple(self.theme.colors.text), font=_load_font(18))
                y += 34

    def _render_closing(self, draw: ImageDraw.ImageDraw, image: Image.Image, slide_spec: Slide, meta: PresentationMeta) -> None:
        quote = slide_spec.quote or slide_spec.title or ""
        asset = self.resolve_asset(slide_spec.image_path)
        focal_x, focal_y = infer_contextual_image_focal_point(slide_spec)
        quote_box = (160, 220, 1120, 420)
        if asset:
            loaded = Image.open(asset).convert("RGB")
            resampling = getattr(getattr(Image, "Resampling", Image), "LANCZOS")
            fitted = ImageOps.fit(
                loaded,
                (328, 344),
                method=resampling,
                centering=(focal_x, focal_y),
            )
            image.paste(fitted, (844, 168))
            quote_box = (140, 220, 790, 430)
        self._draw_text_block(draw, quote, quote_box, font=_load_font(34, bold=True), fill=_rgb_tuple(self.theme.colors.navy), line_gap=14)
        if slide_spec.attribution:
            draw.text((160, 450), slide_spec.attribution, fill=_rgb_tuple(self.theme.colors.muted), font=_load_font(20))

    def resolve_asset(self, image_path: str | None) -> Path | None:
        if not image_path:
            return None
        candidate = Path(image_path)
        if not candidate.is_absolute():
            candidate = self.asset_root / candidate
        return candidate if candidate.exists() else None


def render_previews(
    spec: PresentationInput,
    output_dir: str | Path,
    *,
    theme_name: str | None = None,
    asset_root: str | Path | None = None,
    primary_color: str | None = None,
    secondary_color: str | None = None,
    basename: str | None = None,
    debug_grid: bool = False,
    debug_safe_areas: bool = False,
    backend: str = "auto",
    baseline_dir: str | Path | None = None,
    diff_threshold: float = 0.01,
    write_diff_images: bool = False,
    require_real_previews: bool = False,
) -> dict[str, object]:
    renderer = PreviewRenderer(
        theme_name=theme_name,
        asset_root=asset_root,
        primary_color=primary_color,
        secondary_color=secondary_color,
        debug_grid=debug_grid,
        debug_safe_areas=debug_safe_areas,
        backend=backend,
        baseline_dir=baseline_dir,
        diff_threshold=diff_threshold,
        write_diff_images=write_diff_images,
        require_real_previews=require_real_previews,
    )
    return renderer.render(spec, output_dir, basename=basename)


def render_previews_from_pptx(
    input_pptx: str | Path,
    output_dir: str | Path,
    *,
    theme_name: str | None = None,
    basename: str | None = None,
    baseline_dir: str | Path | None = None,
    diff_threshold: float = 0.01,
    write_diff_images: bool = False,
    require_real_previews: bool = False,
) -> dict[str, object]:
    renderer = PreviewRenderer(
        theme_name=theme_name,
        baseline_dir=baseline_dir,
        diff_threshold=diff_threshold,
        write_diff_images=write_diff_images,
        backend="office",
        require_real_previews=require_real_previews,
    )
    return renderer.render_pptx_previews(input_pptx, output_dir, basename=basename)


def review_pptx_artifact(
    input_pptx: str | Path,
    output_dir: str | Path,
    *,
    theme_name: str | None = None,
    basename: str | None = None,
    baseline_dir: str | Path | None = None,
    diff_threshold: float = 0.01,
    write_diff_images: bool = False,
    require_real_previews: bool = False,
) -> dict[str, object]:
    preview_result = render_previews_from_pptx(
        input_pptx,
        output_dir,
        theme_name=theme_name,
        basename=basename,
        baseline_dir=baseline_dir,
        diff_threshold=diff_threshold,
        write_diff_images=write_diff_images,
        require_real_previews=require_real_previews,
    )
    review = review_preview_artifacts(preview_result, input_pptx=str(Path(input_pptx).resolve()))
    return {
        "mode": "review-pptx",
        "input_pptx": str(Path(input_pptx).resolve()),
        "output_dir": str(Path(output_dir).resolve()),
        "preview_result": preview_result,
        **{key: value for key, value in review.items() if key != "mode"},
    }


def compare_preview_directories(
    current_dir: str | Path,
    baseline_dir: str | Path,
    output_dir: str | Path,
    *,
    theme_name: str | None = None,
    basename: str | None = None,
    diff_threshold: float = 0.01,
    write_diff_images: bool = False,
    require_real_previews: bool = False,
) -> dict[str, object]:
    current_root = Path(current_dir).resolve()
    baseline_root = Path(baseline_dir).resolve()
    destination = Path(output_dir).resolve()

    if not current_root.exists() or not current_root.is_dir():
        raise FileNotFoundError(f"Current preview directory not found: {current_root}")
    if not baseline_root.exists() or not baseline_root.is_dir():
        raise FileNotFoundError(f"Baseline preview directory not found: {baseline_root}")

    current_images_resolved, current_manifest = _resolve_preview_sequence(current_root)
    baseline_images, baseline_manifest = _resolve_preview_sequence(baseline_root)
    current_images = [str(path) for path in current_images_resolved]
    if not current_images:
        raise FileNotFoundError(f"No preview PNGs found in current directory: {current_root}")
    if not baseline_images:
        raise FileNotFoundError(f"No preview PNGs found in baseline directory: {baseline_root}")

    destination.mkdir(parents=True, exist_ok=True)
    base = basename or f"{_safe_basename(current_root.name)}-vs-{_safe_basename(baseline_root.name)}"
    current_label = str((current_manifest or {}).get("presentation_title") or current_root.name)
    baseline_label = str((baseline_manifest or {}).get("presentation_title") or baseline_root.name)
    renderer = PreviewRenderer(
        theme_name=theme_name,
        baseline_dir=baseline_root,
        diff_threshold=diff_threshold,
        write_diff_images=write_diff_images,
        backend="synthetic",
        require_real_previews=require_real_previews,
    )
    comparison = renderer.build_visual_regression_report(
        current_images,
        destination,
        base,
        current_preview_source=(
            str(current_manifest.get("preview_source"))
            if current_manifest and current_manifest.get("preview_source")
            else None
        ),
        current_manifest=current_manifest,
        current_manifest_path=_preview_manifest_path(current_root),
    )
    return {
        "mode": "compare-previews",
        "current_dir": str(current_root),
        "baseline_dir": str(baseline_root),
        "output_dir": str(destination),
        "current_label": current_label,
        "baseline_label": baseline_label,
        "current_preview_count": len(current_images),
        "baseline_preview_count": len(baseline_images),
        "current_manifest": str(_preview_manifest_path(current_root)) if _preview_manifest_path(current_root).exists() else None,
        "baseline_manifest": str(_preview_manifest_path(baseline_root)) if _preview_manifest_path(baseline_root).exists() else None,
        "current_preview_source": current_manifest.get("preview_source") if current_manifest else None,
        "baseline_preview_source": baseline_manifest.get("preview_source") if baseline_manifest else None,
        "comparison": comparison,
    }


def promote_preview_baseline(
    source_dir: str | Path,
    baseline_dir: str | Path,
    *,
    clean: bool = True,
    include_thumbnail_sheet: bool = True,
) -> dict[str, object]:
    source_root = Path(source_dir).resolve()
    baseline_root = Path(baseline_dir).resolve()

    if not source_root.exists() or not source_root.is_dir():
        raise FileNotFoundError(f"Preview source directory not found: {source_root}")

    preview_paths, manifest = _resolve_preview_sequence(source_root)
    if not preview_paths:
        raise FileNotFoundError(f"No preview PNGs found in source directory: {source_root}")

    cleaned = False
    if clean and baseline_root.exists():
        shutil.rmtree(baseline_root)
        cleaned = True

    baseline_root.mkdir(parents=True, exist_ok=True)
    copied_previews: list[str] = []
    copied_thumbnail_sheets: list[str] = []

    source_manifest_path = _preview_manifest_path(source_root)
    target_manifest_path = baseline_root / PREVIEW_MANIFEST_FILENAME
    if source_manifest_path.exists():
        shutil.copy2(source_manifest_path, target_manifest_path)

    for preview_path in preview_paths:
        source_path = Path(preview_path)
        target_path = baseline_root / source_path.name
        shutil.copy2(source_path, target_path)
        copied_previews.append(str(target_path))

    if include_thumbnail_sheet:
        for thumbnail_path in sorted(source_root.glob("*-thumbnails.png")):
            target_path = baseline_root / thumbnail_path.name
            shutil.copy2(thumbnail_path, target_path)
            copied_thumbnail_sheets.append(str(target_path))

    return {
        "mode": "promote-baseline",
        "source_dir": str(source_root),
        "baseline_dir": str(baseline_root),
        "clean": clean,
        "cleaned_existing_baseline": cleaned,
        "preview_count": len(copied_previews),
        "copied_previews": copied_previews,
        "copied_thumbnail_sheets": copied_thumbnail_sheets,
        "preview_manifest": str(target_manifest_path) if target_manifest_path.exists() else None,
        "preview_source": (manifest or {}).get("preview_source"),
        "real_preview": (manifest or {}).get("real_preview"),
    }


def compare_pptx_artifacts(
    before_pptx: str | Path,
    after_pptx: str | Path,
    output_dir: str | Path,
    *,
    theme_name: str | None = None,
    basename: str | None = None,
    diff_threshold: float = 0.01,
    write_diff_images: bool = False,
    require_real_previews: bool = False,
) -> dict[str, object]:
    before_path = Path(before_pptx).resolve()
    after_path = Path(after_pptx).resolve()
    if not before_path.exists():
        raise FileNotFoundError(f"Before PPTX not found: {before_path}")
    if not after_path.exists():
        raise FileNotFoundError(f"After PPTX not found: {after_path}")
    if before_path.suffix.lower() != ".pptx":
        raise ValueError(f"Before PPTX path must end with .pptx: {before_path}")
    if after_path.suffix.lower() != ".pptx":
        raise ValueError(f"After PPTX path must end with .pptx: {after_path}")

    with TemporaryDirectory(prefix="ppt_creator_compare_pptx_") as tmpdir:
        temp_root = Path(tmpdir)
        before_preview_dir = temp_root / "before"
        after_preview_dir = temp_root / "after"
        before_result = render_previews_from_pptx(
            before_path,
            before_preview_dir,
            theme_name=theme_name,
            basename="before",
            diff_threshold=diff_threshold,
            write_diff_images=False,
            require_real_previews=True,
        )
        after_result = render_previews_from_pptx(
            after_path,
            after_preview_dir,
            theme_name=theme_name,
            basename="after",
            diff_threshold=diff_threshold,
            write_diff_images=False,
            require_real_previews=True,
        )
        comparison = compare_preview_directories(
            after_preview_dir,
            before_preview_dir,
            output_dir,
            theme_name=theme_name,
            basename=basename or f"{_safe_basename(after_path.stem)}-vs-{_safe_basename(before_path.stem)}",
            diff_threshold=diff_threshold,
            write_diff_images=write_diff_images,
            require_real_previews=True,
        )

    return {
        "mode": "compare-pptx",
        "before_pptx": str(before_path),
        "after_pptx": str(after_path),
        "output_dir": str(Path(output_dir).resolve()),
        "before_preview_count": before_result["preview_count"],
        "after_preview_count": after_result["preview_count"],
        "before_preview_manifest": before_result.get("preview_manifest"),
        "after_preview_manifest": after_result.get("preview_manifest"),
        "before_office_conversion_strategy": before_result.get("office_conversion_strategy"),
        "after_office_conversion_strategy": after_result.get("office_conversion_strategy"),
        "comparison": comparison["comparison"],
    }


def render_previews_for_rendered_artifact(
    spec: PresentationInput,
    output_dir: str | Path,
    *,
    rendered_pptx: str | Path | None = None,
    theme_name: str | None = None,
    asset_root: str | Path | None = None,
    primary_color: str | None = None,
    secondary_color: str | None = None,
    basename: str | None = None,
    debug_grid: bool = False,
    debug_safe_areas: bool = False,
    backend: str = "auto",
    baseline_dir: str | Path | None = None,
    diff_threshold: float = 0.01,
    write_diff_images: bool = False,
    require_real_previews: bool = False,
) -> tuple[dict[str, object], str]:
    if rendered_pptx is None:
        return (
            render_previews(
                spec,
                output_dir,
                theme_name=theme_name,
                asset_root=asset_root,
                primary_color=primary_color,
                secondary_color=secondary_color,
                basename=basename,
                debug_grid=debug_grid,
                debug_safe_areas=debug_safe_areas,
                backend=backend,
                baseline_dir=baseline_dir,
                diff_threshold=diff_threshold,
                write_diff_images=write_diff_images,
                require_real_previews=require_real_previews,
            ),
            "spec",
        )

    if backend == "synthetic":
        return (
            render_previews(
                spec,
                output_dir,
                theme_name=theme_name,
                asset_root=asset_root,
                primary_color=primary_color,
                secondary_color=secondary_color,
                basename=basename,
                debug_grid=debug_grid,
                debug_safe_areas=debug_safe_areas,
                backend="synthetic",
                baseline_dir=baseline_dir,
                diff_threshold=diff_threshold,
                write_diff_images=write_diff_images,
                require_real_previews=require_real_previews,
            ),
            "spec",
        )

    if backend == "office":
        return (
            render_previews_from_pptx(
                rendered_pptx,
                output_dir,
                theme_name=theme_name,
                basename=basename,
                baseline_dir=baseline_dir,
                diff_threshold=diff_threshold,
                write_diff_images=write_diff_images,
                require_real_previews=require_real_previews,
            ),
            "rendered_pptx",
        )

    runtime = find_office_runtime()
    if runtime:
        try:
            return (
                render_previews_from_pptx(
                    rendered_pptx,
                    output_dir,
                    theme_name=theme_name,
                    basename=basename,
                    baseline_dir=baseline_dir,
                    diff_threshold=diff_threshold,
                    write_diff_images=write_diff_images,
                    require_real_previews=require_real_previews,
                ),
                "rendered_pptx",
            )
        except Exception:
            pass

    return (
        render_previews(
            spec,
            output_dir,
            theme_name=theme_name,
            asset_root=asset_root,
            primary_color=primary_color,
            secondary_color=secondary_color,
            basename=basename,
            debug_grid=debug_grid,
            debug_safe_areas=debug_safe_areas,
            backend="synthetic",
            baseline_dir=baseline_dir,
            diff_threshold=diff_threshold,
            write_diff_images=write_diff_images,
            require_real_previews=require_real_previews,
        ),
        "spec",
    )

from __future__ import annotations

import math
from pathlib import Path
from typing import TYPE_CHECKING

from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ppt_creator.layouts import LAYOUT_RENDERERS
from ppt_creator.schema import PresentationInput, PresentationMeta, Slide, SlideType
from ppt_creator.theme import SemanticLayoutPreset, get_theme, rgb

if TYPE_CHECKING:
    from pptx.slide import Slide as PptxSlide
    from pptx.text.text import Font, TextFrame


def infer_visual_placeholder_kind(
    *,
    slide_type: str,
    title: str | None = None,
    body: str | None = None,
    caption: str | None = None,
    image_path: str | None = None,
) -> str:
    text = " ".join(part for part in [title, body, caption, image_path] if part).lower()
    if any(keyword in text for keyword in ["screenshot", "screen", "ui", "dashboard", "mockup", "product"]):
        return "screenshot"
    if any(keyword in text for keyword in ["diagram", "workflow", "process", "architecture", "roadmap", "timeline"]):
        return "diagram"
    if any(keyword in text for keyword in ["chart", "metric", "kpi", "trend", "table", "analytics", "analysis"]):
        return "analytical_visual"
    if slide_type in {"title", "closing", "summary"}:
        return "photo"
    return "photo"


def infer_contextual_image_focal_point(slide_spec: Slide) -> tuple[float, float]:
    kind = infer_visual_placeholder_kind(
        slide_type=slide_spec.type.value,
        title=slide_spec.title,
        body=slide_spec.body,
        caption=slide_spec.image_caption,
        image_path=slide_spec.image_path,
    )
    fallback_by_kind = {
        "photo": (0.52, 0.38),
        "screenshot": (0.50, 0.46),
        "diagram": (0.50, 0.42),
        "analytical_visual": (0.50, 0.44),
    }
    type_defaults = {
        SlideType.TITLE: {
            "photo": (0.52, 0.32 if slide_spec.layout_variant == "hero_cover" else 0.38),
            "screenshot": (0.50, 0.40),
            "diagram": (0.50, 0.36),
            "analytical_visual": (0.50, 0.40),
        },
        SlideType.IMAGE_TEXT: {
            "photo": (0.52, 0.40),
            "screenshot": (0.50, 0.46),
            "diagram": (0.50, 0.42),
            "analytical_visual": (0.50, 0.45),
        },
        SlideType.SUMMARY: {
            "photo": (0.54, 0.34),
            "screenshot": (0.50, 0.42),
            "diagram": (0.50, 0.38),
            "analytical_visual": (0.50, 0.41),
        },
        SlideType.CLOSING: {
            "photo": (0.55, 0.34),
            "screenshot": (0.50, 0.42),
            "diagram": (0.50, 0.38),
            "analytical_visual": (0.50, 0.40),
        },
    }
    default_x, default_y = type_defaults.get(slide_spec.type, {}).get(kind, fallback_by_kind[kind])
    return (
        slide_spec.image_focal_x if slide_spec.image_focal_x is not None else default_x,
        slide_spec.image_focal_y if slide_spec.image_focal_y is not None else default_y,
    )


class PresentationRenderer:
    def __init__(
        self,
        theme_name: str | None = None,
        asset_root: str | Path | None = None,
        primary_color: str | None = None,
        secondary_color: str | None = None,
    ):
        self.requested_theme_name = theme_name
        self.requested_primary_color = primary_color
        self.requested_secondary_color = secondary_color
        self.theme = get_theme(
            theme_name,
            primary_color=primary_color,
            secondary_color=secondary_color,
        )
        self.asset_root = Path(asset_root or ".").resolve()

    def render(self, spec: PresentationInput, output_path: str | Path) -> Path:
        self.theme = get_theme(
            self.requested_theme_name or spec.presentation.theme or self.theme.name,
            primary_color=self.requested_primary_color or spec.presentation.primary_color,
            secondary_color=self.requested_secondary_color or spec.presentation.secondary_color,
        )
        destination = self.validate_output_path(output_path)
        presentation = Presentation()
        presentation.slide_width = Inches(self.theme.canvas.width)
        presentation.slide_height = Inches(self.theme.canvas.height)

        blank_layout = presentation.slide_layouts[6]
        total_slides = len(spec.slides)

        for index, slide_spec in enumerate(spec.slides, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            self.apply_background(slide)
            self.render_slide(slide, slide_spec, spec.presentation, index, total_slides)
            self.add_footer(slide, spec.presentation, index, total_slides)
            self.add_speaker_notes(slide, slide_spec.speaker_notes)

        destination.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(destination))
        return destination

    def validate_output_path(self, output_path: str | Path) -> Path:
        destination = Path(output_path)
        if destination.suffix.lower() != ".pptx":
            raise ValueError(f"Output path must end with .pptx: {destination}")
        return destination

    def content_bounds(self) -> tuple[float, float, float]:
        grid = self.theme.grid
        return grid.content_left, grid.content_right, grid.content_width

    def resolve_semantic_layout(
        self,
        slide_type: SlideType | str | None = None,
        layout_variant: str | None = None,
    ) -> SemanticLayoutPreset:
        normalized_type = slide_type.value if isinstance(slide_type, SlideType) else str(slide_type or "").strip().lower()
        normalized_variant = str(layout_variant or "").strip().lower().replace("-", "_")

        if normalized_type == SlideType.TITLE.value:
            if normalized_variant == "hero_cover":
                return SemanticLayoutPreset(
                    heading_top=1.55,
                    body_top=3.55,
                    panel_top=1.60,
                    panel_title_height=0.44,
                )
            return SemanticLayoutPreset(
                heading_top=1.25,
                body_top=3.35,
                panel_top=1.60,
                panel_title_height=0.42,
            )
        if normalized_type == SlideType.SECTION.value:
            return SemanticLayoutPreset(
                heading_top=2.40,
                body_top=3.05,
                panel_top=2.15,
                footer_boundary=5.65,
                panel_title_height=0.42,
            )
        if normalized_type == SlideType.METRICS.value:
            return SemanticLayoutPreset(
                heading_top=1.02,
                body_top=2.42,
                panel_top=2.42,
                panel_title_height=0.42,
            )
        if normalized_type in {
            SlideType.COMPARISON.value,
            SlideType.TWO_COLUMN.value,
            SlideType.FAQ.value,
            SlideType.SUMMARY.value,
        }:
            return SemanticLayoutPreset(
                heading_top=1.02,
                body_top=2.24,
                panel_top=2.24,
                panel_title_height=0.43,
            )
        if normalized_type == SlideType.TABLE.value:
            return SemanticLayoutPreset(
                heading_top=1.02,
                body_top=2.10,
                panel_top=2.10,
                panel_title_height=0.40,
            )
        return SemanticLayoutPreset()

    def _build_named_region_map(
        self,
        *,
        regions: list[dict[str, float | str]],
        bounds: list[object],
        default_prefix: str,
    ) -> dict[str, object]:
        named: dict[str, object] = {}
        for index, (region, bound) in enumerate(zip(regions, bounds, strict=True), start=1):
            kind = str(region.get("kind") or f"{default_prefix}_{index}")
            named[kind] = bound
        return named

    def resolve_layout_variant(self, slide_spec: Slide, default: str) -> str:
        return slide_spec.layout_variant or default

    def render_slide(
        self,
        slide: "PptxSlide",
        slide_spec: Slide,
        meta: PresentationMeta,
        index: int,
        total_slides: int,
    ) -> None:
        LAYOUT_RENDERERS[slide_spec.type](self, slide, slide_spec, meta, index, total_slides)

    def apply_background(self, slide: "PptxSlide") -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(self.theme.colors.background)

    def textbox(
        self,
        slide: "PptxSlide",
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        margin: float = 0.0,
        vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    ):
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = vertical_anchor
        text_frame.margin_left = Pt(margin)
        text_frame.margin_right = Pt(margin)
        text_frame.margin_top = Pt(margin)
        text_frame.margin_bottom = Pt(margin)
        return shape

    def add_eyebrow(
        self,
        slide: "PptxSlide",
        text: str,
        *,
        left: float,
        top: float,
        width: float = 4.6,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        uppercase: bool = True,
    ):
        shape = self.textbox(slide, left, top, width, 0.25)
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text.upper() if uppercase else text
        self.set_run_style(run, size=self.theme.typography.eyebrow_size, color=self.theme.colors.accent, bold=True)
        return shape

    def add_heading(
        self,
        slide: "PptxSlide",
        *,
        title: str,
        left: float,
        top: float,
        width: float,
        subtitle: str | None = None,
        eyebrow: str | None = None,
        title_size: int | None = None,
        subtitle_width: float | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        slide_type: SlideType | str | None = None,
        layout_variant: str | None = None,
        max_title_lines: int | None = None,
        max_subtitle_lines: int | None = 2,
        min_title_size: int | None = None,
        min_subtitle_size: int | None = None,
        title_subtitle_gap: float = 0.10,
    ) -> float:
        def _truncate_to_chars(text: str, limit: int) -> str:
            if len(text) <= limit:
                return text
            clipped = text[: max(1, limit - 1)].rstrip()
            last_space = clipped.rfind(" ")
            if last_space >= max(12, int(limit * 0.55)):
                clipped = clipped[:last_space].rstrip()
            return clipped.rstrip("-–—,:; ") + "…"

        requested_title_size = title_size or self.theme.typography.title_size
        resolved_title_size = requested_title_size
        resolved_min_title_size = min_title_size or max(self.theme.typography.body_size + 4, requested_title_size - 8)
        display_title = title

        def _estimate_line_count(text: str | None, box_width: float, font_size: int, *, min_chars_per_inch: float, base_chars_per_inch: float, slope: float, min_chars_per_line: int) -> int:
            if not text:
                return 1
            chars_per_inch = max(min_chars_per_inch, base_chars_per_inch - (font_size * slope))
            chars_per_line = max(min_chars_per_line, int(box_width * chars_per_inch))
            return max(1, math.ceil(len(text) / max(1, chars_per_line)))

        def _conservative_char_budget(box_width: float, *, lines: int, chars_per_inch: float) -> int:
            return max(18, int(box_width * chars_per_inch * max(1, lines)))

        allowed_title_lines = max_title_lines or 2

        estimated_lines = 1
        while True:
            estimated_lines = _estimate_line_count(
                display_title,
                width,
                resolved_title_size,
                min_chars_per_inch=3.8,
                base_chars_per_inch=10.6,
                slope=0.16,
                min_chars_per_line=14,
            )
            if estimated_lines <= allowed_title_lines or resolved_title_size <= resolved_min_title_size:
                break
            resolved_title_size -= 1

        if max_title_lines == 1 and display_title:
            display_title = _truncate_to_chars(
                display_title,
                _conservative_char_budget(
                    width,
                    lines=1,
                    chars_per_inch=5.35 if resolved_title_size >= 30 else 5.8,
                ),
            )
            estimated_lines = 1

        if display_title and estimated_lines > allowed_title_lines:
            chars_per_inch = max(3.8, 10.6 - (resolved_title_size * 0.16))
            chars_per_line = max(14, int(width * chars_per_inch))
            display_title = _truncate_to_chars(display_title, chars_per_line * allowed_title_lines)
            estimated_lines = _estimate_line_count(
                display_title,
                width,
                resolved_title_size,
                min_chars_per_inch=3.8,
                base_chars_per_inch=10.6,
                slope=0.16,
                min_chars_per_line=14,
            )

        title_box_height = min(1.95, 0.72 + (0.32 * estimated_lines))

        if eyebrow:
            self.add_eyebrow(slide, eyebrow, left=left, top=top - 0.27, width=min(width, 4.8), align=align)

        title_box = self.textbox(slide, left, top, width, title_box_height)
        self.write_paragraph(
            title_box.text_frame,
            display_title,
            size=resolved_title_size,
            color=self.theme.colors.navy,
            bold=True,
            align=align,
        )
        self.fit_text_frame(
            title_box.text_frame,
            max_size=resolved_title_size,
            min_size=max(14, resolved_min_title_size),
            bold=True,
        )

        heading_bottom = top + title_box_height

        if subtitle:
            resolved_subtitle_size = self.theme.typography.subtitle_size
            resolved_min_subtitle_size = min_subtitle_size or max(self.theme.typography.small_size + 1, resolved_subtitle_size - 2)
            subtitle_box_width = subtitle_width or width
            allowed_subtitle_lines = max_subtitle_lines or 2
            subtitle_lines = 1
            display_subtitle = subtitle
            while True:
                subtitle_lines = _estimate_line_count(
                    display_subtitle,
                    subtitle_box_width,
                    resolved_subtitle_size,
                    min_chars_per_inch=5.0,
                    base_chars_per_inch=13.0,
                    slope=0.20,
                    min_chars_per_line=18,
                )
                if subtitle_lines <= allowed_subtitle_lines or resolved_subtitle_size <= resolved_min_subtitle_size:
                    break
                resolved_subtitle_size -= 1

            display_subtitle = _truncate_to_chars(
                display_subtitle,
                _conservative_char_budget(
                    subtitle_box_width,
                    lines=allowed_subtitle_lines,
                    chars_per_inch=8.9,
                ),
            )
            subtitle_lines = min(
                allowed_subtitle_lines,
                _estimate_line_count(
                    display_subtitle,
                    subtitle_box_width,
                    resolved_subtitle_size,
                    min_chars_per_inch=5.0,
                    base_chars_per_inch=13.0,
                    slope=0.20,
                    min_chars_per_line=18,
                ),
            )

            if display_subtitle and subtitle_lines > allowed_subtitle_lines:
                subtitle_chars_per_inch = max(5.0, 13.0 - (resolved_subtitle_size * 0.20))
                subtitle_chars_per_line = max(18, int(subtitle_box_width * subtitle_chars_per_inch))
                display_subtitle = _truncate_to_chars(display_subtitle, subtitle_chars_per_line * allowed_subtitle_lines)
                subtitle_lines = _estimate_line_count(
                    display_subtitle,
                    subtitle_box_width,
                    resolved_subtitle_size,
                    min_chars_per_inch=5.0,
                    base_chars_per_inch=13.0,
                    slope=0.20,
                    min_chars_per_line=18,
                )

            subtitle_top = heading_bottom + title_subtitle_gap + (0.02 if estimated_lines > 1 else 0.0)
            subtitle_height = min(0.88, 0.34 + (0.18 * subtitle_lines))
            subtitle_box = self.textbox(slide, left, subtitle_top, subtitle_width or width, subtitle_height)
            self.write_paragraph(
                subtitle_box.text_frame,
                display_subtitle,
                size=resolved_subtitle_size,
                color=self.theme.colors.muted,
                align=align,
            )
            self.fit_text_frame(
                subtitle_box.text_frame,
                max_size=resolved_subtitle_size,
                min_size=resolved_min_subtitle_size,
            )
            heading_bottom = subtitle_top + subtitle_height

        return heading_bottom

    def fit_text_frame(
        self,
        text_frame: "TextFrame",
        *,
        max_size: int,
        min_size: int = 10,
        bold: bool = False,
        italic: bool = False,
    ) -> None:
        text_frame.word_wrap = True
        try:
            text_frame.fit_text(
                font_family=self.theme.typography.font_name,
                max_size=max_size,
                bold=bold,
                italic=italic,
            )
        except Exception:
            pass

        self._apply_manual_text_frame_fit(
            text_frame,
            max_size=max_size,
            min_size=min_size,
            bold=bold,
            italic=italic,
        )

    def _apply_manual_text_frame_fit(
        self,
        text_frame: "TextFrame",
        *,
        max_size: int,
        min_size: int,
        bold: bool = False,
        italic: bool = False,
    ) -> None:
        stats = self._estimate_text_frame_pressure(text_frame)
        recommended_size = self._recommended_font_size(
            max_size=max_size,
            min_size=min_size,
            char_count=stats["char_count"],
            paragraph_count=stats["paragraph_count"],
            area=stats["area"],
            longest_line=stats["longest_line"],
        )

        for paragraph in text_frame.paragraphs:
            if paragraph.level > 0 and recommended_size <= max_size:
                paragraph.space_after = Pt(min(4, max(1, recommended_size // 4)))
            elif paragraph.level == 0 and stats["paragraph_count"] > 1:
                paragraph.space_after = Pt(min(6, max(2, recommended_size // 3)))

            if stats["paragraph_count"] > 2:
                paragraph.line_spacing = 1.0

            for run in paragraph.runs:
                font = run.font
                if font.name is None:
                    font.name = self.theme.typography.font_name
                if font.size is None or font.size.pt > recommended_size:
                    font.size = Pt(recommended_size)
                if bold:
                    font.bold = True
                if italic:
                    font.italic = True

    def _estimate_text_frame_pressure(self, text_frame: "TextFrame") -> dict[str, float | int]:
        shape = getattr(text_frame, "_parent", None)
        width_emu = getattr(shape, "width", 0) or 0
        height_emu = getattr(shape, "height", 0) or 0
        width_inches = float(width_emu) / 914400 if width_emu else 4.0
        height_inches = float(height_emu) / 914400 if height_emu else 1.0
        area = max(0.6, width_inches * height_inches)

        paragraph_count = 0
        char_count = 0
        longest_line = 0
        for paragraph in text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs).strip()
            if not line:
                continue
            paragraph_count += 1
            char_count += len(line)
            longest_line = max(longest_line, len(line))

        if char_count == 0 and getattr(text_frame, "text", ""):
            text = text_frame.text.strip()
            char_count = len(text)
            paragraph_count = max(1, len([segment for segment in text.splitlines() if segment.strip()]))
            longest_line = max((len(segment.strip()) for segment in text.splitlines() if segment.strip()), default=0)

        return {
            "width_inches": width_inches,
            "height_inches": height_inches,
            "area": area,
            "paragraph_count": paragraph_count,
            "char_count": char_count,
            "longest_line": longest_line,
        }

    def _recommended_font_size(
        self,
        *,
        max_size: int,
        min_size: int,
        char_count: int,
        paragraph_count: int,
        area: float,
        longest_line: int,
    ) -> int:
        density = char_count / area if area else float(char_count)
        reduction = 0
        if density > 30:
            reduction += 1
        if density > 42:
            reduction += 1
        if density > 56:
            reduction += 1
        if density > 74:
            reduction += 2
        if paragraph_count >= 3:
            reduction += 1
        if paragraph_count >= 5:
            reduction += 1
        if longest_line >= 70:
            reduction += 1
        if longest_line >= 110:
            reduction += 1
        return max(min_size, max_size - reduction)

    def panel_content_box(
        self,
        slide: "PptxSlide",
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        padding: float | None = None,
    ):
        pad = padding if padding is not None else self.theme.components.panel_padding
        return self.textbox(slide, left + pad, top + pad, width - (pad * 2), height - (pad * 2))

    def panel_inner_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        padding: float | None = None,
    ) -> tuple[float, float, float, float]:
        pad = padding if padding is not None else self.theme.components.panel_padding
        return left + pad, top + pad, width - (pad * 2), height - (pad * 2)

    def compute_cover_crop(
        self,
        *,
        image_width_px: int,
        image_height_px: int,
        box_width: float,
        box_height: float,
        focal_x: float | None = None,
        focal_y: float | None = None,
    ) -> tuple[float, float, float, float]:
        if image_width_px <= 0 or image_height_px <= 0 or box_width <= 0 or box_height <= 0:
            return 0.0, 0.0, 0.0, 0.0

        image_aspect = image_width_px / image_height_px
        box_aspect = box_width / box_height
        if abs(image_aspect - box_aspect) <= 1e-6:
            return 0.0, 0.0, 0.0, 0.0

        if image_aspect > box_aspect:
            keep_fraction = box_aspect / image_aspect
            center = 0.5 if focal_x is None else focal_x
            crop_left = min(max(0.0, center - (keep_fraction / 2.0)), 1.0 - keep_fraction)
            crop_right = max(0.0, 1.0 - keep_fraction - crop_left)
            return crop_left, 0.0, crop_right, 0.0

        keep_fraction = image_aspect / box_aspect
        center = 0.5 if focal_y is None else focal_y
        crop_top = min(max(0.0, center - (keep_fraction / 2.0)), 1.0 - keep_fraction)
        crop_bottom = max(0.0, 1.0 - keep_fraction - crop_top)
        return 0.0, crop_top, 0.0, crop_bottom

    def add_image_cover(
        self,
        slide: "PptxSlide",
        image_path: str | Path,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        focal_x: float | None = None,
        focal_y: float | None = None,
    ):
        resolved_path = Path(image_path)
        picture = slide.shapes.add_picture(
            str(resolved_path),
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height),
        )
        with PILImage.open(resolved_path) as image:
            crop_left, crop_top, crop_right, crop_bottom = self.compute_cover_crop(
                image_width_px=image.width,
                image_height_px=image.height,
                box_width=width,
                box_height=height,
                focal_x=focal_x,
                focal_y=focal_y,
            )
        picture.crop_left = crop_left
        picture.crop_top = crop_top
        picture.crop_right = crop_right
        picture.crop_bottom = crop_bottom
        return picture

    def stack_vertical_regions(
        self,
        *,
        top: float,
        height: float,
        regions: list[dict[str, float | str]],
        gap: float,
    ) -> list[tuple[dict[str, float | str], tuple[float, float]]]:
        if not regions:
            return []

        base_heights: list[float] = []
        flex_total = 0.0
        for region in regions:
            base_height = float(region.get("height") or region.get("min_height") or 0.0)
            base_heights.append(base_height)
            flex_total += float(region.get("flex") or 0.0)

        gap_total = gap * max(0, len(regions) - 1)
        remaining = max(0.0, height - sum(base_heights) - gap_total)

        bounds: list[tuple[dict[str, float | str], tuple[float, float]]] = []
        cursor_top = top
        for region, base_height in zip(regions, base_heights, strict=True):
            extra = 0.0
            if flex_total > 0 and float(region.get("flex") or 0.0) > 0:
                extra = remaining * (float(region.get("flex") or 0.0) / flex_total)
            region_height = base_height + extra
            bounds.append((region, (cursor_top, region_height)))
            cursor_top += region_height + gap
        return bounds

    def stack_horizontal_regions(
        self,
        *,
        left: float,
        width: float,
        regions: list[dict[str, float | str]],
        gap: float,
    ) -> list[tuple[dict[str, float | str], tuple[float, float]]]:
        if not regions:
            return []

        base_widths: list[float] = []
        flex_total = 0.0
        for region in regions:
            base_width = float(region.get("width") or region.get("min_width") or 0.0)
            base_widths.append(base_width)
            flex_total += float(region.get("flex") or 0.0)

        gap_total = gap * max(0, len(regions) - 1)
        remaining = max(0.0, width - sum(base_widths) - gap_total)

        bounds: list[tuple[dict[str, float | str], tuple[float, float]]] = []
        cursor_left = left
        for region, base_width in zip(regions, base_widths, strict=True):
            extra = 0.0
            if flex_total > 0 and float(region.get("flex") or 0.0) > 0:
                extra = remaining * (float(region.get("flex") or 0.0) / flex_total)
            region_width = base_width + extra
            bounds.append((region, (cursor_left, region_width)))
            cursor_left += region_width + gap
        return bounds

    def _allocate_constrained_region_sizes(
        self,
        *,
        total_size: float,
        gap: float,
        regions: list[dict[str, float | str]],
        fixed_key: str,
        min_key: str,
        max_key: str,
    ) -> list[float]:
        if not regions:
            return []

        gap_total = gap * max(0, len(regions) - 1)
        available = max(0.0, total_size - gap_total)
        sizes = [0.0 for _ in regions]
        remaining = available
        flexible_indices: list[int] = []

        for index, region in enumerate(regions):
            fixed_size = region.get(fixed_key)
            min_size = float(region.get(min_key) or 0.0)
            max_size = float(region.get(max_key)) if region.get(max_key) is not None else None
            if fixed_size is not None:
                resolved = max(min_size, float(fixed_size))
                if max_size is not None:
                    resolved = min(resolved, max_size)
                sizes[index] = resolved
                remaining -= resolved
            else:
                sizes[index] = min_size
                remaining -= min_size
                flexible_indices.append(index)

        remaining = max(0.0, remaining)
        active = flexible_indices.copy()
        while active and remaining > 1e-6:
            target_units = []
            for index in active:
                target_share = regions[index].get("target_share")
                if target_share is not None:
                    target_units.append(max(0.0, float(target_share)))
                else:
                    target_units.append(max(0.0, float(regions[index].get("content_weight") or regions[index].get("flex") or 1.0)))

            target_total = sum(target_units) or float(len(active))
            clamped = False
            next_active: list[int] = []
            consumed = 0.0

            for index, unit in zip(active, target_units, strict=True):
                share = (unit / target_total) if target_total else (1.0 / len(active))
                proposed_extra = remaining * share
                max_size = float(regions[index].get(max_key)) if regions[index].get(max_key) is not None else None
                if max_size is not None and sizes[index] + proposed_extra > max_size:
                    extra = max(0.0, max_size - sizes[index])
                    sizes[index] = max_size
                    consumed += extra
                    clamped = True
                else:
                    next_active.append(index)

            if not clamped:
                for index, unit in zip(active, target_units, strict=True):
                    share = (unit / target_total) if target_total else (1.0 / len(active))
                    sizes[index] += remaining * share
                remaining = 0.0
                break

            remaining = max(0.0, remaining - consumed)
            active = next_active

        return sizes

    def build_constrained_columns(
        self,
        *,
        left: float,
        width: float,
        gap: float,
        regions: list[dict[str, float | str]],
    ) -> list[tuple[float, float]]:
        column_widths = self._allocate_constrained_region_sizes(
            total_size=width,
            gap=gap,
            regions=regions,
            fixed_key="width",
            min_key="min_width",
            max_key="max_width",
        )
        bounds: list[tuple[float, float]] = []
        cursor_left = left
        for column_width in column_widths:
            bounds.append((cursor_left, column_width))
            cursor_left += column_width + gap
        return bounds

    def build_constrained_rows(
        self,
        *,
        top: float,
        height: float,
        gap: float,
        regions: list[dict[str, float | str]],
    ) -> list[tuple[float, float]]:
        row_heights = self._allocate_constrained_region_sizes(
            total_size=height,
            gap=gap,
            regions=regions,
            fixed_key="height",
            min_key="min_height",
            max_key="max_height",
        )
        bounds: list[tuple[float, float]] = []
        cursor_top = top
        for row_height in row_heights:
            bounds.append((cursor_top, row_height))
            cursor_top += row_height + gap
        return bounds

    def build_columns(
        self,
        *,
        left: float,
        width: float,
        gap: float,
        count: int | None = None,
        min_width: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
    ) -> list[tuple[float, float]]:
        if regions is None:
            if count is None:
                raise ValueError("build_columns requires count or regions")
            regions = [
                {"kind": f"column_{index + 1}", "min_width": min_width, "flex": 1.0}
                for index in range(count)
            ]
        if any(region.get("target_share") is not None or region.get("max_width") is not None for region in regions):
            return self.build_constrained_columns(left=left, width=width, gap=gap, regions=regions)
        return [bounds for _, bounds in self.stack_horizontal_regions(left=left, width=width, regions=regions, gap=gap)]

    def build_named_columns(
        self,
        *,
        left: float,
        width: float,
        gap: float,
        count: int | None = None,
        min_width: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
        kind_prefix: str = "column",
    ) -> dict[str, tuple[float, float]]:
        resolved_regions = regions
        if resolved_regions is None:
            if count is None:
                raise ValueError("build_named_columns requires count or regions")
            resolved_regions = [
                {"kind": f"{kind_prefix}_{index + 1}", "min_width": min_width, "flex": 1.0}
                for index in range(count)
            ]
        bounds = self.build_columns(
            left=left,
            width=width,
            gap=gap,
            count=count,
            min_width=min_width,
            regions=resolved_regions,
        )
        return self._build_named_region_map(regions=resolved_regions, bounds=bounds, default_prefix=kind_prefix)

    def build_weighted_columns(
        self,
        *,
        left: float,
        width: float,
        gap: float,
        weights: list[float],
        min_width: float = 0.0,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
        kind_prefix: str = "column",
    ) -> list[tuple[float, float]]:
        flexes = self.normalize_content_flexes(weights, min_flex=min_flex, max_flex=max_flex)
        return self.build_columns(
            left=left,
            width=width,
            gap=gap,
            min_width=min_width,
            regions=[
                {"kind": f"{kind_prefix}_{index + 1}", "min_width": min_width, "flex": flex}
                for index, flex in enumerate(flexes)
            ],
        )

    def build_rows(
        self,
        *,
        top: float,
        height: float,
        gap: float,
        count: int | None = None,
        min_height: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
    ) -> list[tuple[float, float]]:
        if regions is None:
            if count is None:
                raise ValueError("build_rows requires count or regions")
            regions = [
                {"kind": f"row_{index + 1}", "min_height": min_height, "flex": 1.0}
                for index in range(count)
            ]
        if any(region.get("target_share") is not None or region.get("max_height") is not None for region in regions):
            return self.build_constrained_rows(top=top, height=height, gap=gap, regions=regions)
        return [bounds for _, bounds in self.stack_vertical_regions(top=top, height=height, regions=regions, gap=gap)]

    def build_named_rows(
        self,
        *,
        top: float,
        height: float,
        gap: float,
        count: int | None = None,
        min_height: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
        kind_prefix: str = "row",
    ) -> dict[str, tuple[float, float]]:
        resolved_regions = regions
        if resolved_regions is None:
            if count is None:
                raise ValueError("build_named_rows requires count or regions")
            resolved_regions = [
                {"kind": f"{kind_prefix}_{index + 1}", "min_height": min_height, "flex": 1.0}
                for index in range(count)
            ]
        bounds = self.build_rows(
            top=top,
            height=height,
            gap=gap,
            count=count,
            min_height=min_height,
            regions=resolved_regions,
        )
        return self._build_named_region_map(regions=resolved_regions, bounds=bounds, default_prefix=kind_prefix)

    def build_weighted_rows(
        self,
        *,
        top: float,
        height: float,
        gap: float,
        weights: list[float],
        min_height: float = 0.0,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
        kind_prefix: str = "row",
    ) -> list[tuple[float, float]]:
        flexes = self.normalize_content_flexes(weights, min_flex=min_flex, max_flex=max_flex)
        return self.build_rows(
            top=top,
            height=height,
            gap=gap,
            min_height=min_height,
            regions=[
                {"kind": f"{kind_prefix}_{index + 1}", "min_height": min_height, "flex": flex}
                for index, flex in enumerate(flexes)
            ],
        )

    def build_panel_row_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        count: int | None = None,
        min_width: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
    ) -> list[tuple[float, float, float, float]]:
        columns = self.build_columns(
            left=left,
            width=width,
            gap=gap,
            count=count,
            min_width=min_width,
            regions=regions,
        )
        return [(column_left, top, column_width, height) for column_left, column_width in columns]

    def build_weighted_panel_row_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        weights: list[float],
        min_width: float = 0.0,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
        kind_prefix: str = "panel",
    ) -> list[tuple[float, float, float, float]]:
        columns = self.build_weighted_columns(
            left=left,
            width=width,
            gap=gap,
            weights=weights,
            min_width=min_width,
            min_flex=min_flex,
            max_flex=max_flex,
            kind_prefix=kind_prefix,
        )
        return [(column_left, top, column_width, height) for column_left, column_width in columns]

    def build_constrained_panel_row_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        regions: list[dict[str, float | str]],
    ) -> list[tuple[float, float, float, float]]:
        columns = self.build_constrained_columns(left=left, width=width, gap=gap, regions=regions)
        return [(column_left, top, column_width, height) for column_left, column_width in columns]

    def build_panel_row_content_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        count: int | None = None,
        min_width: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
        padding: float | None = None,
    ) -> list[tuple[tuple[float, float, float, float], tuple[float, float, float, float]]]:
        panel_bounds = self.build_panel_row_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            gap=gap,
            count=count,
            min_width=min_width,
            regions=regions,
        )
        return [
            (
                bounds,
                self.panel_inner_bounds(
                    left=bounds[0],
                    top=bounds[1],
                    width=bounds[2],
                    height=bounds[3],
                    padding=padding,
                ),
            )
            for bounds in panel_bounds
        ]

    def build_named_panel_row_content_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        count: int | None = None,
        min_width: float = 0.0,
        regions: list[dict[str, float | str]] | None = None,
        padding: float | None = None,
        kind_prefix: str = "panel",
    ) -> dict[str, dict[str, tuple[float, float, float, float]]]:
        resolved_regions = regions
        if resolved_regions is None:
            if count is None:
                raise ValueError("build_named_panel_row_content_bounds requires count or regions")
            resolved_regions = [
                {"kind": f"{kind_prefix}_{index + 1}", "min_width": min_width, "flex": 1.0}
                for index in range(count)
            ]
        bounds = self.build_panel_row_content_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            gap=gap,
            count=count,
            min_width=min_width,
            regions=resolved_regions,
            padding=padding,
        )
        named: dict[str, dict[str, tuple[float, float, float, float]]] = {}
        for index, (region, (panel_bounds, content_bounds)) in enumerate(zip(resolved_regions, bounds, strict=True), start=1):
            kind = str(region.get("kind") or f"{kind_prefix}_{index}")
            named[kind] = {"panel": panel_bounds, "content": content_bounds}
        return named

    def build_weighted_panel_row_content_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        weights: list[float],
        min_width: float = 0.0,
        padding: float | None = None,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
        kind_prefix: str = "panel",
    ) -> list[tuple[tuple[float, float, float, float], tuple[float, float, float, float]]]:
        panel_bounds = self.build_weighted_panel_row_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            gap=gap,
            weights=weights,
            min_width=min_width,
            min_flex=min_flex,
            max_flex=max_flex,
            kind_prefix=kind_prefix,
        )
        return [
            (
                bounds,
                self.panel_inner_bounds(
                    left=bounds[0],
                    top=bounds[1],
                    width=bounds[2],
                    height=bounds[3],
                    padding=padding,
                ),
            )
            for bounds in panel_bounds
        ]

    def build_constrained_panel_row_content_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        gap: float,
        regions: list[dict[str, float | str]],
        padding: float | None = None,
    ) -> list[tuple[tuple[float, float, float, float], tuple[float, float, float, float]]]:
        panel_bounds = self.build_constrained_panel_row_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            gap=gap,
            regions=regions,
        )
        return [
            (
                bounds,
                self.panel_inner_bounds(
                    left=bounds[0],
                    top=bounds[1],
                    width=bounds[2],
                    height=bounds[3],
                    padding=padding,
                ),
            )
            for bounds in panel_bounds
        ]

    def build_weighted_panel_grid(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_gap: float,
        row_gap: float,
        column_weights: list[float],
        row_weights: list[float],
        column_min_width: float = 0.0,
        row_min_height: float = 0.0,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
        column_kind_prefix: str = "column",
        row_kind_prefix: str = "row",
    ) -> list[list[tuple[float, float, float, float]]]:
        return self.build_grid_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            column_regions=[
                {
                    "kind": f"{column_kind_prefix}_{index + 1}",
                    "min_width": column_min_width,
                    "flex": flex,
                }
                for index, flex in enumerate(
                    self.normalize_content_flexes(column_weights, min_flex=min_flex, max_flex=max_flex)
                )
            ],
            row_regions=[
                {
                    "kind": f"{row_kind_prefix}_{index + 1}",
                    "min_height": row_min_height,
                    "flex": flex,
                }
                for index, flex in enumerate(
                    self.normalize_content_flexes(row_weights, min_flex=min_flex, max_flex=max_flex)
                )
            ],
            column_gap=column_gap,
            row_gap=row_gap,
        )

    def build_panel_content_stack_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        regions: list[dict[str, float | str]],
        gap: float,
        padding: float | None = None,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
    ) -> list[tuple[dict[str, float | str], tuple[float, float, float, float]]]:
        content_left, content_top, content_width, content_height = self.panel_inner_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            padding=padding,
        )
        return [
            (region, (content_left, region_top, content_width, region_height))
            for region, (region_top, region_height) in self.build_content_stack(
                top=content_top,
                height=content_height,
                regions=regions,
                gap=gap,
                min_flex=min_flex,
                max_flex=max_flex,
            )
        ]

    def build_named_panel_content_stack_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        regions: list[dict[str, float | str]],
        gap: float,
        padding: float | None = None,
        constrained: bool = False,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
        kind_prefix: str = "region",
    ) -> dict[str, tuple[float, float, float, float]]:
        bounds = (
            self.build_constrained_panel_content_stack_bounds(
                left=left,
                top=top,
                width=width,
                height=height,
                regions=regions,
                gap=gap,
                padding=padding,
            )
            if constrained
            else self.build_panel_content_stack_bounds(
                left=left,
                top=top,
                width=width,
                height=height,
                regions=regions,
                gap=gap,
                padding=padding,
                min_flex=min_flex,
                max_flex=max_flex,
            )
        )
        return self._build_named_region_map(
            regions=[region for region, _ in bounds],
            bounds=[region_bounds for _, region_bounds in bounds],
            default_prefix=kind_prefix,
        )

    def add_structured_panel(
        self,
        slide: "PptxSlide",
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        regions: list[dict[str, float | str]] | None = None,
        gap: float = 0.08,
        padding: float | None = None,
        fill_color: str | None = None,
        line_color: str | None = None,
        accent_color: str | None = None,
        accent_height: float | None = None,
        constrained: bool = False,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
    ) -> dict[str, object]:
        panel = self.add_panel(
            slide,
            left,
            top,
            width,
            height,
            fill_color=fill_color or self.theme.colors.surface,
            line_color=line_color or self.theme.colors.line,
        )
        if accent_color:
            self.add_accent_bar(
                slide,
                left,
                top,
                width,
                accent_height or self.theme.components.accent_bar_height,
                color=accent_color,
            )
        inner_bounds = self.panel_inner_bounds(left=left, top=top, width=width, height=height, padding=padding)
        content_regions = {}
        if regions:
            content_regions = self.build_named_panel_content_stack_bounds(
                left=left,
                top=top,
                width=width,
                height=height,
                regions=regions,
                gap=gap,
                padding=padding,
                constrained=constrained,
                min_flex=min_flex,
                max_flex=max_flex,
            )
        return {
            "panel": panel,
            "panel_bounds": (left, top, width, height),
            "inner_bounds": inner_bounds,
            "content_regions": content_regions,
        }

    def add_accent_panel(
        self,
        slide: "PptxSlide",
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        accent_color: str,
        padding: float | None = None,
        fill_color: str | None = None,
        line_color: str | None = None,
        accent_height: float | None = None,
    ) -> tuple[object, tuple[float, float, float, float]]:
        panel = self.add_panel(
            slide,
            left,
            top,
            width,
            height,
            fill_color=fill_color or self.theme.colors.surface,
            line_color=line_color or self.theme.colors.line,
        )
        self.add_accent_bar(
            slide,
            left,
            top,
            width,
            accent_height or self.theme.components.accent_bar_height,
            color=accent_color,
        )
        inner_bounds = self.panel_inner_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            padding=padding,
        )
        return panel, inner_bounds

    def add_visual_slot(
        self,
        slide: "PptxSlide",
        slide_spec: Slide,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        accent_color: str | None = None,
        padding: float | None = None,
        caption_text: str | None = None,
    ) -> dict[str, object]:
        asset = self.resolve_asset(slide_spec.image_path)
        focal_x, focal_y = self.resolve_image_focal_point(slide_spec)
        caption = caption_text if caption_text is not None else slide_spec.image_caption

        if asset:
            picture = self.add_image_cover(
                slide,
                asset,
                left=left,
                top=top,
                width=width,
                height=height,
                focal_x=focal_x,
                focal_y=focal_y,
            )
            if caption:
                caption_box = self.textbox(slide, left, top + height + 0.08, width, 0.24)
                self.write_paragraph(
                    caption_box.text_frame,
                    caption,
                    size=self.theme.typography.small_size,
                    color=self.theme.colors.muted,
                )
                self.fit_text_frame(caption_box.text_frame, max_size=self.theme.typography.small_size)
            return {
                "asset": str(asset),
                "used_asset": True,
                "shape": picture,
                "bounds": (left, top, width, height),
            }

        self.add_panel(
            slide,
            left,
            top,
            width,
            height,
            fill_color=self.theme.colors.soft_fill,
            line_color=self.theme.colors.line,
        )
        self.add_accent_bar(
            slide,
            left,
            top,
            width,
            self.theme.components.accent_bar_height,
            color=accent_color or self.theme.colors.accent,
        )
        content_left, content_top, content_width, content_height = self.panel_inner_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            padding=padding,
        )
        placeholder = self.describe_visual_placeholder(slide_spec)
        label_box = self.textbox(slide, content_left, content_top, min(content_width, 2.4), 0.24)
        self.write_paragraph(
            label_box.text_frame,
            placeholder["label"],
            size=self.theme.typography.small_size,
            color=self.theme.colors.accent,
            bold=True,
        )
        self.fit_text_frame(label_box.text_frame, max_size=self.theme.typography.small_size, bold=True)

        headline_box = self.textbox(slide, content_left, content_top + 0.30, content_width, 0.3)
        self.write_paragraph(
            headline_box.text_frame,
            placeholder["headline"],
            size=self.theme.typography.small_size + 2,
            color=self.theme.colors.navy,
            bold=True,
        )
        self.fit_text_frame(
            headline_box.text_frame,
            max_size=self.theme.typography.small_size + 2,
            min_size=self.theme.typography.small_size,
            bold=True,
        )

        missing_asset_note = f"Missing asset: {slide_spec.image_path}" if slide_spec.image_path else None
        guidance_text = caption or placeholder["guidance"]
        detail_lines = [line for line in [missing_asset_note, guidance_text] if line]
        guidance_box = self.textbox(slide, content_left, content_top + 0.72, content_width, max(0.3, content_height - 0.72))
        self.write_paragraph(
            guidance_box.text_frame,
            "\n".join(detail_lines),
            size=self.theme.typography.small_size + 1,
            color=self.theme.colors.text,
        )
        self.fit_text_frame(
            guidance_box.text_frame,
            max_size=self.theme.typography.small_size + 1,
            min_size=self.theme.typography.small_size,
        )
        return {
            "asset": None,
            "used_asset": False,
            "shape": None,
            "bounds": (left, top, width, height),
            "content_bounds": (content_left, content_top, content_width, content_height),
        }

    def build_panel_grid(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_gap: float,
        row_gap: float,
        column_count: int | None = None,
        row_count: int | None = None,
        column_min_width: float = 0.0,
        row_min_height: float = 0.0,
        column_regions: list[dict[str, float | str]] | None = None,
        row_regions: list[dict[str, float | str]] | None = None,
    ) -> list[list[tuple[float, float, float, float]]]:
        resolved_column_regions = column_regions
        if resolved_column_regions is None:
            if column_count is None:
                raise ValueError("build_panel_grid requires column_count or column_regions")
            resolved_column_regions = [
                {"kind": f"column_{index + 1}", "min_width": column_min_width, "flex": 1.0}
                for index in range(column_count)
            ]

        resolved_row_regions = row_regions
        if resolved_row_regions is None:
            if row_count is None:
                raise ValueError("build_panel_grid requires row_count or row_regions")
            resolved_row_regions = [
                {"kind": f"row_{index + 1}", "min_height": row_min_height, "flex": 1.0}
                for index in range(row_count)
            ]

        return self.build_grid_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            column_regions=resolved_column_regions,
            row_regions=resolved_row_regions,
            column_gap=column_gap,
            row_gap=row_gap,
        )

    def build_panel_grid_content_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_gap: float,
        row_gap: float,
        column_count: int | None = None,
        row_count: int | None = None,
        column_min_width: float = 0.0,
        row_min_height: float = 0.0,
        column_regions: list[dict[str, float | str]] | None = None,
        row_regions: list[dict[str, float | str]] | None = None,
        padding: float | None = None,
    ) -> list[list[tuple[tuple[float, float, float, float], tuple[float, float, float, float]]]]:
        grid = self.build_panel_grid(
            left=left,
            top=top,
            width=width,
            height=height,
            column_gap=column_gap,
            row_gap=row_gap,
            column_count=column_count,
            row_count=row_count,
            column_min_width=column_min_width,
            row_min_height=row_min_height,
            column_regions=column_regions,
            row_regions=row_regions,
        )
        return [
            [
                (
                    bounds,
                    self.panel_inner_bounds(
                        left=bounds[0],
                        top=bounds[1],
                        width=bounds[2],
                        height=bounds[3],
                        padding=padding,
                    ),
                )
                for bounds in row
            ]
            for row in grid
        ]

    def estimate_content_weight(
        self,
        *,
        title: str | None = None,
        body: str | None = None,
        bullets: list[str] | None = None,
        footer: str | None = None,
        tag: str | None = None,
    ) -> float:
        bullet_items = bullets or []
        weight = 1.0
        if title:
            weight += min(0.8, len(title.split()) / 7)
        if body:
            weight += min(1.4, len(body.split()) / 18)
        if footer:
            weight += min(0.35, len(footer.split()) / 10)
        if tag:
            weight += min(0.25, len(tag.split()) / 6)
        for bullet in bullet_items:
            weight += min(0.3, len(bullet.split()) / 14)
        return weight

    def normalize_content_flexes(
        self,
        weights: list[float],
        *,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
    ) -> list[float]:
        if not weights:
            return []
        min_weight = min(weights)
        max_weight = max(weights)
        if max_weight == min_weight:
            return [1.0 for _ in weights]
        span = max_flex - min_flex
        return [
            min_flex + ((weight - min_weight) / (max_weight - min_weight)) * span
            for weight in weights
        ]

    def rebalance_flexible_regions(
        self,
        regions: list[dict[str, float | str]],
        *,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
    ) -> list[dict[str, float | str]]:
        balanced_regions = [dict(region) for region in regions]
        flexible_indices = [
            index for index, region in enumerate(balanced_regions) if float(region.get("flex") or 0.0) > 0.0
        ]
        if len(flexible_indices) <= 1:
            return balanced_regions

        weights = [
            float(balanced_regions[index].get("content_weight") or balanced_regions[index].get("flex") or 1.0)
            for index in flexible_indices
        ]
        flexes = self.normalize_content_flexes(weights, min_flex=min_flex, max_flex=max_flex)
        if weights and max(weights) / max(0.1, min(weights)) >= 1.85:
            flexes = self.normalize_content_flexes(
                weights,
                min_flex=max(0.82, min_flex - 0.04),
                max_flex=min(1.65, max_flex + 0.16),
            )
        for index, flex in zip(flexible_indices, flexes, strict=True):
            balanced_regions[index]["flex"] = flex
            balanced_regions[index].pop("content_weight", None)
        return balanced_regions

    def infer_visual_placeholder_kind(
        self,
        *,
        slide_type: str,
        title: str | None = None,
        body: str | None = None,
        caption: str | None = None,
        image_path: str | None = None,
    ) -> str:
        return infer_visual_placeholder_kind(
            slide_type=slide_type,
            title=title,
            body=body,
            caption=caption,
            image_path=image_path,
        )

    def resolve_image_focal_point(self, slide_spec: Slide) -> tuple[float, float]:
        return infer_contextual_image_focal_point(slide_spec)

    def describe_visual_placeholder(
        self,
        slide_spec: Slide,
    ) -> dict[str, str]:
        kind = self.infer_visual_placeholder_kind(
            slide_type=slide_spec.type.value,
            title=slide_spec.title,
            body=slide_spec.body,
            caption=slide_spec.image_caption,
            image_path=slide_spec.image_path,
        )
        mapping = {
            "photo": {
                "label": "VISUAL",
                "headline": "Editorial visual pending" if slide_spec.image_path else "Editorial visual suggested",
                "guidance": "Use a premium contextual photo with negative space to reinforce the narrative once the final asset is available.",
            },
            "screenshot": {
                "label": "PRODUCT VISUAL",
                "headline": "Product screenshot pending" if slide_spec.image_path else "Product visual suggested",
                "guidance": "Use a crisp product screenshot, dashboard capture, or UI mockup that supports the point being made on this slide.",
            },
            "diagram": {
                "label": "WORKFLOW VISUAL",
                "headline": "Workflow visual pending" if slide_spec.image_path else "Workflow visual suggested",
                "guidance": "Use a structured workflow, architecture, or process visual that makes the sequence and relationships easier to scan.",
            },
            "analytical_visual": {
                "label": "ANALYTICAL VISUAL",
                "headline": "Analytical visual pending" if slide_spec.image_path else "Analytical visual suggested",
                "guidance": "Use a chart, scorecard, dashboard, or other analytical visual that makes the evidence easier to read at a glance.",
            },
        }
        return {"kind": kind, **mapping[kind]}

    def build_content_stack(
        self,
        *,
        top: float,
        height: float,
        regions: list[dict[str, float | str]],
        gap: float,
        min_flex: float = 0.9,
        max_flex: float = 1.35,
    ) -> list[tuple[dict[str, float | str], tuple[float, float]]]:
        return self.stack_vertical_regions(
            top=top,
            height=height,
            regions=self.rebalance_flexible_regions(regions, min_flex=min_flex, max_flex=max_flex),
            gap=gap,
        )

    def build_constrained_content_stack(
        self,
        *,
        top: float,
        height: float,
        regions: list[dict[str, float | str]],
        gap: float,
    ) -> list[tuple[dict[str, float | str], tuple[float, float]]]:
        rows = self.build_constrained_rows(top=top, height=height, gap=gap, regions=regions)
        return list(zip(regions, rows, strict=True))

    def build_constrained_panel_content_stack_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        regions: list[dict[str, float | str]],
        gap: float,
        padding: float | None = None,
    ) -> list[tuple[dict[str, float | str], tuple[float, float, float, float]]]:
        content_left, content_top, content_width, content_height = self.panel_inner_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            padding=padding,
        )
        return [
            (region, (content_left, region_top, content_width, region_height))
            for region, (region_top, region_height) in self.build_constrained_content_stack(
                top=content_top,
                height=content_height,
                regions=regions,
                gap=gap,
            )
        ]

    def build_constrained_grid_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_regions: list[dict[str, float | str]],
        row_regions: list[dict[str, float | str]],
        column_gap: float,
        row_gap: float,
    ) -> list[list[tuple[float, float, float, float]]]:
        columns = self.build_constrained_columns(
            left=left,
            width=width,
            gap=column_gap,
            regions=column_regions,
        )
        rows = self.build_constrained_rows(
            top=top,
            height=height,
            gap=row_gap,
            regions=row_regions,
        )

        return [
            [
                (column_left, row_top, column_width, row_height)
                for column_left, column_width in columns
            ]
            for row_top, row_height in rows
        ]

    def build_constrained_panel_grid(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_gap: float,
        row_gap: float,
        column_regions: list[dict[str, float | str]],
        row_regions: list[dict[str, float | str]],
    ) -> list[list[tuple[float, float, float, float]]]:
        return self.build_constrained_grid_bounds(
            left=left,
            top=top,
            width=width,
            height=height,
            column_regions=column_regions,
            row_regions=row_regions,
            column_gap=column_gap,
            row_gap=row_gap,
        )

    def build_constrained_panel_grid_content_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_gap: float,
        row_gap: float,
        column_regions: list[dict[str, float | str]],
        row_regions: list[dict[str, float | str]],
        padding: float | None = None,
    ) -> list[list[tuple[tuple[float, float, float, float], tuple[float, float, float, float]]]]:
        grid = self.build_constrained_panel_grid(
            left=left,
            top=top,
            width=width,
            height=height,
            column_gap=column_gap,
            row_gap=row_gap,
            column_regions=column_regions,
            row_regions=row_regions,
        )
        return [
            [
                (
                    bounds,
                    self.panel_inner_bounds(
                        left=bounds[0],
                        top=bounds[1],
                        width=bounds[2],
                        height=bounds[3],
                        padding=padding,
                    ),
                )
                for bounds in row
            ]
            for row in grid
        ]

    def build_grid_bounds(
        self,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        column_regions: list[dict[str, float | str]],
        row_regions: list[dict[str, float | str]],
        column_gap: float,
        row_gap: float,
    ) -> list[list[tuple[float, float, float, float]]]:
        columns = self.stack_horizontal_regions(
            left=left,
            width=width,
            regions=column_regions,
            gap=column_gap,
        )
        rows = self.stack_vertical_regions(
            top=top,
            height=height,
            regions=row_regions,
            gap=row_gap,
        )

        return [
            [
                (column_left, row_top, column_width, row_height)
                for _, (column_left, column_width) in columns
            ]
            for _, (row_top, row_height) in rows
        ]

    def add_quote_block(
        self,
        slide: "PptxSlide",
        *,
        quote: str,
        left: float,
        top: float,
        width: float,
        height: float,
        attribution: str | None = None,
    ) -> None:
        quote_box = self.textbox(slide, left, top, width, height)
        paragraph = quote_box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = quote
        self.set_run_style(run, size=self.theme.typography.quote_size, color=self.theme.colors.navy, bold=True, italic=True)
        self.fit_text_frame(
            quote_box.text_frame,
            max_size=self.theme.typography.quote_size,
            bold=True,
            italic=True,
        )

        if attribution:
            attribution_box = self.textbox(slide, left, top + height + 0.3, min(width, 4.2), 0.35)
            self.write_paragraph(
                attribution_box.text_frame,
                attribution,
                size=self.theme.typography.body_size - 1,
                color=self.theme.colors.muted,
            )
            self.fit_text_frame(
                attribution_box.text_frame,
                max_size=self.theme.typography.body_size - 1,
            )

    def set_run_style(
        self,
        run,
        *,
        size: int,
        color: str,
        bold: bool = False,
        italic: bool = False,
        font_name: str | None = None,
    ) -> None:
        font = run.font
        self.set_font_style(
            font,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            font_name=font_name,
        )

    def set_font_style(
        self,
        font: "Font",
        *,
        size: int,
        color: str,
        bold: bool = False,
        italic: bool = False,
        font_name: str | None = None,
    ) -> None:
        font.name = font_name or self.theme.typography.font_name
        font.size = Pt(size)
        font.bold = bold
        font.italic = italic
        font.color.rgb = rgb(color)

    def write_paragraph(
        self,
        text_frame: "TextFrame",
        text: str,
        *,
        size: int,
        color: str,
        bold: bool = False,
        level: int = 0,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        space_after: int = 0,
        italic: bool = False,
    ):
        paragraph = text_frame.paragraphs[0] if not text_frame.text and len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].runs else text_frame.add_paragraph()
        paragraph.text = ""
        paragraph.level = level
        paragraph.alignment = align
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = text
        self.set_run_style(run, size=size, color=color, bold=bold, italic=italic)
        return paragraph

    def add_rule(self, slide: "PptxSlide", x1: float, y1: float, x2: float, y2: float, *, color: str, width_pt: float = 1.25):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        line.line.color.rgb = rgb(color)
        line.line.width = Pt(width_pt)
        return line

    def add_panel(
        self,
        slide: "PptxSlide",
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        fill_color: str,
        line_color: str,
    ):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(self.theme.components.panel_border_width_pt)
        return shape

    def add_chart_panel(
        self,
        slide: "PptxSlide",
        left: float,
        top: float,
        width: float,
        height: float,
    ):
        return self.add_panel(
            slide,
            left,
            top,
            width,
            height,
            fill_color=self.theme.colors.surface,
            line_color=self.theme.colors.line,
        )

    def add_accent_bar(self, slide: "PptxSlide", left: float, top: float, width: float, height: float, *, color: str):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        return shape

    def add_footer(self, slide: "PptxSlide", meta: PresentationMeta, index: int, total_slides: int) -> None:
        grid = self.theme.grid
        left_text = meta.footer_text or meta.client_name or meta.author or meta.title
        if meta.date:
            left_text = f"{left_text}  •  {meta.date}" if left_text else meta.date

        left_box = self.textbox(slide, grid.content_left, grid.footer_top, 6.0, 0.2)
        left_paragraph = left_box.text_frame.paragraphs[0]
        left_paragraph.alignment = PP_ALIGN.LEFT
        left_run = left_paragraph.add_run()
        left_run.text = left_text or ""
        self.set_run_style(left_run, size=self.theme.typography.small_size, color=self.theme.colors.muted)

        right_box = self.textbox(slide, 10.9, grid.footer_top, 1.35, 0.2)
        right_paragraph = right_box.text_frame.paragraphs[0]
        right_paragraph.alignment = PP_ALIGN.RIGHT
        right_run = right_paragraph.add_run()
        right_run.text = f"{index:02d} / {total_slides:02d}"
        self.set_run_style(right_run, size=self.theme.typography.small_size, color=self.theme.colors.muted)

        self.add_rule(
            slide,
            grid.content_left,
            grid.footer_line_y,
            grid.content_right,
            grid.footer_line_y,
            color=self.theme.colors.line,
            width_pt=self.theme.components.footer_rule_width_pt,
        )

    def add_speaker_notes(self, slide: "PptxSlide", notes: str | None) -> None:
        if not notes:
            return

        notes_slide = slide.notes_slide
        text_frame = getattr(notes_slide, "notes_text_frame", None)
        if text_frame is None:
            fallback_placeholders = [placeholder for placeholder in notes_slide.placeholders if hasattr(placeholder, "text_frame")]
            if not fallback_placeholders:
                return
            text_frame = fallback_placeholders[-1].text_frame

        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = notes.strip()
        self.set_run_style(run, size=11, color=self.theme.colors.text)

    def resolve_asset(self, image_path: str | None) -> Path | None:
        if not image_path:
            return None
        candidate = Path(image_path)
        if not candidate.is_absolute():
            candidate = self.asset_root / candidate
        return candidate if candidate.exists() else None

    def resolve_brand_logo(self, meta: PresentationMeta) -> Path | None:
        return self.resolve_asset(meta.logo_path)

    def collect_missing_assets(self, spec: PresentationInput) -> list[str]:
        missing_assets: list[str] = []
        if spec.presentation.logo_path and self.resolve_brand_logo(spec.presentation) is None:
            missing_assets.append(
                f"presentation branding: missing asset '{spec.presentation.logo_path}'"
            )

        for index, slide_spec in enumerate(spec.slides, start=1):
            if not slide_spec.image_path:
                continue

            if self.resolve_asset(slide_spec.image_path) is None:
                slide_label = slide_spec.title or slide_spec.type.value
                missing_assets.append(
                    f"slide {index:02d} ({slide_label}): missing asset '{slide_spec.image_path}'"
                )
        return missing_assets


def render_presentation(
    spec: PresentationInput,
    output_path: str | Path,
    *,
    theme_name: str | None = None,
    asset_root: str | Path | None = None,
    primary_color: str | None = None,
    secondary_color: str | None = None,
) -> Path:
    renderer = PresentationRenderer(
        theme_name=theme_name,
        asset_root=asset_root,
        primary_color=primary_color,
        secondary_color=secondary_color,
    )
    return renderer.render(spec, output_path)

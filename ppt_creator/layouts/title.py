from __future__ import annotations

from pptx.util import Inches

from ppt_creator.theme import theme_display_name


def _is_dense_cover_copy(title: str | None, subtitle: str | None) -> bool:
    return len(title or "") > 82 or len(subtitle or "") > 78


def _render_body_panel(renderer, slide, *, left, top, width, text, typography, colors, dense: bool) -> None:
    panel_width = min(width, 5.35 if dense else 5.75)
    panel_height = 1.02 if dense else 0.92
    renderer.add_panel(
        slide,
        left,
        top,
        panel_width,
        panel_height,
        fill_color=colors.surface,
        line_color=colors.line,
    )
    content_left, content_top, content_width, content_height = renderer.panel_inner_bounds(
        left=left,
        top=top,
        width=panel_width,
        height=panel_height,
        padding=0.18 if dense else 0.20,
    )
    body_box = renderer.textbox(slide, content_left, content_top + 0.02, content_width, content_height - 0.04)
    renderer.write_paragraph(
        body_box.text_frame,
        text,
        size=typography.body_size - (1 if dense else 0),
        color=colors.text,
    )
    renderer.fit_text_frame(
        body_box.text_frame,
        max_size=typography.body_size - (1 if dense else 0),
        min_size=typography.small_size + 1,
    )


def _render_meta_blocks(renderer, slide, *, left, top, width, height, blocks, typography, colors) -> None:
    if not blocks:
        return

    regions = [
        {
            "kind": f"block_{index + 1}",
            "min_height": block.get("min_height", 0.34),
            "flex": 1.0,
            "content_weight": renderer.estimate_content_weight(
                title=block.get("label"),
                body=block.get("value"),
            ),
        }
        for index, block in enumerate(blocks)
    ]

    for block, (_, (region_top, region_height)) in zip(
        blocks,
        renderer.build_content_stack(
            top=top,
            height=height,
            regions=regions,
            gap=0.08,
            min_flex=0.9,
            max_flex=1.25,
        ),
        strict=True,
    ):
        box = renderer.textbox(slide, left, region_top, width, region_height)
        tf = box.text_frame
        if block.get("label"):
            renderer.write_paragraph(
                tf,
                block["label"],
                size=typography.small_size,
                color=colors.muted,
                bold=True,
                space_after=2,
            )
        renderer.write_paragraph(
            tf,
            block["value"],
            size=block.get("size") or (typography.body_size - 1),
            color=block.get("color") or colors.text,
            bold=bool(block.get("bold", False)),
        )


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    g = renderer.theme.grid
    t = renderer.theme.typography
    colors = renderer.theme.colors
    variant = renderer.resolve_layout_variant(slide_spec, "split_panel")

    logo_asset = renderer.resolve_brand_logo(meta)
    cover_asset = renderer.resolve_asset(slide_spec.image_path)
    focal_x, focal_y = renderer.resolve_image_focal_point(slide_spec)
    dense_cover = _is_dense_cover_copy(slide_spec.title or meta.title, slide_spec.subtitle)

    if variant == "hero_cover":
        renderer.add_accent_bar(slide, g.content_left, 0.78, g.content_width, 0.08, color=colors.accent)

        eyebrow_text = slide_spec.eyebrow or meta.client_name or meta.subtitle or theme_display_name(renderer.theme.name)
        hero_columns = renderer.build_constrained_columns(
            left=g.content_left,
            width=g.content_width,
            gap=0.38,
            regions=[
                {
                    "kind": "hero_main",
                    "min_width": 6.0,
                    "target_share": renderer.estimate_content_weight(
                        title=slide_spec.title or meta.title,
                        body=slide_spec.body,
                        footer=slide_spec.subtitle,
                        tag=eyebrow_text,
                    ),
                },
                {
                    "kind": "hero_context",
                    "width": 3.15 if cover_asset else 3.35,
                    "min_width": 2.85,
                },
            ],
        )
        main_left, main_width = hero_columns[0]
        panel_left, panel_width = hero_columns[1]

        if logo_asset:
            slide.shapes.add_picture(
                str(logo_asset),
                Inches(main_left),
                Inches(0.55),
                width=Inches(1.8),
                height=Inches(0.34),
            )

        if eyebrow_text:
            renderer.add_eyebrow(slide, eyebrow_text, left=main_left, top=1.2, width=min(main_width, 5.4), uppercase=False)

        heading_bottom = renderer.add_heading(
            slide,
            title=slide_spec.title or meta.title,
            subtitle=slide_spec.subtitle,
            left=main_left,
            top=1.55,
            width=min(main_width, 6.85 if dense_cover else 7.2),
            subtitle_width=min(main_width, 5.9 if dense_cover else 6.4),
            title_size=t.title_size + (4 if dense_cover else 6),
            max_title_lines=2,
            max_subtitle_lines=2,
            min_title_size=22 if dense_cover else 24,
            min_subtitle_size=11,
            title_subtitle_gap=0.16,
        )

        if slide_spec.body:
            body_top = max(3.72, heading_bottom + 0.54)
            _render_body_panel(
                renderer,
                slide,
                left=main_left,
                top=body_top,
                width=main_width,
                text=slide_spec.body,
                typography=t,
                colors=colors,
                dense=dense_cover,
            )

        if cover_asset:
            renderer.add_image_cover(
                slide,
                cover_asset,
                left=panel_left,
                top=1.6,
                width=panel_width,
                height=3.7,
                focal_x=focal_x,
                focal_y=focal_y,
            )
            renderer.add_panel(
                slide,
                panel_left + 0.16,
                4.12,
                panel_width - 0.32,
                1.0,
                fill_color=colors.surface,
                line_color=colors.line,
            )
            content_left, content_top, content_width, content_height = renderer.panel_inner_bounds(
                left=panel_left + 0.16,
                top=4.12,
                width=panel_width - 0.32,
                height=1.0,
                padding=0.18,
            )
        else:
            renderer.add_panel(
                slide,
                panel_left,
                1.6,
                panel_width,
                3.7,
                fill_color=colors.surface,
                line_color=colors.line,
            )
            renderer.add_accent_bar(slide, panel_left, 1.6, panel_width, renderer.theme.components.accent_bar_height, color=colors.navy)
            content_left, content_top, content_width, content_height = renderer.panel_inner_bounds(
                left=panel_left,
                top=1.6,
                width=panel_width,
                height=3.7,
                padding=0.28,
            )
        _render_meta_blocks(
            renderer,
            slide,
            left=content_left,
            top=content_top,
            width=content_width,
            height=content_height,
            blocks=[
                *(
                    [{"label": "Client", "value": meta.client_name, "min_height": 0.34}]
                    if meta.client_name
                    else []
                ),
                *(
                    [{"label": "Author", "value": meta.author, "min_height": 0.28}]
                    if meta.author
                    else []
                ),
                *(
                    [{"label": "Date", "value": meta.date, "min_height": 0.28}]
                    if meta.date
                    else []
                ),
            ],
            typography=t,
            colors=colors,
        )
        return

    renderer.add_accent_bar(slide, g.content_left, 0.72, 0.1, 1.15, color=colors.accent)

    heading_left = g.content_left + 0.33
    split_columns = renderer.build_constrained_columns(
        left=heading_left,
        width=g.content_right - heading_left,
        gap=0.42,
        regions=[
            {
                "kind": "split_main",
                "min_width": 5.6,
                "target_share": renderer.estimate_content_weight(
                    title=slide_spec.title or meta.title,
                    body=slide_spec.body,
                    footer=slide_spec.subtitle,
                    tag=slide_spec.eyebrow or meta.subtitle,
                ),
            },
            {
                "kind": "split_panel",
                "width": 3.35 if cover_asset else 3.6,
                "min_width": 3.0,
            },
        ],
    )
    main_left, main_width = split_columns[0]
    panel_left, panel_width = split_columns[1]

    if logo_asset:
        slide.shapes.add_picture(
            str(logo_asset),
            Inches(panel_left),
            Inches(0.55),
            width=Inches(min(2.75, panel_width)),
            height=Inches(0.48),
        )

    heading_bottom = renderer.add_heading(
        slide,
        title=slide_spec.title or meta.title,
        subtitle=slide_spec.subtitle,
        eyebrow=slide_spec.eyebrow or meta.subtitle,
        left=main_left,
        top=1.25,
        width=min(main_width, 6.15 if dense_cover else 6.7),
        subtitle_width=min(main_width, 5.35 if dense_cover else 5.9),
        title_size=t.title_size + (2 if dense_cover else 4),
        max_title_lines=2,
        max_subtitle_lines=2,
        min_title_size=21 if dense_cover else 23,
        min_subtitle_size=11,
        title_subtitle_gap=0.16,
    )

    if slide_spec.body:
        body_top = max(3.52, heading_bottom + 0.54)
        _render_body_panel(
            renderer,
            slide,
            left=main_left,
            top=body_top,
            width=main_width,
            text=slide_spec.body,
            typography=t,
            colors=colors,
            dense=dense_cover,
        )

    if cover_asset:
        renderer.add_image_cover(
            slide,
            cover_asset,
            left=panel_left,
            top=1.05,
            width=panel_width,
            height=4.65,
            focal_x=focal_x,
            focal_y=focal_y,
        )
        renderer.add_panel(
            slide,
            panel_left + 0.2,
            3.95,
            panel_width - 0.4,
            1.35,
            fill_color=colors.surface,
            line_color=colors.line,
        )
        content_left, content_top, content_width, content_height = renderer.panel_inner_bounds(
            left=panel_left + 0.2,
            top=3.95,
            width=panel_width - 0.4,
            height=1.35,
            padding=0.22,
        )
    else:
        renderer.add_panel(
            slide,
            panel_left,
            1.05,
            panel_width,
            4.65,
            fill_color=colors.surface,
            line_color=colors.line,
        )

        content_left, content_top, content_width, content_height = renderer.panel_inner_bounds(
            left=panel_left,
            top=1.05,
            width=panel_width,
            height=4.65,
            padding=0.40,
        )
    _render_meta_blocks(
        renderer,
        slide,
        left=content_left,
        top=content_top,
        width=content_width,
        height=content_height,
        blocks=[
            *(
                [{"label": "Client", "value": meta.client_name, "min_height": 0.52}]
                if meta.client_name
                else []
            ),
            *(
                [{"label": "Author", "value": meta.author, "min_height": 0.52}]
                if meta.author
                else []
            ),
            *(
                [{"label": "Date", "value": meta.date, "min_height": 0.52}]
                if meta.date
                else []
            ),
        ],
        typography=t,
        colors=colors,
    )

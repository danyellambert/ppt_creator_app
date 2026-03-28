from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from ppt_creator.theme import rgb

CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
}


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    g = renderer.theme.grid
    t = renderer.theme.typography
    colors = renderer.theme.colors
    variant = renderer.resolve_layout_variant(slide_spec, "column")

    renderer.add_heading(
        slide,
        title=slide_spec.title or "",
        subtitle=slide_spec.subtitle,
        eyebrow=slide_spec.eyebrow or "Data view",
        left=g.content_left,
        top=g.title_top,
        width=8.2,
        subtitle_width=7.0,
    )

    chart_regions: list[dict[str, float | str]] = []
    if slide_spec.body:
        chart_regions.append(
            {
                "kind": "body",
                "min_height": 0.52,
                "flex": 1.0,
                "content_weight": renderer.estimate_content_weight(body=slide_spec.body),
            }
        )
    chart_regions.append(
        {
            "kind": "chart",
            "min_height": 2.85,
            "flex": 1.0,
            "content_weight": max(1.0, len(slide_spec.chart_categories) * 0.45 + len(slide_spec.chart_series) * 0.8),
        }
    )

    chart_top = 2.2
    chart_height = 4.05
    chart_bounds: tuple[float, float, float, float] | None = None
    for region, (region_top, region_height) in renderer.build_constrained_content_stack(
        top=chart_top,
        height=chart_height,
        regions=[
            {
                **region,
                **(
                    {
                        "target_share": 1.0,
                        "max_height": 0.9,
                    }
                    if region["kind"] == "body"
                    else {
                        "target_share": max(
                            2.4,
                            len(slide_spec.chart_categories) * 0.5 + len(slide_spec.chart_series) * 0.85,
                        )
                    }
                ),
            }
            for region in chart_regions
        ],
        gap=0.18,
    ):
        if region["kind"] == "body":
            body_box = renderer.textbox(slide, g.content_left, region_top, g.content_width, region_height)
            renderer.write_paragraph(
                body_box.text_frame,
                slide_spec.body or "",
                size=t.body_size - 1,
                color=colors.text,
            )
            renderer.fit_text_frame(body_box.text_frame, max_size=t.body_size - 1)
        else:
            chart_bounds = (g.content_left, region_top, g.content_width, region_height)

    chart_data = CategoryChartData()
    chart_data.categories = slide_spec.chart_categories
    for series in slide_spec.chart_series:
        chart_data.add_series(series.name, series.values)

    assert chart_bounds is not None
    chart_left, chart_top, chart_width, chart_height = chart_bounds

    chart_frame = slide.shapes.add_chart(
        CHART_TYPES[variant],
        Inches(chart_left),
        Inches(chart_top),
        Inches(chart_width),
        Inches(chart_height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(slide_spec.chart_series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(t.small_size)
        chart.legend.font.color.rgb = rgb(colors.muted)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(t.small_size)
    category_axis.tick_labels.font.color.rgb = rgb(colors.muted)
    category_axis.format.line.color.rgb = rgb(colors.line)

    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(t.small_size)
    value_axis.tick_labels.font.color.rgb = rgb(colors.muted)
    value_axis.format.line.color.rgb = rgb(colors.line)
    value_axis.major_gridlines.format.line.color.rgb = rgb(colors.line)

    plot = chart.plots[0]
    if hasattr(plot, "gap_width"):
        category_count = len(slide_spec.chart_categories)
        plot.gap_width = 70 if category_count <= 4 else 55 if category_count <= 6 else 38

    palette = [colors.navy, colors.accent, colors.text, colors.muted]
    for idx, series in enumerate(chart.series):
        color = palette[idx % len(palette)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb(color)
        series.format.line.color.rgb = rgb(color)
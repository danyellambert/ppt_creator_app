from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from ppt_creator.theme import rgb


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    g = renderer.theme.grid
    t = renderer.theme.typography
    colors = renderer.theme.colors

    renderer.add_heading(
        slide,
        title=slide_spec.title or "",
        subtitle=slide_spec.subtitle,
        eyebrow=slide_spec.eyebrow,
        left=g.content_left,
        top=g.title_top,
        width=8.0,
        subtitle_width=7.4,
    )

    items = slide_spec.timeline_items
    gap = 0.24
    marker_y = 2.28
    panel_top = 2.55
    panel_height = 2.9

    panel_bounds = renderer.build_panel_row_bounds(
        left=g.content_left,
        top=panel_top,
        width=g.content_width,
        height=panel_height,
        gap=gap,
        regions=[
            {
                "kind": f"timeline_{index + 1}",
                "min_width": 1.85,
                "target_share": renderer.estimate_content_weight(
                    title=item.title,
                    body=item.body,
                    footer=item.footer,
                    tag=item.tag,
                ),
            }
            for index, item in enumerate(items)
        ],
    )

    line_start = panel_bounds[0][0] + (panel_bounds[0][2] / 2)
    line_end = panel_bounds[-1][0] + (panel_bounds[-1][2] / 2)

    renderer.add_rule(
        slide,
        line_start,
        marker_y,
        line_end,
        marker_y,
        color=colors.line,
        width_pt=1.4,
    )

    for idx, (item, (left, _, item_width, item_height)) in enumerate(zip(items, panel_bounds, strict=True)):
        marker_color = colors.navy if idx % 2 == 0 else colors.accent
        marker_center = left + (item_width / 2)

        marker = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(marker_center - 0.12),
            Inches(marker_y - 0.12),
            Inches(0.24),
            Inches(0.24),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = rgb(marker_color)
        marker.line.fill.background()

        label_box = renderer.textbox(slide, left, 1.95, item_width, 0.22)
        renderer.write_paragraph(
            label_box.text_frame,
            item.tag or f"Step {idx + 1:02d}",
            size=t.small_size,
            color=colors.muted,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        renderer.fit_text_frame(label_box.text_frame, max_size=t.small_size, bold=True)

        renderer.add_panel(
            slide,
            left,
            panel_top,
            item_width,
            item_height,
            fill_color=colors.surface,
            line_color=colors.line,
        )
        renderer.add_accent_bar(
            slide,
            left,
            panel_top,
            item_width,
            renderer.theme.components.accent_bar_height,
            color=marker_color,
        )

        regions = [{"kind": "title", "height": 0.34}]
        if item.body:
            regions.append(
                {
                    "kind": "body",
                    "min_height": 0.92,
                    "flex": 1.0,
                    "content_weight": renderer.estimate_content_weight(body=item.body),
                }
            )
        if item.footer:
            regions.append({"kind": "footer", "height": 0.22})

        for region, (content_left, region_top, content_width, region_height) in renderer.build_panel_content_stack_bounds(
            left=left,
            top=panel_top,
            width=item_width,
            height=item_height,
            regions=regions,
            gap=0.08,
            padding=0.22,
            min_flex=0.9,
            max_flex=1.25,
        ):
            box = renderer.textbox(slide, content_left, region_top, content_width, region_height)
            if region["kind"] == "title":
                renderer.write_paragraph(
                    box.text_frame,
                    item.title,
                    size=t.body_size - 1,
                    color=colors.navy,
                    bold=True,
                )
                renderer.fit_text_frame(box.text_frame, max_size=t.body_size - 1, bold=True)
            elif region["kind"] == "body":
                renderer.write_paragraph(
                    box.text_frame,
                    item.body or "",
                    size=t.small_size + 1,
                    color=colors.text,
                )
                renderer.fit_text_frame(box.text_frame, max_size=t.small_size + 1)
            else:
                renderer.write_paragraph(
                    box.text_frame,
                    item.footer or "",
                    size=t.small_size,
                    color=marker_color,
                    bold=True,
                )
                renderer.fit_text_frame(box.text_frame, max_size=t.small_size, bold=True)
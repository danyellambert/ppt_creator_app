from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    c = renderer.theme.canvas
    t = renderer.theme.typography
    components = renderer.theme.components
    colors = renderer.theme.colors
    variant = renderer.resolve_layout_variant(slide_spec, "image_right")

    if slide_spec.eyebrow:
        eyebrow_box = renderer.textbox(slide, c.margin_x, 0.78, 4.6, 0.25)
        renderer.write_paragraph(eyebrow_box.text_frame, slide_spec.eyebrow.upper(), size=t.eyebrow_size, color=colors.accent, bold=True)
        renderer.fit_text_frame(eyebrow_box.text_frame, max_size=t.eyebrow_size, bold=True)

    title_box = renderer.textbox(slide, c.margin_x, 1.02, 5.4, 0.8)
    renderer.write_paragraph(title_box.text_frame, slide_spec.title or "", size=t.title_size, color=colors.navy, bold=True)
    renderer.fit_text_frame(title_box.text_frame, max_size=t.title_size, bold=True)

    if slide_spec.subtitle:
        subtitle_box = renderer.textbox(slide, c.margin_x, 1.8, 5.4, 0.45)
        renderer.write_paragraph(subtitle_box.text_frame, slide_spec.subtitle, size=t.subtitle_size, color=colors.muted)
        renderer.fit_text_frame(subtitle_box.text_frame, max_size=t.subtitle_size)

    image_left = 7.1 if variant == "image_right" else c.margin_x
    image_top = 1.28
    image_width = 5.15
    image_height = 4.95
    text_left = c.margin_x if variant == "image_right" else 7.1
    text_width = 5.4 if variant == "image_right" else 5.15
    asset = renderer.resolve_asset(slide_spec.image_path)
    focal_x, focal_y = renderer.resolve_image_focal_point(slide_spec)
    placeholder_copy = renderer.describe_visual_placeholder(slide_spec)
    if asset:
        renderer.add_image_cover(
            slide,
            asset,
            left=image_left,
            top=image_top,
            width=image_width,
            height=image_height,
            focal_x=focal_x,
            focal_y=focal_y,
        )
    else:
        panel = renderer.add_panel(slide, image_left, image_top, image_width, image_height, fill_color=colors.surface, line_color=colors.line)
        renderer.add_accent_bar(slide, image_left, image_top, image_width, components.accent_bar_height, color=colors.accent)
        placeholder = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            panel.left,
            panel.top,
            panel.width,
            panel.height,
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = panel.fill.fore_color.rgb
        placeholder.line.fill.background()

        inner_left = image_left + 0.55
        inner_top = image_top + 0.72
        inner_width = image_width - 1.1
        inner_height = 1.42
        renderer.add_panel(
            slide,
            inner_left,
            inner_top,
            inner_width,
            inner_height,
            fill_color=colors.soft_fill,
            line_color=colors.line,
        )
        renderer.add_rule(
            slide,
            inner_left + 0.18,
            inner_top + 0.18,
            inner_left + inner_width - 0.18,
            inner_top + inner_height - 0.18,
            color=colors.line,
            width_pt=1.0,
        )
        renderer.add_rule(
            slide,
            inner_left + inner_width - 0.18,
            inner_top + 0.18,
            inner_left + 0.18,
            inner_top + inner_height - 0.18,
            color=colors.line,
            width_pt=1.0,
        )

        label_box = renderer.textbox(slide, image_left + 0.55, image_top + 0.35, 2.1, 0.26)
        renderer.write_paragraph(
            label_box.text_frame,
            placeholder_copy["label"],
            size=t.small_size,
            color=colors.accent,
            bold=True,
        )
        renderer.fit_text_frame(label_box.text_frame, max_size=t.small_size, bold=True)

        heading_box = renderer.textbox(slide, image_left + 0.55, image_top + 2.36, image_width - 1.1, 0.35)
        renderer.write_paragraph(
            heading_box.text_frame,
            placeholder_copy["headline"],
            size=t.body_size,
            color=colors.muted,
            bold=True,
        )
        renderer.fit_text_frame(heading_box.text_frame, max_size=t.body_size, bold=True)

        caption_box = renderer.textbox(slide, image_left + 0.55, image_top + 2.78, image_width - 1.1, 0.9)
        renderer.write_paragraph(
            caption_box.text_frame,
            slide_spec.image_caption
            or placeholder_copy["guidance"],
            size=t.small_size + 1,
            color=colors.text,
        )
        renderer.fit_text_frame(caption_box.text_frame, max_size=t.small_size + 1)

        if slide_spec.image_path:
            path_box = renderer.textbox(slide, image_left + 0.55, image_top + 3.82, image_width - 1.1, 0.55)
            renderer.write_paragraph(
                path_box.text_frame,
                f"Missing asset: {slide_spec.image_path}",
                size=t.small_size,
                color=colors.muted,
            )
            renderer.fit_text_frame(path_box.text_frame, max_size=t.small_size)

    if slide_spec.body:
        pass

    if slide_spec.body or slide_spec.bullets:
        text_regions: list[dict[str, float | str]] = []
        if slide_spec.body:
            if slide_spec.bullets:
                text_regions.append(
                    {
                        "kind": "body",
                        "min_height": 0.82,
                        "flex": 1.0,
                        "content_weight": renderer.estimate_content_weight(body=slide_spec.body),
                    }
                )
            else:
                text_regions.append({"kind": "body", "height": 2.35})
        if slide_spec.bullets:
            text_regions.append(
                {
                    "kind": "bullets",
                    "min_height": 1.5 if slide_spec.body else 3.1,
                    "flex": 1.0,
                    "content_weight": renderer.estimate_content_weight(bullets=slide_spec.bullets),
                }
            )

        for region, (region_top, region_height) in renderer.build_content_stack(
            top=2.45,
            height=3.1,
            regions=text_regions,
            gap=0.18,
            min_flex=0.9,
            max_flex=1.35,
        ):
            if region["kind"] == "body":
                body_box = renderer.textbox(slide, text_left, region_top, text_width, region_height)
                renderer.write_paragraph(body_box.text_frame, slide_spec.body or "", size=t.body_size, color=colors.text)
                renderer.fit_text_frame(body_box.text_frame, max_size=t.body_size)
            else:
                bullets_box = renderer.textbox(slide, text_left, region_top, text_width, region_height)
                tf = bullets_box.text_frame
                for bullet in slide_spec.bullets:
                    paragraph = tf.add_paragraph()
                    run = paragraph.add_run()
                    run.text = f"• {bullet}"
                    renderer.set_run_style(run, size=t.body_size - 1, color=colors.text)
                renderer.fit_text_frame(tf, max_size=t.body_size - 1)

    if slide_spec.image_caption and asset:
        caption_box = renderer.textbox(slide, image_left, 6.02, image_width, 0.24)
        renderer.write_paragraph(caption_box.text_frame, slide_spec.image_caption, size=t.small_size, color=colors.muted)
        renderer.fit_text_frame(caption_box.text_frame, max_size=t.small_size)

from __future__ import annotations

from pptx.enum.text import PP_ALIGN


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    c = renderer.theme.canvas
    t = renderer.theme.typography
    colors = renderer.theme.colors

    renderer.add_accent_bar(slide, c.margin_x, 1.0, 0.1, 4.8, color=colors.accent)

    if slide_spec.title:
        title_box = renderer.textbox(slide, 1.28, 1.15, 5.0, 0.7)
        renderer.write_paragraph(title_box.text_frame, slide_spec.title, size=t.eyebrow_size + 1, color=colors.accent, bold=True)

    quote_box = renderer.textbox(slide, 1.28, 2.0, 8.8, 1.7)
    paragraph = quote_box.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = slide_spec.quote or slide_spec.body or ""
    renderer.set_run_style(run, size=t.quote_size, color=colors.navy, bold=True, italic=True)

    if slide_spec.attribution:
        attribution_box = renderer.textbox(slide, 1.28, 4.0, 4.2, 0.45)
        renderer.write_paragraph(attribution_box.text_frame, slide_spec.attribution, size=t.body_size - 1, color=colors.muted)

    renderer.add_panel(slide, 9.2, 1.55, 3.05, 3.0, fill_color=colors.surface, line_color=colors.line)
    box = renderer.textbox(slide, 9.55, 1.95, 2.35, 2.2)
    tf = box.text_frame
    renderer.write_paragraph(tf, "Next", size=t.eyebrow_size, color=colors.accent, bold=True, space_after=8)
    renderer.write_paragraph(tf, "Approve the narrative, connect your content pipeline, and reuse the same renderer across future decks.", size=t.body_size - 1, color=colors.text)

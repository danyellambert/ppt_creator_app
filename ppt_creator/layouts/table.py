from __future__ import annotations

from pptx.enum.text import PP_ALIGN


def _table_char_count(columns, rows) -> int:
    return sum(len(text) for text in [*columns, *(cell for row in rows for cell in row)] if text)


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    g = renderer.theme.grid
    t = renderer.theme.typography
    colors = renderer.theme.colors
    columns = slide_spec.table_columns
    rows = slide_spec.table_rows
    dense_table = bool(
        len(columns) >= 4
        or len(rows) >= 4
        or any(len(column) > 18 for column in columns)
        or any(len(cell) > 34 for row in rows for cell in row)
        or _table_char_count(columns, rows) > 220
    )

    renderer.add_heading(
        slide,
        title=slide_spec.title or "",
        subtitle=slide_spec.subtitle,
        eyebrow=slide_spec.eyebrow or "Executive table",
        left=g.content_left,
        top=g.title_top,
        width=7.8,
        subtitle_width=6.3,
    )

    top = 2.10
    if slide_spec.body:
        intro_width = min(g.content_width, 6.8)
        renderer.add_panel(
            slide,
            g.content_left,
            top,
            intro_width,
            0.64 if dense_table else 0.58,
            fill_color=colors.surface,
            line_color=colors.line,
        )
        intro_left, intro_top, intro_inner_width, intro_inner_height = renderer.panel_inner_bounds(
            left=g.content_left,
            top=top,
            width=intro_width,
            height=0.64 if dense_table else 0.58,
            padding=0.14,
        )
        intro_box = renderer.textbox(slide, intro_left, intro_top, intro_inner_width, intro_inner_height)
        renderer.write_paragraph(
            intro_box.text_frame,
            slide_spec.body,
            size=t.body_size - 2 if dense_table else t.body_size - 1,
            color=colors.text,
        )
        renderer.fit_text_frame(
            intro_box.text_frame,
            max_size=t.body_size - 2 if dense_table else t.body_size - 1,
            min_size=t.small_size + 1,
        )
        top += 0.74 if dense_table else 0.68

    gap = 0.07 if dense_table else 0.08
    header_height = 0.60 if any(len(column) > 16 for column in columns) else 0.54
    row_height = 0.50 if dense_table else 0.52

    column_weights = [
        renderer.estimate_content_weight(
            title=column,
            body=" ".join(row[index] for row in rows if index < len(row)),
        )
        for index, column in enumerate(columns)
    ]
    column_regions = [
        {
            "kind": f"column_{index + 1}",
            "min_width": 1.5 if index == 0 else 1.05,
            "target_share": weight,
            **({"max_width": 4.6} if index == 0 and len(columns) >= 3 else {}),
        }
        for index, weight in enumerate(column_weights)
    ]

    column_header_bounds = renderer.build_constrained_panel_row_content_bounds(
        left=g.content_left,
        top=top,
        width=g.content_width,
        height=header_height,
        gap=gap,
        regions=column_regions,
        padding=0.09 if dense_table else 0.08,
    )

    available_rows_height = max(0.8, (g.footer_line_y - 0.30) - (top + header_height + 0.10))
    row_min_height = max(
        0.3 if dense_table else 0.28,
        min(
            row_height,
            (available_rows_height - (0.10 * max(0, len(rows) - 1))) / max(1, len(rows)),
        ),
    )
    row_weights = [renderer.estimate_content_weight(body=" ".join(row)) for row in rows]
    row_regions = [
        {
            "kind": f"row_{index + 1}",
            "min_height": row_min_height,
            "target_share": weight,
            **({"max_height": row_min_height * 1.24} if dense_table else {}),
        }
        for index, weight in enumerate(row_weights)
    ]

    cell_grid = renderer.build_constrained_panel_grid_content_bounds(
        left=g.content_left,
        top=top + header_height + 0.08,
        width=g.content_width,
        height=available_rows_height,
        column_gap=gap,
        row_gap=0.08 if dense_table else 0.09,
        column_regions=column_regions,
        row_regions=row_regions,
        padding=0.09 if dense_table else 0.08,
    )

    header_size = t.small_size + 1
    body_size = t.small_size + (0 if dense_table else 1)
    for column, (panel_bounds, content_bounds) in zip(columns, column_header_bounds, strict=True):
        left, panel_top, column_width, panel_height = panel_bounds
        content_left, content_top, content_width, content_height = content_bounds
        renderer.add_panel(
            slide,
            left,
            panel_top,
            column_width,
            panel_height,
            fill_color=colors.soft_fill,
            line_color=colors.line,
        )
        renderer.add_accent_bar(slide, left, panel_top, column_width, 0.04, color=colors.accent if len(columns) <= 4 else colors.navy)
        header_box = renderer.textbox(slide, content_left, content_top, content_width, content_height)
        renderer.write_paragraph(
            header_box.text_frame,
            column,
            size=header_size,
            color=colors.navy,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        renderer.fit_text_frame(header_box.text_frame, max_size=header_size, min_size=t.small_size - 1, bold=True)

    for row, row_cells in zip(rows, cell_grid, strict=True):
        for col_index, (cell, (panel_bounds, content_bounds)) in enumerate(zip(row, row_cells, strict=True), start=1):
            left, panel_top, column_width, panel_height = panel_bounds
            content_left, content_top, content_width, content_height = content_bounds
            renderer.add_panel(
                slide,
                left,
                panel_top,
                column_width,
                panel_height,
                fill_color=colors.surface if (col_index + rows.index(row)) % 2 else colors.background,
                line_color=colors.line,
            )
            cell_box = renderer.textbox(slide, content_left, content_top, content_width, content_height)
            renderer.write_paragraph(
                cell_box.text_frame,
                cell,
                size=body_size,
                color=colors.text,
                align=PP_ALIGN.LEFT if col_index == 1 or len(cell) > 18 else PP_ALIGN.CENTER,
            )
            renderer.fit_text_frame(cell_box.text_frame, max_size=body_size, min_size=t.small_size - 1)
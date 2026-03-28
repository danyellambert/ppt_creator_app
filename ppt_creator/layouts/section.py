from __future__ import annotations

from pptx.enum.text import PP_ALIGN


def render(renderer, slide, slide_spec, meta, index, total_slides) -> None:
    g = renderer.theme.grid
    t = renderer.theme.typography
    colors = renderer.theme.colors

    section_columns = renderer.build_constrained_columns(
        left=g.content_left,
        width=g.content_width,
        gap=0.42,
        regions=[
            {
                "kind": "section_content",
                "min_width": 7.2,
                "target_share": renderer.estimate_content_weight(
                    title=slide_spec.title,
                    body=slide_spec.subtitle,
                    tag=slide_spec.section_label or slide_spec.eyebrow or "SECTION",
                ),
            },
            {
                "kind": "section_marker",
                "width": 1.65,
                "min_width": 1.5,
            },
        ],
    )
    content_left, content_width = section_columns[0]
    marker_left, marker_width = section_columns[1]

    renderer.add_accent_bar(slide, g.content_left, 1.15, g.content_width, 0.09, color=colors.navy)
    renderer.add_accent_bar(slide, g.content_left, 5.65, min(2.0, max(1.35, content_width * 0.28)), 0.09, color=colors.accent)

    renderer.add_heading(
        slide,
        title=slide_spec.title or "",
        subtitle=slide_spec.subtitle,
        eyebrow=slide_spec.section_label or slide_spec.eyebrow or "SECTION",
        left=content_left,
        top=2.4,
        width=content_width,
        subtitle_width=min(content_width, 6.4),
        title_size=t.section_size,
    )

    renderer.add_panel(slide, marker_left, 2.15, marker_width, 1.55, fill_color=colors.surface, line_color=colors.line)
    marker_box = renderer.textbox(slide, marker_left + 0.12, 2.48, marker_width - 0.24, 0.8)
    marker_p = marker_box.text_frame.paragraphs[0]
    marker_p.alignment = PP_ALIGN.CENTER
    marker_run = marker_p.add_run()
    marker_run.text = f"{index:02d}"
    renderer.set_run_style(marker_run, size=t.section_size - 4, color=colors.navy, bold=True)
    renderer.fit_text_frame(marker_box.text_frame, max_size=t.section_size - 4, bold=True)
